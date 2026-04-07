# -*- coding: utf-8 -*-
"""
Generateur PPTX pour import CANVA V2 - Rapport U62
BAHAFID Mohamed - BTS MEC Session 2026
Format A4 portrait | Design corporate ameliore | Compatible Canva
V2: ombres, images, donut charts, transparence
"""
import os, re
from pptx import Presentation
from pptx.util import Cm, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData
from lxml import etree

from contenu_v2 import (
    CANDIDAT, PAGE_01, PAGE_02, PAGE_03, PAGE_04, PAGE_05, PAGE_06, PAGE_07,
    PAGE_08, PAGE_09, PAGE_10, PAGE_15, PAGE_16, PAGE_17, PAGE_18,
    PAGE_20, PAGE_21, PAGE_22, PAGE_23, PAGE_24, PAGE_25, PAGE_26,
    PAGE_27, PAGE_28, PAGE_29, PAGE_30,
    SITUATION_1, SITUATION_2, SITUATION_3, SITUATION_4, SITUATION_5,
    TABLE_CORPS_ETAT, TABLE_BUDGET_ROUTE, TABLE_AUTRES_MARCHES, TABLE_COMPARAISON_REG,
)

BASE_DIR = r"D:\PREPA BTS MEC\08_U62_Rapport_Activites"
OUT = os.path.join(BASE_DIR, "Rapport_Redaction", "RAPPORT_U62_CANVA_V7.pptx")
IMG = os.path.join(BASE_DIR, "Annexes", "Images")
PHO = os.path.join(BASE_DIR, "Documents_Entreprise", "Photos_Chantier")

# Image paths
IMG_PHOTO = os.path.join(IMG, "photo_bahafid.jpg")
IMG_LOGO = os.path.join(IMG, "logo_bimco.png")
IMG_BANNER = os.path.join(IMG, "banniere_bimco.png")
DOC_IMGS = [os.path.join(IMG, f) for f in ["cps_marche46_signe.jpg", "convocation_cao.jpg", "notification_rejet.jpg"]]
CHANTIER_IMGS = [os.path.join(PHO, f) for f in [
    "20180523_140034.jpg", "20180605_122447.jpg", "20181206_130506.jpg",
    "20180503_164033.jpg", "20180330_154009.jpg", "20180228_144120.jpg",
]]

# XML namespace
_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'

# Colors
OG = RGBColor(0xF3, 0x92, 0x00)
NV = RGBColor(0x1E, 0x3A, 0x5F)
TQ = RGBColor(0x5C, 0xC8, 0xC0)
WH = RGBColor(0xFF, 0xFF, 0xFF)
LG = RGBColor(0xF5, 0xF5, 0xF5)
GY = RGBColor(0x88, 0x88, 0x88)
DK = RGBColor(0x2D, 0x2D, 0x2D)
PR = RGBColor(0x8B, 0x5C, 0xF6)
GN = RGBColor(0x10, 0xB9, 0x81)
RD = RGBColor(0xEF, 0x44, 0x44)
LO = RGBColor(0xFF, 0xF8, 0xF0)  # light orange bg

# A4 portrait
W = Cm(21)
H = Cm(29.7)

# ---- helpers ----

def add_shadow(shape, blur=4, dist=3, alpha=18):
    """Ombre douce sur une forme via XML."""
    spPr = shape._element.spPr
    for old in spPr.findall(f'{{{_a}}}effectLst'):
        spPr.remove(old)
    el = etree.SubElement(spPr, f'{{{_a}}}effectLst')
    sh = etree.SubElement(el, f'{{{_a}}}outerShdw', blurRad=str(blur*12700),
                          dist=str(dist*12700), dir='5400000', algn='tl', rotWithShape='0')
    clr = etree.SubElement(sh, f'{{{_a}}}srgbClr', val='000000')
    etree.SubElement(clr, f'{{{_a}}}alpha', val=str(alpha*1000))

def set_alpha(shape, opacity):
    """Transparence (0=invisible, 100=opaque)."""
    sf = shape._element.spPr.find(f'.//{{{_a}}}solidFill')
    if sf is not None and len(sf) > 0:
        c = sf[0]
        for old in c.findall(f'{{{_a}}}alpha'):
            c.remove(old)
        etree.SubElement(c, f'{{{_a}}}alpha', val=str(opacity*1000))

def pic(slide, path, left, top, w=None, h=None):
    """Insert image (retourne None si fichier absent)."""
    if not os.path.exists(path):
        return None
    kw = {'image_file': path, 'left': left, 'top': top}
    if w: kw['width'] = w
    if h: kw['height'] = h
    return slide.shapes.add_picture(**kw)

def add_rect(slide, left, top, w, h, fill, radius=None, shadow=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shape_type, left, top, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    if radius and hasattr(s, 'adjustments') and len(s.adjustments) > 0:
        s.adjustments[0] = radius
    if shadow:
        add_shadow(s)
    return s

def add_circle(slide, left, top, sz, fill, shadow=False, alpha=100):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, sz, sz)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    if shadow:
        add_shadow(s)
    if alpha < 100:
        set_alpha(s, alpha)
    return s

def add_text(slide, left, top, w, h, text, size=10, bold=False, color=DK, align=PP_ALIGN.LEFT, font_name="Montserrat"):
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return tb

def add_rich_text(slide, left, top, w, h, text, size=10, bold_color=NV,
                   color=RGBColor(0x55,0x55,0x55), align=PP_ALIGN.JUSTIFY, font_name="Inter"):
    """Textbox avec marqueurs **gras** rendus en bold_color."""
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    parts = re.split(r'(\*\*.*?\*\*)', text)
    for part in parts:
        if part.startswith('**') and part.endswith('**'):
            run = p.add_run()
            run.text = part[2:-2]
            run.font.bold = True
            run.font.color.rgb = bold_color
        else:
            run = p.add_run()
            run.text = part
            run.font.bold = False
            run.font.color.rgb = color
        run.font.size = Pt(size)
        run.font.name = font_name
    return tb

def add_para(tf, text, size=10, bold=False, color=DK, align=PP_ALIGN.LEFT, font_name="Inter", space_after=Pt(4)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    p.space_after = space_after
    return p

def footer(slide, num):
    band = add_rect(slide, Cm(0), Cm(28.8), W, Cm(0.9), NV)
    set_gradient(band, NV, _rgb_shift(NV, -15), 0)
    add_circle(slide, Cm(0.3), Cm(28.65), Cm(0.6), OG)
    add_rect(slide, Cm(19), Cm(28.8), Cm(2), Cm(0.12), OG)
    add_text(slide, Cm(0), Cm(28.9), W, Cm(0.7),
             f"BAHAFID Mohamed  \u2502  Rapport U62  \u2502  BTS MEC 2026  \u2502  {num}/30",
             7, False, WH, PP_ALIGN.CENTER, "Inter")

def add_badge(slide, left, top, text, bg_color, text_color=WH, w=Cm(3.5)):
    s = add_rect(slide, left, top, w, Cm(0.7), bg_color, 0.5, shadow=True)
    s.text_frame.word_wrap = False
    p = s.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(6.5)
    p.font.bold = True
    p.font.color.rgb = text_color
    p.font.name = "Montserrat"
    p.alignment = PP_ALIGN.CENTER
    return s

def add_kpi_card(slide, left, top, value, label, accent_color, w=Cm(4)):
    card = add_rect(slide, left, top, w, Cm(2.4), WH, 0.15, shadow=True)
    add_rect(slide, left, top, w, Cm(0.15), accent_color)
    add_text(slide, left+Cm(0.3), top+Cm(0.3), w-Cm(0.6), Cm(1),
             value, 18, True, accent_color, PP_ALIGN.CENTER, "Montserrat")
    add_text(slide, left+Cm(0.3), top+Cm(1.5), w-Cm(0.6), Cm(0.6),
             label, 7, False, GY, PP_ALIGN.CENTER, "Inter")

def orange_bar(slide, left, top, w=Cm(2.5)):
    add_rect(slide, left, top, w, Cm(0.12), OG)

def draw_table(slide, left, top, col_defs, headers, rows, header_fill=NV):
    """Tableau stylise corporate : header navy, lignes alternees, separateurs fins.
    col_defs : liste de (width, align)
    Retourne : y bas du tableau"""
    HDR_H = Cm(0.75)
    ROW_H = Cm(1.05)
    total_w = sum(w for w, _ in col_defs)
    # Header band
    h_rect = add_rect(slide, left, top, total_w, HDR_H, header_fill, 0.05, shadow=True)
    cx = left
    for (w, align), label in zip(col_defs, headers):
        add_text(slide, cx + Cm(0.15), top + Cm(0.1), w - Cm(0.3), HDR_H - Cm(0.1),
                 label, 6.5, True, WH, align, "Montserrat")
        cx += w
    # Data rows
    y = top + HDR_H
    for idx, row in enumerate(rows):
        is_total = str(row[0]).upper().startswith("TOTAL")
        bg = RGBColor(0xFD, 0xF3, 0xE4) if is_total else (LG if idx % 2 == 0 else WH)
        add_rect(slide, left, y, total_w, ROW_H, bg, 0.02)
        add_rect(slide, left, y + ROW_H - Cm(0.018), total_w, Cm(0.018),
                 RGBColor(0xDD, 0xDD, 0xDD))
        cx = left
        for (w, align), cell in zip(col_defs, row):
            fc = OG if is_total else DK
            add_text(slide, cx + Cm(0.15), y + Cm(0.12), w - Cm(0.3), ROW_H - Cm(0.12),
                     str(cell), 6.5, is_total, fc, align, "Inter")
            cx += w
        y += ROW_H
    # Bottom border
    add_rect(slide, left, y, total_w, Cm(0.04), NV)
    return y + Cm(0.04)

def add_donut(slide, left, top, w, h, cats, vals, colors, center_text=""):
    """Donut chart avec couleurs personnalisees."""
    try:
        cd = CategoryChartData()
        cd.categories = cats
        cd.add_series('Budget', vals)
        cf = slide.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, left, top, w, h, cd)
        chart = cf.chart
        chart.has_legend = False
        series = chart.series[0]
        for i, clr in enumerate(colors):
            pt = series.points[i]
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = clr
        if center_text:
            cx = left + w//2 - Cm(1.8)
            cy = top + h//2 - Cm(0.5)
            add_text(slide, cx, cy, Cm(3.6), Cm(1), center_text, 11, True, NV, PP_ALIGN.CENTER)
        return cf
    except Exception:
        return None

def _rgb_shift(color, delta):
    """Shift an RGBColor by delta (clamped 0-255)."""
    h = str(color)
    r, g, b = int(h[0:2],16), int(h[2:4],16), int(h[4:6],16)
    return RGBColor(max(0,min(255,r+delta)), max(0,min(255,g+delta)), max(0,min(255,b+delta)))

def constructys_header(slide, title, label="", bg=TQ, accent=OG, tall=True):
    """En-tete Constructys : bande degradee + cercles + dots + pill label."""
    hc = 5.5 if tall else 3.5
    band = add_rect(slide, Cm(0), Cm(0), W, Cm(hc), bg, 0.05)
    set_gradient(band, bg, _rgb_shift(bg, -25), 2700000)
    add_circle(slide, Cm(15.5), Cm(-2 if tall else -1), Cm(7), _rgb_shift(bg, -15), alpha=30)
    add_circle(slide, Cm(-1.5), Cm(hc - 2.5), Cm(4), _rgb_shift(bg, 20), alpha=20)
    add_circle(slide, Cm(17), Cm(hc - 2), Cm(3), _rgb_shift(bg, -25), alpha=25)
    add_dot_grid(slide, Cm(15), Cm(0.5 if tall else 0.3), 2, 3, WH, 15)
    if label:
        pw = max(Cm(3), Cm(len(label) * 0.18 + 1))
        pill = add_rect(slide, Cm(1.8), Cm(0.7 if tall else 0.5), pw, Cm(0.6), WH, 0.5)
        pill.text_frame.word_wrap = False
        p = pill.text_frame.paragraphs[0]
        p.text = label
        p.font.size = Pt(7)
        p.font.bold = True
        p.font.color.rgb = bg
        p.font.name = "Montserrat"
        p.alignment = PP_ALIGN.CENTER
    ty = Cm(1.8) if tall else Cm(1.2)
    tsz = 22 if tall else 18
    add_text(slide, Cm(1.8), ty, Cm(14), Cm(2.5), title, tsz, True, WH, font_name="Montserrat")
    by = Cm(hc + 0.3)
    add_rect(slide, Cm(1.8), by, Cm(3), Cm(0.12), accent)
    add_circle(slide, Cm(5.2), by - Cm(0.15), Cm(0.4), accent)

def set_gradient(shape, color1, color2, angle=2700000):
    """Degrade lineaire via XML. angle en 60000e de degre (0=droite, 5400000=bas)."""
    spPr = shape._element.spPr
    for old in spPr.findall(f'{{{_a}}}solidFill'):
        spPr.remove(old)
    for old in spPr.findall(f'{{{_a}}}gradFill'):
        spPr.remove(old)
    gf = etree.SubElement(spPr, f'{{{_a}}}gradFill')
    gsLst = etree.SubElement(gf, f'{{{_a}}}gsLst')
    gs1 = etree.SubElement(gsLst, f'{{{_a}}}gs', pos='0')
    etree.SubElement(gs1, f'{{{_a}}}srgbClr', val=str(color1))
    gs2 = etree.SubElement(gsLst, f'{{{_a}}}gs', pos='100000')
    etree.SubElement(gs2, f'{{{_a}}}srgbClr', val=str(color2))
    etree.SubElement(gf, f'{{{_a}}}lin', ang=str(angle), scaled='1')

def add_dot_grid(slide, x, y, rows=2, cols=3, color=WH, alpha=20):
    """Grille de points decoratifs Constructys."""
    sp = Cm(0.5)
    for r in range(rows):
        for c in range(cols):
            add_circle(slide, x + sp*c, y + sp*r, Cm(0.12), color, alpha=alpha)

def constructys_deco(slide, color=TQ):
    """Touches decoratives Constructys sur slides a en-tete leger."""
    add_circle(slide, Cm(18), Cm(-0.8), Cm(2.5), _rgb_shift(color, -15), alpha=25)
    add_circle(slide, Cm(19.5), Cm(1.5), Cm(1.8), _rgb_shift(color, 20), alpha=20)
    add_dot_grid(slide, Cm(17), Cm(0.6), 2, 3, _rgb_shift(color, -10), 18)

# ---- SLIDES ----
def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    blank = prs.slide_layouts[6]  # blank layout

    # ═══ S1 COVER ═══
    s = prs.slides.add_slide(blank)
    # Right panel turquoise with gradient
    rp = add_rect(s, Cm(12.5), Cm(0), Cm(8.5), H, TQ)
    set_gradient(rp, TQ, _rgb_shift(TQ, -30), 5400000)
    # Decorative transparent circles on TQ panel
    add_circle(s, Cm(13), Cm(2), Cm(5), RGBColor(0x4A,0xB0,0xA8), alpha=40)
    add_circle(s, Cm(16), Cm(8), Cm(3), RGBColor(0x7D,0xD8,0xD2), alpha=30)
    add_circle(s, Cm(14), Cm(16), Cm(6), RGBColor(0x4A,0xB0,0xA8), alpha=25)
    # Dot grid + accent stripe
    add_dot_grid(s, Cm(1.8), Cm(19.5), 2, 4, RGBColor(0xEE,0xEE,0xEE), 40)
    add_rect(s, Cm(0), Cm(20.5), Cm(12), Cm(0.08), OG)
    # Photo candidat on right panel
    p_img = pic(s, IMG_PHOTO, Cm(14), Cm(4), Cm(5.5), Cm(7))
    if p_img:
        add_shadow(p_img)
    else:
        add_rect(s, Cm(14), Cm(4), Cm(5.5), Cm(7), RGBColor(0x4A,0xB0,0xA8), 0.15)
        add_text(s, Cm(14), Cm(7), Cm(5.5), Cm(1), "Photo candidat", 9, False, WH, PP_ALIGN.CENTER)
    # Large orange block bottom right
    add_rect(s, Cm(14), Cm(24), Cm(7), Cm(5.7), OG, 0.15)
    # Navy block accent
    add_rect(s, Cm(17), Cm(21.5), Cm(4), Cm(2.8), NV, 0.1)
    # Badge SESSION
    add_badge(s, Cm(1.8), Cm(1.5), "SESSION 2026", OG, WH, Cm(3))
    add_text(s, Cm(5.2), Cm(1.5), Cm(5), Cm(0.7), "Academie de Lyon", 8, False, GY)
    # Label
    add_text(s, Cm(1.8), Cm(7), Cm(9), Cm(0.5), "DOSSIER DE SYNTHESE", 7, True, OG, font_name="Montserrat")
    # Title
    add_text(s, Cm(1.8), Cm(8), Cm(10), Cm(2), "RAPPORT\nD'ACTIVITES", 32, True, NV, font_name="Montserrat")
    add_text(s, Cm(1.8), Cm(11.5), Cm(10), Cm(1.2), "PROFESSIONNELLES", 32, True, OG, font_name="Montserrat")
    # Orange bar
    orange_bar(s, Cm(1.8), Cm(13.2), Cm(4))
    # Candidat
    add_text(s, Cm(1.8), Cm(22), Cm(5), Cm(0.4), "CANDIDAT", 7, True, OG, font_name="Montserrat")
    add_text(s, Cm(1.8), Cm(22.7), Cm(5), Cm(1), "BAHAFID", 24, True, NV, font_name="Montserrat")
    add_text(s, Cm(1.8), Cm(24.5), Cm(5), Cm(1), "Mohamed", 24, True, OG, font_name="Montserrat")
    add_text(s, Cm(1.8), Cm(26.2), Cm(5), Cm(0.5), f"N. {CANDIDAT['numero']}", 8, False, GY, font_name="Inter")
    add_badge(s, Cm(8), Cm(23.5), "BTS MEC", NV, WH, Cm(2.5))
    add_text(s, Cm(7), Cm(24.5), Cm(4.5), Cm(1), "Management Economique\nde la Construction", 8, False, GY, PP_ALIGN.CENTER, "Inter")

    # ═══ S2 FICHE CANDIDAT ═══
    s = prs.slides.add_slide(blank)
    constructys_header(s, "FICHE D'IDENTITE\nDU CANDIDAT", "IDENTITE", TQ)
    p2 = pic(s, IMG_PHOTO, Cm(16), Cm(1), Cm(3.2), Cm(4))
    if p2:
        add_shadow(p2)
    add_rich_text(s, Cm(1.8), Cm(7), Cm(17), Cm(1.3),
             "Un profil construit en trois phases : maitrise d'ouvrage publique (**Maroc**, "
             "estimation, commissions d'AO, suivi financier), execution terrain (**France**, "
             "chef d'equipe et de chantier gros oeuvre), puis formation certifiante BIM. "
             "Chaque phase a enrichi la suivante — voir le projet d'un autre angle change "
             "profondement la facon d'estimer.",
             8, TQ)
    # Table rows
    y = Cm(8.5)
    for i, (k, v) in enumerate(PAGE_02["champs"]):
        bg = LG if i % 2 == 0 else WH
        add_rect(s, Cm(1.8), y, Cm(17.4), Cm(1.2), bg, 0.08)
        add_text(s, Cm(2.2), y+Cm(0.15), Cm(4.5), Cm(0.9), k, 8, True, OG, font_name="Inter")
        add_text(s, Cm(7), y+Cm(0.15), Cm(11.5), Cm(0.9), v, 9, False, DK, font_name="Inter")
        y += Cm(1.3)
    # Stats row
    kpi_data = [("8 ans", "Experience BTP", TQ), ("2 pays", "Maroc + France", OG), ("7", "Marches publics", NV)]
    for i, (val, lab, clr) in enumerate(kpi_data):
        add_kpi_card(s, Cm(1.8) + Cm(6)*i, Cm(24), val, lab, clr, Cm(5.5))
    # Footer capsule
    add_rect(s, Cm(1.8), Cm(26.8), Cm(17.4), Cm(1), OG, 0.5, shadow=True)
    add_text(s, Cm(2.5), Cm(26.9), Cm(16), Cm(0.8), PAGE_02["pied"], 8, True, WH, PP_ALIGN.CENTER, "Inter")
    footer(s, 2)

    # ═══ S3 SOMMAIRE ═══
    s = prs.slides.add_slide(blank)
    bg3 = add_rect(s, Cm(0), Cm(0), W, H, TQ)
    set_gradient(bg3, TQ, _rgb_shift(TQ, -20), 5400000)
    add_circle(s, Cm(-2), Cm(20), Cm(8), _rgb_shift(TQ, -15), alpha=30)
    add_circle(s, Cm(16), Cm(-3), Cm(6), _rgb_shift(TQ, 15), alpha=25)
    add_circle(s, Cm(18), Cm(24), Cm(4), _rgb_shift(TQ, -10), alpha=20)
    add_dot_grid(s, Cm(15.5), Cm(2), 3, 4, WH, 15)
    add_rect(s, Cm(2.5), Cm(6), Cm(16), Cm(22), WH, 0.1, shadow=True)
    # White pill + title
    add_rect(s, Cm(2.5), Cm(2.5), Cm(2), Cm(0.8), WH, 0.5)
    add_text(s, Cm(5), Cm(1.8), Cm(12), Cm(2.5), "Sommaire", 30, True, WH, font_name="Montserrat")
    y = Cm(6.5)
    _sc = [OG, TQ, NV, OG]
    _si = 0  # index for numbered sections only
    for num, titre, desc, page in PAGE_03["sections"]:
        if num == "—":
            # Section non-numérotée : intro/conclusion — style sobre
            add_circle(s, Cm(3.2), y + Cm(0.08), Cm(0.65), GY, alpha=50)
            add_text(s, Cm(3.2), y+Cm(0.05), Cm(0.65), Cm(0.65), "–", 10, False, WH, PP_ALIGN.CENTER, "Montserrat")
            add_text(s, Cm(4.3), y+Cm(0.1), Cm(13), Cm(0.7), titre, 9, False, RGBColor(0x88,0x88,0x88), font_name="Inter")
        else:
            # Section numérotée : 01/02/03/04 correspondant aux slides séparateurs
            sc = _sc[_si % len(_sc)]
            _si += 1
            add_circle(s, Cm(3.2), y + Cm(0.05), Cm(0.75), sc)
            add_text(s, Cm(3.2), y+Cm(0.1), Cm(0.75), Cm(0.75), num, 8, True, WH, PP_ALIGN.CENTER, "Montserrat")
            add_text(s, Cm(4.3), y, Cm(13), Cm(0.85), titre, 11, True, NV, font_name="Montserrat")
        add_text(s, Cm(4.3), y+Cm(0.85), Cm(13), Cm(0.5), desc, 7, False, GY, font_name="Inter")
        add_text(s, Cm(4.3), y+Cm(1.45), Cm(13), Cm(0.45), page, 7, True, OG, font_name="Montserrat")
        add_rect(s, Cm(3.5), y+Cm(2.1), Cm(14), Cm(0.03), RGBColor(0xDD,0xDD,0xDD))
        y += Cm(2.4)
    footer(s, 3)

    # ═══ S4 INTRODUCTION ═══
    s = prs.slides.add_slide(blank)
    add_rect(s, Cm(0), Cm(0), W, Cm(0.4), TQ)
    constructys_deco(s, TQ)
    # Icon circle
    add_circle(s, Cm(1.5), Cm(1), Cm(1.2), OG)
    add_text(s, Cm(3), Cm(0.8), Cm(16), Cm(2), "INTRODUCTION\nET PARCOURS", 22, True, NV, font_name="Montserrat")
    orange_bar(s, Cm(1.8), Cm(3.5))
    add_text(s, Cm(1.8), Cm(4), Cm(17.4), Cm(1.8),
             PAGE_04["intro_texte"], 8, False, RGBColor(0x44,0x44,0x44), PP_ALIGN.JUSTIFY, "Inter")
    # Texte de synthese (deplace depuis la couverture)
    add_rich_text(s, Cm(1.8), Cm(6.0), Cm(17.4), Cm(2.0),
             "**8 annees** de BTP : 3 ans MOA publique au Maroc, 5 ans execution + BIM en France. "
             "**5 situations professionnelles** analysees selon la demarche CPAR. "
             "**7 marches publics** et **100 M DH** pilotes : estimation confidentielle, "
             "analyse des offres en CAO, suivi financier mensuel, ordres de service, decomptes definitifs. "
             "Sur le terrain : banches, beton, ferraillage, encadrement d'equipe — "
             "les **couts reels de production** vus de l'interieur.",
             7.5, OG)
    # Timeline left
    add_text(s, Cm(1.8), Cm(8.5), Cm(8), Cm(0.5), "MON PARCOURS", 7, True, OG, font_name="Montserrat")
    tl_colors = [TQ, NV, OG, TQ, OG]
    y = Cm(9.3)
    for i, (date, title, desc) in enumerate(PAGE_04["phases"]):
        c = tl_colors[i]
        add_circle(s, Cm(2), y, Cm(0.6), c)
        if i < len(PAGE_04["phases"])-1:
            add_rect(s, Cm(2.22), y+Cm(0.6), Cm(0.15), Cm(2.2), RGBColor(0xEE,0xEE,0xEE))
        add_text(s, Cm(3.2), y-Cm(0.1), Cm(6), Cm(0.4), date, 7, True, OG, font_name="Montserrat")
        add_text(s, Cm(3.2), y+Cm(0.35), Cm(6), Cm(0.5), title.split(" \u2013 ")[0] if " \u2013 " in title else title, 9, True, NV, font_name="Montserrat")
        add_text(s, Cm(3.2), y+Cm(0.85), Cm(6), Cm(0.5), desc, 7, False, GY, font_name="Inter")
        y += Cm(2.8)
    # KPIs right - progress rings as circles with text
    add_text(s, Cm(11), Cm(8.5), Cm(8), Cm(0.5), "CHIFFRES CLES", 7, True, OG, font_name="Montserrat")
    kpis = PAGE_04["chiffres_cles"]
    kpi_c = [NV, OG, TQ, PR]
    for i, (v, l) in enumerate(kpis):
        cx = Cm(11.5) + Cm(2.2)*i
        add_circle(s, cx, Cm(9.5), Cm(1.8), kpi_c[i], shadow=True)
        add_text(s, cx, Cm(9.8), Cm(1.8), Cm(0.8), v, 9, True, WH, PP_ALIGN.CENTER, "Montserrat")
        add_text(s, cx-Cm(0.2), Cm(11.5), Cm(2.2), Cm(0.5), l, 7, False, GY, PP_ALIGN.CENTER, "Inter")
    # Capsule CPAR
    add_rect(s, Cm(11), Cm(13), Cm(8.5), Cm(1.5), OG, 0.5, shadow=True)
    add_text(s, Cm(11.5), Cm(13.1), Cm(7.5), Cm(0.6), "5 SITUATIONS PROFESSIONNELLES", 8, True, WH, PP_ALIGN.CENTER, "Montserrat")
    add_text(s, Cm(11.5), Cm(13.75), Cm(7.5), Cm(0.55), "Contexte | Probleme | Action | Resultat", 7, False, RGBColor(0xFF,0xDD,0xAA), PP_ALIGN.CENTER, "Inter")
    # Methode CPAR reminder
    add_text(s, Cm(11), Cm(15.0), Cm(8.5), Cm(2.5),
             "Chaque situation est analysee selon la demarche CPAR : "
             "le Contexte pose la mission, le Probleme identifie l'obstacle, "
             "l'Action decrit la reponse apportee, le Resultat mesure l'impact reel.",
             7, False, RGBColor(0x55,0x55,0x55), PP_ALIGN.JUSTIFY, "Inter")
    # Project cards bottom
    add_text(s, Cm(1.8), Cm(24.0), Cm(17.4), Cm(0.5), "PROJETS ANALYSES", 7, True, OG, PP_ALIGN.CENTER, "Montserrat")
    proj_data = [("Projet 1", "4 communes - 53,5 M DH TTC - VRD", OG), ("Projet 2", "Lehri-Kerrouchen - 29 M DH - 25 km", TQ)]
    for i, (t, d, c) in enumerate(proj_data):
        x = Cm(1.8) + Cm(9)*i
        card = add_rect(s, x, Cm(24.8), Cm(8.5), Cm(2), WH, 0.1, shadow=True)
        add_rect(s, x, Cm(24.8), Cm(0.15), Cm(2), c)
        add_circle(s, x+Cm(0.5), Cm(25.2), Cm(1), c)
        add_text(s, x+Cm(1.8), Cm(25.1), Cm(6), Cm(0.6), t, 10, True, NV, font_name="Montserrat")
        add_text(s, x+Cm(1.8), Cm(25.8), Cm(6), Cm(0.6), d, 7, False, GY, font_name="Inter")
    footer(s, 4)

    # ═══ SEPARATOR helper ═══
    def sep_slide(numero, titre, sous_titre, bg=TQ):
        s = prs.slides.add_slide(blank)
        bg_r = add_rect(s, Cm(0), Cm(0), W, H, bg)
        set_gradient(bg_r, bg, _rgb_shift(bg, -25), 5400000)
        # Constructys decorative geometry
        add_circle(s, Cm(-3), Cm(5), Cm(10), _rgb_shift(bg, -20), alpha=40)
        add_circle(s, Cm(16), Cm(20), Cm(8), _rgb_shift(bg, 15), alpha=35)
        add_circle(s, Cm(14), Cm(-3), Cm(6), _rgb_shift(bg, -10), alpha=25)
        add_circle(s, Cm(-2), Cm(22), Cm(5), _rgb_shift(bg, 25), alpha=20)
        # Pill badge
        pw = Cm(3)
        pill = add_rect(s, Cm(9), Cm(9), pw, Cm(0.7), WH, 0.5)
        pill.text_frame.word_wrap = False
        p = pill.text_frame.paragraphs[0]
        p.text = numero
        p.font.size = Pt(9)
        p.font.bold = True
        p.font.color.rgb = bg
        p.font.name = "Montserrat"
        p.alignment = PP_ALIGN.CENTER
        add_text(s, Cm(1), Cm(10.5), W-Cm(2), Cm(3), numero, 72, True,
                 _rgb_shift(bg, 30), PP_ALIGN.CENTER, "Montserrat")
        t = titre.replace('\n', ' ')
        add_text(s, Cm(2), Cm(14), Cm(17), Cm(3), t, 26, True, WH, PP_ALIGN.CENTER, "Montserrat")
        orange_bar(s, Cm(8.5), Cm(17.5), Cm(4))
        add_text(s, Cm(2), Cm(18.5), Cm(17), Cm(2), sous_titre.replace('\n', ' '), 10, False,
                 _rgb_shift(bg, 60), PP_ALIGN.CENTER, "Inter")
        return s

    # ═══ S5 SEP ═══
    sep_slide(PAGE_05["numero"], PAGE_05["titre"], PAGE_05["sous_titre"])

    # ═══ S6 CONSEIL REGIONAL ═══
    s = prs.slides.add_slide(blank)
    constructys_header(s, PAGE_06["titre"].upper(), "STRUCTURE D'ACCUEIL", TQ)
    add_rich_text(s, Cm(1.8), Cm(6.2), Cm(10), Cm(3),
             "Collectivite territoriale couvrant **5 provinces** et **28 374 km2** pour **2,5 millions** "
             "d'habitants. L'ensemble des marches de travaux est regi par le **Decret n°2-12-349** "
             "du 20 mars 2013 relatif aux marches publics — appel d'offres ouvert au-dessus de 500 000 DH, "
             "pieces constitutives : CPS, RC, BPDE, plans et estimation confidentielle obligatoire. "
             "En tant que technicien de l'Agence d'Execution des Projets (dir. DOGHMANI), "
             "j'assurais la chaine technique complete — de l'estimation confidentielle avant AO "
             "jusqu'au decompte definitif — sur **7 marches publics** representant **100 M DH** : "
             "routes rurales, VRD, adduction d'eau potable, voirie urbaine et equipements collectifs.",
             8, NV)
    # Fact cards 2x2
    facts = [("5 Provinces", "28 374 km2", TQ), ("2,5 M", "Habitants", OG), ("Routes", "VRD, AEP", NV), ("7 Marches", "Publics", TQ)]
    for i, (v, l, c) in enumerate(facts):
        x = Cm(1.8) + Cm(5)*(i%2)
        yy = Cm(9.5) + Cm(2.5)*(i//2)
        add_rect(s, x, yy, Cm(4.5), Cm(2), WH, 0.12, shadow=True)
        add_circle(s, x+Cm(0.3), yy+Cm(0.4), Cm(0.9), c, shadow=True)
        add_text(s, x+Cm(1.5), yy+Cm(0.3), Cm(2.8), Cm(0.6), v, 9, True, NV, font_name="Montserrat")
        add_text(s, x+Cm(1.5), yy+Cm(1), Cm(2.8), Cm(0.5), l, 7, False, GY, font_name="Inter")
    # Organigramme right
    add_rect(s, Cm(12), Cm(6), Cm(7.5), Cm(9), LG, 0.1)
    add_text(s, Cm(12.3), Cm(6.5), Cm(7), Cm(0.6), "President du Conseil Regional", 8, False, GY, PP_ALIGN.CENTER, "Inter")
    add_rect(s, Cm(15.5), Cm(7.2), Cm(0.1), Cm(0.5), RGBColor(0xDD,0xDD,0xDD))
    add_text(s, Cm(12.3), Cm(7.8), Cm(7), Cm(0.6), "Directeur de l'Agence - M. DOGHMANI", 9, True, NV, PP_ALIGN.CENTER, "Montserrat")
    add_rect(s, Cm(15.5), Cm(8.5), Cm(0.1), Cm(0.5), RGBColor(0xDD,0xDD,0xDD))
    svcs = " | ".join(s.replace('\n',' ') for s in PAGE_06["organigramme"]["services"])
    add_text(s, Cm(12.3), Cm(9.1), Cm(7), Cm(0.5), svcs, 7, False, GY, PP_ALIGN.CENTER, "Inter")
    add_rect(s, Cm(13), Cm(10), Cm(5.5), Cm(1), OG, 0.5)
    add_text(s, Cm(13.2), Cm(10.1), Cm(5), Cm(0.8), "TECHNICIEN DE SUIVI\nBAHAFID Mohamed", 7, True, WH, PP_ALIGN.CENTER, "Montserrat")
    # Missions
    add_text(s, Cm(1.8), Cm(15), Cm(17.4), Cm(0.5), "MISSIONS PRINCIPALES", 7, True, OG, PP_ALIGN.CENTER, "Montserrat")
    mission_colors = [OG, TQ, NV, OG, TQ]
    y = Cm(16)
    for i, m in enumerate(PAGE_06["missions"]):
        x = Cm(1.8) + Cm(9.2)*(i%2)
        yy = y + Cm(2)*(i//2)
        add_circle(s, x, yy+Cm(0.1), Cm(0.8), mission_colors[i])
        add_text(s, x+Cm(1.2), yy, Cm(7.5), Cm(1.5), m, 8, False, RGBColor(0x55,0x55,0x55), font_name="Inter")
    footer(s, 6)

    # ═══ S7 BIMCO ═══
    s = prs.slides.add_slide(blank)
    add_rect(s, Cm(0), Cm(0), W, Cm(0.4), TQ)
    constructys_deco(s, TQ)
    # Logo BIMCO
    logo = pic(s, IMG_LOGO, Cm(1.2), Cm(0.8), Cm(1.5), Cm(1.5))
    if not logo:
        add_circle(s, Cm(1.5), Cm(1), Cm(1.2), OG)
    add_text(s, Cm(3), Cm(0.7), Cm(10), Cm(0.6), "BIMCO", 12, True, OG, font_name="Montserrat")
    add_text(s, Cm(3.2), Cm(1.2), Cm(16), Cm(1.5), "EXPERT BIM &\nECONOMIE DE LA CONSTRUCTION", 18, True, NV, font_name="Montserrat")
    orange_bar(s, Cm(1.8), Cm(3.5))
    add_rich_text(s, Cm(1.8), Cm(4), Cm(17.4), Cm(1.3),
             "Cree en **janvier 2026** a Bussieres (Loire), BIMCO occupe un creneau precis : "
             "appliquer le BIM aux metiers de l'economiste de la construction. "
             "La plupart des outils BIM sont concus pour les architectes et ingenieurs — "
             "pas pour le professionnel du chiffrage. BIMCO corrige ce manque : "
             "metres par extraction de maquette numerique, etudes de prix ancrees dans "
             "les couts reels, suivi economique en temps reel, et developpement de plugins "
             "Revit/Dynamo pour automatiser la chaine **metre → DPGF**. "
             "Il est le prolongement direct de **8 ans** d'experience MOA + execution + BIM.",
             8, OG)
    # Expertise left
    add_text(s, Cm(1.8), Cm(5.5), Cm(8), Cm(0.5), "EXPERTISE", 7, True, OG, PP_ALIGN.CENTER, "Montserrat")
    dom_icons = [TQ, OG, NV, TQ]
    for i, (d, desc) in enumerate(PAGE_07["domaines"][:4]):
        yy = Cm(6.5) + Cm(2)*i
        add_circle(s, Cm(2), yy, Cm(0.8), dom_icons[i])
        add_text(s, Cm(3.2), yy-Cm(0.1), Cm(6), Cm(0.5), d, 9, True, NV, font_name="Montserrat")
        add_text(s, Cm(3.2), yy+Cm(0.5), Cm(6), Cm(0.5), desc, 7, False, GY, font_name="Inter")
    # App capsule
    add_rect(s, Cm(1.8), Cm(15), Cm(8), Cm(2.5), LG, 0.15, shadow=True)
    add_text(s, Cm(2.3), Cm(15.3), Cm(7), Cm(0.6), PAGE_07["app"]["titre"], 9, True, NV, PP_ALIGN.CENTER, "Montserrat")
    add_text(s, Cm(2.3), Cm(16), Cm(7), Cm(0.5), PAGE_07["app"]["url"], 8, True, OG, PP_ALIGN.CENTER, "Inter")
    add_text(s, Cm(2.3), Cm(16.6), Cm(7), Cm(0.5), PAGE_07["app"]["stack"], 7, False, GY, PP_ALIGN.CENTER, "Inter")
    # Skills right - horizontal bars
    add_text(s, Cm(11), Cm(5.5), Cm(8), Cm(0.5), "COMPETENCES", 7, True, OG, PP_ALIGN.CENTER, "Montserrat")
    for i, (name, pct) in enumerate(PAGE_07["outils_principaux"]):
        yy = Cm(6.5) + Cm(1.5)*i
        add_text(s, Cm(11), yy, Cm(3.5), Cm(0.5), name, 8, True, NV, PP_ALIGN.RIGHT, "Inter")
        # Bar bg
        add_rect(s, Cm(15), yy+Cm(0.1), Cm(4), Cm(0.6), RGBColor(0xEE,0xEE,0xEE), 0.5)
        # Bar fill
        bar_c = TQ if i%2==0 else OG
        add_rect(s, Cm(15), yy+Cm(0.1), Cm(4*pct/100), Cm(0.6), bar_c, 0.5)
        add_text(s, Cm(15.2), yy+Cm(0.1), Cm(2), Cm(0.6), f"{pct}%", 7, True, WH, font_name="Montserrat")
    # Triple competence capsule
    add_rect(s, Cm(11), Cm(16), Cm(8.5), Cm(1.2), NV, 0.5, shadow=True)
    add_text(s, Cm(11.5), Cm(16.1), Cm(7.5), Cm(1), "TRIPLE COMPETENCE\nMOA + Execution + BIM", 8, True, WH, PP_ALIGN.CENTER, "Montserrat")
    # Bottom geometric shapes
    add_rect(s, Cm(18), Cm(24), Cm(3), Cm(5.7), TQ, 0.15)
    add_rect(s, Cm(17), Cm(27), Cm(2), Cm(2.7), OG, 0.1)
    footer(s, 7)

    # ═══ S8 SEP P1 ═══
    sep_slide(PAGE_08["numero"], PAGE_08["titre"], PAGE_08["sous_titre"])

    # ═══ S9 FICHE P1 ═══
    s = prs.slides.add_slide(blank)
    constructys_header(s, "MISE A NIVEAU DES\nCENTRES DE 4 COMMUNES", "PROJET 1", TQ)
    add_rich_text(s, Cm(1.8), Cm(6.2), Cm(10), Cm(2.2),
             "Programme d'amenagement urbain et VRD des centres emergents de **4 communes** "
             "de la province de Khenifra. Le marche unique couvre **8 corps d'etat** — "
             "de l'assainissement a l'eclairage public — sur quatre sites distincts distants "
             "de **20 a 80 km**. Montant : **53,5 M DH TTC**. Sur ce projet, j'ai etabli "
             "les avant-metres (**112 prix par commune**), realise l'estimation confidentielle "
             "de l'administration, participe a la commission d'appel d'offres et pilote "
             "le suivi financier mensuel sur **18 mois** avec un tableau de bord consolide.",
             8, OG)
    # Fiche table
    y = Cm(8.5)
    for k, v in PAGE_09["fiche"]:
        add_text(s, Cm(2.2), y, Cm(3.5), Cm(0.8), k, 8, True, OG, font_name="Inter")
        add_text(s, Cm(6), y, Cm(4.5), Cm(0.8), v.replace('\n', ', '), 8, False, DK, font_name="Inter")
        add_rect(s, Cm(2.2), y+Cm(0.9), Cm(8), Cm(0.03), RGBColor(0xEE,0xEE,0xEE))
        y += Cm(1.1)
    # Montant capsule right
    add_rect(s, Cm(12), Cm(6.5), Cm(7), Cm(5), OG, 0.15, shadow=True)
    add_text(s, Cm(12.5), Cm(7.5), Cm(6), Cm(2), PAGE_09["montant"], 28, True, WH, PP_ALIGN.CENTER, "Montserrat")
    add_text(s, Cm(12.5), Cm(10), Cm(6), Cm(1), PAGE_09["detail_montant"], 8, False, WH, PP_ALIGN.CENTER, "Inter")
    # Photo chantier reelle
    ch1 = pic(s, CHANTIER_IMGS[0], Cm(12), Cm(12), Cm(7), Cm(4))
    if ch1:
        add_shadow(ch1)
    else:
        add_rect(s, Cm(12), Cm(12), Cm(7), Cm(4), LG, 0.1)
        add_text(s, Cm(12), Cm(13.5), Cm(7), Cm(1), "Photo chantier urbain", 8, False, GY, PP_ALIGN.CENTER, "Inter")
    # Bleed shape
    add_rect(s, Cm(20), Cm(16), Cm(1), Cm(4), OG, 0.2)
    # Footer capsule
    add_rect(s, Cm(1.8), Cm(17.5), Cm(17.4), Cm(1), LG, 0.5)
    add_text(s, Cm(2.5), Cm(17.6), Cm(16), Cm(0.8), PAGE_09["pied"], 7, False, GY, PP_ALIGN.CENTER, "Inter")
    footer(s, 9)

    # ═══ S10 BUDGET ═══
    s = prs.slides.add_slide(blank)
    constructys_header(s, "BUDGET ET 8 CORPS D'ETAT", "BUDGET", OG, TQ, tall=False)
    add_rich_text(s, Cm(1.8), Cm(3.8), Cm(17.4), Cm(0.7),
             "Budget total de **44,6 M DH HT** reparti sur 4 communes selon leur taille et leur "
             "niveau d'equipement. Ouaoumana (**35,5%**) et Sebt Ait Rahou (**33,2%**) concentrent "
             "**68%** des fonds : ces deux communes presentaient les lineaires d'assainissement "
             "et de chaussee les plus importants, avec des surfaces de trottoirs jusqu'a "
             "**11 500 m2** par commune. Cette concentration a impose un suivi financier separe "
             "par commune, avec un tableau de bord hebdomadaire a trois indicateurs.",
             7, OG)
    # Donut chart left
    communes = PAGE_10["communes"]
    c10 = [OG, RGBColor(0xE0,0x78,0x00), TQ, NV]
    cats10 = [c[0] for c in communes]
    vals10 = [float(c[2].replace(",",".").replace("%","")) for c in communes]
    add_donut(s, Cm(1), Cm(4.5), Cm(8), Cm(7.5), cats10, vals10, c10, "44,6 M DH HT")
    # Budget bars right
    add_text(s, Cm(10), Cm(4.5), Cm(9), Cm(0.5), "REPARTITION PAR COMMUNE", 7, True, OG, font_name="Montserrat")
    for i, (name, montant, pct) in enumerate(communes):
        yy = Cm(5.5) + Cm(1.8)*i
        p_val = float(pct.replace(",",".").replace("%",""))
        add_text(s, Cm(10), yy+Cm(0.1), Cm(3.5), Cm(0.6), name, 8, True, NV, PP_ALIGN.RIGHT, "Inter")
        add_rect(s, Cm(13.8), yy+Cm(0.15), Cm(5.5), Cm(0.7), RGBColor(0xF0,0xF0,0xF0), 0.5)
        add_rect(s, Cm(13.8), yy+Cm(0.15), Cm(5.5*p_val/40), Cm(0.7), c10[i], 0.5)
        add_text(s, Cm(14), yy+Cm(0.15), Cm(2), Cm(0.7), f"{p_val:.0f}%", 7, True, WH, font_name="Montserrat")
    # Capsule
    add_rect(s, Cm(1.8), Cm(13), Cm(17.4), Cm(1), OG, 0.5, shadow=True)
    add_text(s, Cm(2.5), Cm(13.1), Cm(16), Cm(0.8), "68% du budget sur 2 communes : Ouaoumana + Sebt Ait Rahou", 8, True, WH, PP_ALIGN.CENTER, "Montserrat")
    # Corps d'etat — tableau stylise
    add_text(s, Cm(1.8), Cm(13.8), Cm(17.4), Cm(0.6),
             "DETAIL DES 8 CORPS D'ETAT DU MARCHE", 8, True, NV, PP_ALIGN.CENTER, "Montserrat")
    orange_bar(s, Cm(8.6), Cm(14.5), Cm(3.8))
    draw_table(s, Cm(1.8), Cm(14.8),
               [(Cm(1.6), PP_ALIGN.CENTER), (Cm(5.6), PP_ALIGN.LEFT), (Cm(10.2), PP_ALIGN.LEFT)],
               ["Partie", "Designation", "Exemple de travaux"],
               TABLE_CORPS_ETAT["lignes"])
    footer(s, 10)

    # ═══ CPAR helper ═══
    def cpar_slide(sit, num):
        s = prs.slides.add_slide(blank)
        nv_band = add_rect(s, Cm(0), Cm(0), W, Cm(5.5), NV)
        set_gradient(nv_band, NV, _rgb_shift(NV, -20), 2700000)
        # Constructys decorative circles
        add_circle(s, Cm(15), Cm(-2), Cm(6), _rgb_shift(NV, 15), alpha=25)
        add_circle(s, Cm(-2), Cm(3), Cm(5), _rgb_shift(NV, 10), alpha=20)
        add_text(s, Cm(1.5), Cm(0.5), Cm(2), Cm(3), str(sit["numero"]), 48, True,
                 RGBColor(0x2A,0x4A,0x70), font_name="Montserrat")
        add_text(s, Cm(3.5), Cm(0.8), Cm(9), Cm(0.5), f"SITUATION {sit['numero']}", 7, True, TQ, font_name="Montserrat")
        add_text(s, Cm(3.5), Cm(1.5), Cm(9), Cm(2), sit["titre"], 16, True, WH, font_name="Montserrat")
        # KPI box
        add_rect(s, Cm(14), Cm(1), Cm(5.5), Cm(3), RGBColor(0x2A,0x4A,0x70), 0.12, shadow=True)
        add_text(s, Cm(14), Cm(1.5), Cm(5.5), Cm(1.2), sit["chiffre_cle"], 18, True, OG, PP_ALIGN.CENTER, "Montserrat")
        add_text(s, Cm(14), Cm(2.8), Cm(5.5), Cm(0.6), sit["chiffre_label"], 7, False, RGBColor(0xAA,0xAA,0xBB), PP_ALIGN.CENTER, "Inter")
        # CPAR blocks
        blocks = [("C","Contexte",sit["contexte"],RGBColor(0x99,0x99,0x99)),
                  ("P","Problematique",sit["probleme"],OG),
                  ("A","Action",sit["action"],TQ),
                  ("R","Resultat",sit["resultat"],NV)]
        y = Cm(6.5)
        for lt, label, text, color in blocks:
            add_rect(s, Cm(1.8), y, Cm(17.4), Cm(4.2), WH, 0.1, shadow=True)
            add_rect(s, Cm(1.8), y, Cm(0.15), Cm(4.2), color)
            add_circle(s, Cm(2.5), y+Cm(0.4), Cm(1.2), color)
            add_text(s, Cm(2.5), y+Cm(0.5), Cm(1.2), Cm(1), lt, 14, True, WH, PP_ALIGN.CENTER, "Montserrat")
            add_text(s, Cm(4.2), y+Cm(0.4), Cm(5), Cm(0.5), label.upper(), 8, True, color, font_name="Montserrat")
            add_rect(s, Cm(4.2), y+Cm(1), Cm(2.5), Cm(0.06), color)  # accent line
            add_text(s, Cm(4.2), y+Cm(1.2), Cm(14), Cm(2.7), text, 9, False, RGBColor(0x55,0x55,0x55), PP_ALIGN.JUSTIFY, "Inter")
            y += Cm(4.6)
        # Footer competence
        comp_bar = add_rect(s, Cm(0), Cm(25.5), W, Cm(1.2), LG)
        add_shadow(comp_bar, blur=2, dist=1, alpha=10)
        add_rect(s, Cm(0), Cm(25.5), W, Cm(0.08), OG)
        add_circle(s, Cm(1.2), Cm(25.65), Cm(0.8), OG)
        add_text(s, Cm(1.2), Cm(25.7), Cm(0.8), Cm(0.7), "C", 10, True, WH, PP_ALIGN.CENTER, "Montserrat")
        add_text(s, Cm(2.3), Cm(25.65), Cm(17), Cm(0.9), sit["competence"].upper(), 7, True, NV, PP_ALIGN.CENTER, "Montserrat")
        footer(s, num)

    # ═══ S11-14 CPAR ═══
    cpar_slide(SITUATION_1, 11)
    cpar_slide(SITUATION_2, 12)
    cpar_slide(SITUATION_3, 13)
    cpar_slide(SITUATION_4, 14)

    # ═══ S15 DEFIS P1 ═══
    s = prs.slides.add_slide(blank)
    add_rect(s, Cm(0), Cm(0), W, Cm(0.4), OG)
    constructys_deco(s, OG)
    add_text(s, Cm(1.8), Cm(1.5), Cm(17), Cm(1.5), "DIFFICULTES ET\nSOLUTIONS - PROJET 1", 20, True, NV, font_name="Montserrat")
    orange_bar(s, Cm(1.8), Cm(4))
    add_rich_text(s, Cm(1.8), Cm(4.5), Cm(17.4), Cm(1.8),
             "**4 chantiers simultanes**, **8 corps d'etat**, des sites distants de **20 a 80 km** "
             "sur la province de Khenifra. Des conditions geologiques imprevues a Kerrouchen "
             "ont genere un surcout de **+12%** sur la chaussee, et des ecarts de quantites de "
             "bordures ont atteint **+15%** a Sebt Ait Rahou. Une gestion rigoureuse etait "
             "indispensable pour contenir les depassements sans recourir a un avenant.",
             8, OG)
    defi_colors = [OG, TQ, PR, NV]
    for i, (df, pb, sl) in enumerate(PAGE_15["defis"]):
        x = Cm(1.2) + Cm(4.8)*i
        c = defi_colors[i]
        add_rect(s, x, Cm(6.5), Cm(4.4), Cm(14), WH, 0.1, shadow=True)
        add_rect(s, x, Cm(6.5), Cm(4.4), Cm(0.15), c)
        add_circle(s, x+Cm(0.3), Cm(7), Cm(0.8), c)
        add_badge(s, x+Cm(1.3), Cm(7.1), df.split('\n')[0][:18], c, WH, Cm(3))
        add_text(s, x+Cm(0.3), Cm(8.2), Cm(3.8), Cm(0.6), df, 9, True, c, font_name="Montserrat")
        pb_text = pb.replace('\n', ' ')
        add_text(s, x+Cm(0.3), Cm(9.2), Cm(3.8), Cm(5), pb_text, 7, False, RGBColor(0x55,0x55,0x55), PP_ALIGN.JUSTIFY, "Inter")
        add_rect(s, x+Cm(0.5), Cm(14), Cm(3.4), Cm(0.03), RGBColor(0xEE,0xEE,0xEE))
        sl_text = sl.replace('\n', ' ')
        add_text(s, x+Cm(0.3), Cm(14.5), Cm(3.8), Cm(4), sl_text, 7, False, c, PP_ALIGN.JUSTIFY, "Inter")
    # Bottom capsule
    cap15 = add_rect(s, Cm(1.8), Cm(21), Cm(17.4), Cm(1.2), NV, 0.5, shadow=True)
    set_gradient(cap15, NV, _rgb_shift(NV, -15), 0)
    add_text(s, Cm(2.5), Cm(21.1), Cm(16), Cm(1), "Resultat : depassement contenu a +0,8% sur Kerrouchen, aucun avenant necessaire sur les 4 communes.", 8, False, WH, PP_ALIGN.CENTER, "Inter")
    footer(s, 15)

    # ═══ S16 SEP P2 ═══
    sep_slide(PAGE_16["numero"], PAGE_16["titre"], PAGE_16["sous_titre"])

    # ═══ S17 FICHE P2 ═══
    s = prs.slides.add_slide(blank)
    constructys_header(s, "ROUTE LEHRI-KERROUCHEN\n25 KM", "PROJET 2", TQ)
    add_rich_text(s, Cm(1.8), Cm(6.2), Cm(10), Cm(1.8),
             "Route rurale strategique de **25 km** desenclavant les communes du Moyen Atlas, "
             "avec un denivele cumule de **400 m** et des pentes atteignant **12%** en lacets. "
             "Le marche se decompose en **53 prix** sur trois sections : lineaire principal "
             "(23 prix), carrefour (11 prix) et bretelles (19 prix). Sur ce projet, j'ai "
             "etabli les avant-metres, suivi l'execution et controle les cubatures. "
             "La decouverte de calcaire fracture a impose la reclassification de **5 000 m3** "
             "et l'ajout de **3 dalots** non prevus au marche initial.",
             8, OG)
    y = Cm(8)
    for k, v in PAGE_17["fiche"]:
        add_text(s, Cm(2.2), y, Cm(3.5), Cm(0.8), k, 8, True, OG, font_name="Inter")
        add_text(s, Cm(6), y, Cm(4.5), Cm(0.8), v, 8, False, DK, font_name="Inter")
        y += Cm(1.2)
    # Montant
    add_rect(s, Cm(12), Cm(6.5), Cm(7), Cm(4), OG, 0.15, shadow=True)
    add_text(s, Cm(12.5), Cm(7.3), Cm(6), Cm(1.5), PAGE_17["montant"], 28, True, WH, PP_ALIGN.CENTER, "Montserrat")
    add_text(s, Cm(12.5), Cm(9.2), Cm(6), Cm(0.7), "Programme PRR3", 9, False, WH, PP_ALIGN.CENTER, "Inter")
    # Photo route
    ch2 = pic(s, CHANTIER_IMGS[1], Cm(12), Cm(11), Cm(7), Cm(3.5))
    if ch2:
        add_shadow(ch2)
    # Metres
    add_text(s, Cm(1.8), Cm(13.5), Cm(17.4), Cm(0.5), "PRINCIPAUX METRES", 9, True, OG, font_name="Montserrat")
    for i, (v, l) in enumerate(PAGE_17["metres"]):
        x = Cm(1.8) + Cm(3.1)*i
        add_rect(s, x, Cm(14.5), Cm(2.8), Cm(2.2), WH, 0.1)
        add_text(s, x, Cm(14.7), Cm(2.8), Cm(0.8), v, 10, True, NV, PP_ALIGN.CENTER, "Montserrat")
        add_text(s, x, Cm(15.7), Cm(2.8), Cm(0.6), l, 7, False, OG, PP_ALIGN.CENTER, "Inter")
    add_rect(s, Cm(20), Cm(13), Cm(1), Cm(4), OG, 0.2)
    footer(s, 17)

    # ═══ S18 BUDGET ROUTE ═══
    s = prs.slides.add_slide(blank)
    constructys_header(s, "REPARTITION BUDGETAIRE\nPROJET 2", "BUDGET", OG, TQ, tall=False)
    add_rich_text(s, Cm(1.8), Cm(3.8), Cm(17.4), Cm(0.7),
             "Budget de **24,2 M DH HT** — soit **29 M DH TTC** avec la TVA. Deux postes dominent : "
             "le corps de chaussee (**30,1%** — GNB + GNF2 sur 25 km a des couts unitaires eleves "
             "en montagne) et les ouvrages hydrauliques (**28,1%** — **3 talwegs** necessitant "
             "des dalots 3x2x2 m et **794 ml** de buses O1000). Ce ratio drainage/chaussee "
             "depasse de **2,5x** celui d'une route en plaine : c'est la signature budgetaire "
             "d'un chantier montagneux au Moyen Atlas.",
             7, OG)
    # Donut chart left
    items = PAGE_18["items"]
    c18 = [OG, RGBColor(0xE0,0x78,0x00), RGBColor(0xC0,0x68,0x00), TQ, RGBColor(0x4A,0xB0,0xA8), NV]
    cats18 = [it[0] for it in items]
    vals18 = [float(it[2].replace(",",".").replace("%","")) for it in items]
    add_donut(s, Cm(1), Cm(4.5), Cm(8), Cm(7.5), cats18, vals18, c18, "24,2 M DH HT")
    # Bars right
    add_text(s, Cm(10), Cm(4.5), Cm(9), Cm(0.5), "REPARTITION PAR POSTE", 7, True, OG, font_name="Montserrat")
    for i, (name, montant, pct) in enumerate(items):
        yy = Cm(5.5) + Cm(1.5)*i
        p_val = float(pct.replace(",",".").replace("%",""))
        add_text(s, Cm(10), yy+Cm(0.1), Cm(3.5), Cm(0.5), name, 7, True, NV, PP_ALIGN.RIGHT, "Inter")
        add_rect(s, Cm(13.8), yy+Cm(0.1), Cm(5.5), Cm(0.6), RGBColor(0xF0,0xF0,0xF0), 0.5)
        add_rect(s, Cm(13.8), yy+Cm(0.1), Cm(5.5*p_val/35), Cm(0.6), c18[i], 0.5)
        add_text(s, Cm(14), yy+Cm(0.1), Cm(2), Cm(0.6), f"{p_val:.0f}%", 7, True, WH, font_name="Montserrat")
    # Capsule callout
    add_rect(s, Cm(1.8), Cm(13.5), Cm(17.4), Cm(1), OG, 0.5, shadow=True)
    add_text(s, Cm(2.5), Cm(13.6), Cm(16), Cm(0.8), PAGE_18["callout"].replace('\n',' | '), 8, True, WH, PP_ALIGN.CENTER, "Montserrat")
    # Budget route — tableau stylise par poste
    add_text(s, Cm(1.8), Cm(15.0), Cm(17.4), Cm(0.5),
             "DETAIL BUDGETAIRE PAR POSTE", 8, True, OG, PP_ALIGN.CENTER, "Montserrat")
    orange_bar(s, Cm(7.9), Cm(15.6), Cm(5.2))
    draw_table(s, Cm(1.8), Cm(15.9),
               [(Cm(9.5), PP_ALIGN.LEFT), (Cm(4.5), PP_ALIGN.RIGHT), (Cm(3.4), PP_ALIGN.CENTER)],
               ["Poste", "Montant HT (DH)", "%"],
               TABLE_BUDGET_ROUTE["lignes"])
    footer(s, 18)

    # ═══ S19 SITUATION 5 ═══
    cpar_slide(SITUATION_5, 19)

    # ═══ S20 DEFIS ROUTE ═══
    s = prs.slides.add_slide(blank)
    add_rect(s, Cm(0), Cm(0), W, Cm(0.4), TQ)
    constructys_deco(s, TQ)
    add_text(s, Cm(1.8), Cm(1.5), Cm(17), Cm(1.5), "DEFIS D'UN CHANTIER ROUTIER\nEN ZONE MONTAGNEUSE", 18, True, NV, font_name="Montserrat")
    orange_bar(s, Cm(1.8), Cm(4))
    add_rich_text(s, Cm(1.8), Cm(4.3), Cm(17.4), Cm(1),
             "**25 km** en zone montagneuse du Moyen Atlas : denivele important, pentes raides, "
             "gel hivernal et precipitations abondantes. Les etudes prealables n'avaient pas "
             "identifie la nature reelle du sous-sol sur une partie du trace. Ces aleas "
             "ont forge ma rigueur dans le controle contradictoire des quantites.",
             8, TQ)
    defi2_colors = [TQ, OG, PR, NV]
    for i, (df, pb, sl) in enumerate(PAGE_20["defis"]):
        x = Cm(1.2) + Cm(4.8)*i
        c = defi2_colors[i]
        add_rect(s, x, Cm(5.5), Cm(4.4), Cm(14), WH, 0.1, shadow=True)
        add_rect(s, x, Cm(5.5), Cm(4.4), Cm(0.15), c)
        add_circle(s, x+Cm(0.3), Cm(6), Cm(0.8), c)
        add_text(s, x+Cm(0.3), Cm(7), Cm(3.8), Cm(0.6), df, 9, True, c, font_name="Montserrat")
        add_text(s, x+Cm(0.3), Cm(7.8), Cm(3.8), Cm(5), pb.replace('\n', ' '), 7, False, RGBColor(0x55,0x55,0x55), PP_ALIGN.JUSTIFY, "Inter")
        add_rect(s, x+Cm(0.5), Cm(13), Cm(3.4), Cm(0.03), RGBColor(0xEE,0xEE,0xEE))
        add_text(s, x+Cm(0.3), Cm(13.5), Cm(3.8), Cm(4), sl.replace('\n', ' '), 7, False, c, PP_ALIGN.JUSTIFY, "Inter")
    cap20 = add_rect(s, Cm(1.8), Cm(20), Cm(17.4), Cm(1.5), TQ, 0.5, shadow=True)
    set_gradient(cap20, TQ, _rgb_shift(TQ, -20), 0)
    add_text(s, Cm(2.5), Cm(20.2), Cm(16), Cm(1), "Bilan : surcoat de 285 000 DH absorbe par compensation, 3 dalots ajoutes, reception dans les delais.", 8, False, WH, PP_ALIGN.CENTER, "Inter")
    footer(s, 20)

    # ═══ S21 ACTIVITES COMPL ═══
    s = prs.slides.add_slide(blank)
    add_rect(s, Cm(0), Cm(0), W, Cm(0.4), TQ)
    constructys_deco(s, TQ)
    add_circle(s, Cm(1.2), Cm(0.8), Cm(1.2), OG)
    add_text(s, Cm(2.8), Cm(0.8), Cm(16), Cm(1.5), "ACTIVITES\nCOMPLEMENTAIRES", 18, True, NV, font_name="Montserrat")
    orange_bar(s, Cm(1.8), Cm(3.2))
    # Maroc left
    add_circle(s, Cm(1.8), Cm(4), Cm(0.8), OG)
    add_text(s, Cm(3), Cm(3.8), Cm(7), Cm(1), "5 AUTRES MARCHES AU\nCONSEIL REGIONAL", 9, True, OG, font_name="Montserrat")
    add_text(s, Cm(1.8), Cm(5.5), Cm(9), Cm(1.3),
             "En parallele des deux projets principaux, j'ai suivi cinq autres marches publics "
             "couvrant des types d'infrastructures differents : routes rurales, pistes, "
             "adduction d'eau potable et voirie urbaine. Ces missions complementaires "
             "m'ont permis de maitriser les specificites de chaque type de marche "
             "— procedures, pieces, suivi terrain — et de construire une vision "
             "transversale de la maitrise d'ouvrage publique en milieu rural.",
             8, False, RGBColor(0x55,0x55,0x55), PP_ALIGN.JUSTIFY, "Inter")
    draw_table(s, Cm(1.8), Cm(7),
               [(Cm(2.2), PP_ALIGN.CENTER), (Cm(4.5), PP_ALIGN.LEFT), (Cm(2.8), PP_ALIGN.LEFT)],
               ["N° Marche", "Objet", "Type / Montant"],
               TABLE_AUTRES_MARCHES["lignes"],
               header_fill=OG)
    # France right
    add_circle(s, Cm(11.5), Cm(4), Cm(0.8), TQ)
    add_text(s, Cm(12.7), Cm(3.8), Cm(7), Cm(1), "EXPERIENCE TERRAIN\nEN FRANCE", 9, True, TQ, font_name="Montserrat")
    add_text(s, Cm(11.5), Cm(5.5), Cm(8), Cm(1.3),
             "Chef d'equipe GO chez Ergalis (Feurs, Loire) puis chef de chantier chez Minssieux et Fils "
             "(Mornant, Rhone). Cette experience terrain m'a appris les couts reels de production — "
             "main-d'oeuvre, rendements, consommation materiaux, couts matériel — "
             "ainsi que les contraintes d'execution qui impactent directement les prix : "
             "acces chantier, intemperies, coordination entre lots, approvisionnements. "
             "Connaitre le chantier de l'interieur change profondement la facon d'estimer.",
             8, False, RGBColor(0x55,0x55,0x55), PP_ALIGN.JUSTIFY, "Inter")
    y = Cm(7)
    for d, p, dt in PAGE_21["france"]:
        add_circle(s, Cm(11.8), y+Cm(0.1), Cm(0.4), TQ)
        add_text(s, Cm(12.6), y, Cm(7), Cm(1.2), f"{d} - {p} : {dt}", 7, False, RGBColor(0x55,0x55,0x55), font_name="Inter")
        y += Cm(1.5)
    # Capsule
    cap21 = add_rect(s, Cm(1.8), Cm(14), Cm(17.4), Cm(1), OG, 0.5, shadow=True)
    set_gradient(cap21, OG, _rgb_shift(OG, -25), 0)
    add_text(s, Cm(2.5), Cm(14.1), Cm(16), Cm(0.8), "4 types d'infrastructures maitrisees : routes, VRD, adduction d'eau potable, voirie urbaine", 8, True, WH, PP_ALIGN.CENTER, "Montserrat")
    footer(s, 21)

    # ═══ S22 SEP BILAN ═══
    sep_slide(PAGE_22["numero"], PAGE_22["titre"], PAGE_22["sous_titre"], NV)

    # ═══ S23 COMPETENCES ═══
    s = prs.slides.add_slide(blank)
    add_rect(s, Cm(0), Cm(0), W, Cm(0.4), OG)
    constructys_deco(s, OG)
    add_circle(s, Cm(1.2), Cm(0.8), Cm(1.2), OG)
    add_text(s, Cm(2.8), Cm(0.8), Cm(16), Cm(1.5), "SYNTHESE DES ACTIVITES\nPROFESSIONNELLES", 18, True, NV, font_name="Montserrat")
    orange_bar(s, Cm(1.8), Cm(3.2))
    add_text(s, Cm(1.8), Cm(3.5), Cm(17.4), Cm(0.9),
             "Ce tableau synthetise les activites realisees sur les deux projets et les cinq situations CPAR, "
             "en les croisant avec les competences du BTS MEC mobilisees. "
             "Chaque ligne renvoie a une situation concrete vecue sous contrainte reelle — delai, budget, alea. "
             "La combinaison triple MOA + Execution + BIM permet d'etablir des estimations realistes, "
             "d'analyser les offres avec pertinence et d'exploiter les maquettes pour des metres precis. "
             "Niveau Maitrise : pratique reguliere et autonome. Niveau Expert : capacite a transmettre.",
             7, False, RGBColor(0x55,0x55,0x55), PP_ALIGN.JUSTIFY, "Inter")
    # Table header
    cols = [("ACTIVITE REALISEE", Cm(5.5)), ("SOUS-COMPETENCE", Cm(3.5)), ("SITUATION", Cm(3)), ("NIVEAU", Cm(4))]
    x = Cm(1.8)
    add_rect(s, x, Cm(4.5), Cm(17.4), Cm(1), NV, 0.08, shadow=True)
    cx = x
    for label, w in cols:
        add_text(s, cx+Cm(0.2), Cm(4.6), w, Cm(0.8), label, 6, True, WH, PP_ALIGN.CENTER, "Montserrat")
        cx += w + Cm(0.3)
    # Table rows
    y = Cm(5.7)
    for act, sc, sit, niv in PAGE_23["tableau"]:
        bg = LG if PAGE_23["tableau"].index((act,sc,sit,niv)) % 2 == 0 else WH
        add_rect(s, Cm(1.8), y, Cm(17.4), Cm(1.5), bg, 0.02)
        cx = Cm(2)
        add_text(s, cx, y+Cm(0.2), Cm(5.5), Cm(1.1), act, 7, False, DK, font_name="Inter")
        cx += Cm(5.8)
        add_text(s, cx, y+Cm(0.2), Cm(3.5), Cm(1.1), sc, 7, False, DK, PP_ALIGN.CENTER, "Inter")
        cx += Cm(3.8)
        add_badge(s, cx, y+Cm(0.4), sit, OG, WH, Cm(2.8))
        cx += Cm(3.3)
        # Progress bar
        bc = NV if niv == "Expert" else TQ
        bp = 100 if niv == "Expert" else 75
        add_rect(s, cx, y+Cm(0.4), Cm(2.2), Cm(0.5), RGBColor(0xEE,0xEE,0xEE), 0.5)
        add_rect(s, cx, y+Cm(0.4), Cm(2.2*bp/100), Cm(0.5), bc, 0.5)
        add_text(s, cx+Cm(2.5), y+Cm(0.3), Cm(1.5), Cm(0.6), niv, 7, True, bc, font_name="Inter")
        y += Cm(1.6)
    add_text(s, Cm(1.8), y+Cm(0.3), Cm(17.4), Cm(0.6), PAGE_23["legende"], 7, False, GY, font_name="Inter")
    footer(s, 23)

    # ═══ S24 COMPARAISON ═══
    s = prs.slides.add_slide(blank)
    _MA = RGBColor(0xC8, 0x10, 0x2E)  # Rouge Maroc
    _FR = RGBColor(0x00, 0x2F, 0x6C)  # Bleu France
    add_rect(s, Cm(0), Cm(0), W, Cm(0.4), OG)
    constructys_deco(s, OG)
    add_circle(s, Cm(1.2), Cm(0.8), Cm(1), OG)
    add_text(s, Cm(1.8), Cm(0.5), Cm(10), Cm(0.5), "ANALYSE COMPARATIVE", 7, True, OG, font_name="Montserrat")
    add_text(s, Cm(1.8), Cm(1.2), Cm(15), Cm(1.5), "MAROC vs FRANCE — CADRES REGLEMENTAIRES", 18, True, NV, font_name="Montserrat")
    orange_bar(s, Cm(1.8), Cm(3.2))
    add_text(s, Cm(1.8), Cm(3.6), Cm(17.4), Cm(1),
             PAGE_24["intro"],
             7.5, False, RGBColor(0x55,0x55,0x55), PP_ALIGN.JUSTIFY, "Inter")
    # Colonnes MAROC / FRANCE — bandeaux de couleur
    add_rect(s, Cm(7.4), Cm(4.8), Cm(6.8), Cm(0.65), _MA, 0.05)
    add_text(s, Cm(7.4), Cm(4.88), Cm(6.8), Cm(0.5), "MAROC", 9, True, WH, PP_ALIGN.CENTER, "Montserrat")
    add_rect(s, Cm(14.4), Cm(4.8), Cm(4.8), Cm(0.65), _FR, 0.05)
    add_text(s, Cm(14.4), Cm(4.88), Cm(4.8), Cm(0.5), "FRANCE", 9, True, WH, PP_ALIGN.CENTER, "Montserrat")
    # Tableau 8 aspects
    bot = draw_table(s, Cm(1.8), Cm(5.45),
                     [(Cm(5.6), PP_ALIGN.LEFT), (Cm(6.8), PP_ALIGN.LEFT), (Cm(5.0), PP_ALIGN.LEFT)],
                     ["Aspect", "Maroc", "France"],
                     TABLE_COMPARAISON_REG["lignes"],
                     header_fill=RGBColor(0x35,0x35,0x55))
    # Synthese capsule — degradé Maroc→France
    syn = add_rect(s, Cm(1.8), bot + Cm(0.5), Cm(17.4), Cm(1.5), NV, 0.5, shadow=True)
    set_gradient(syn, _MA, _FR, 0)
    add_text(s, Cm(2.5), bot + Cm(0.65), Cm(16), Cm(1.1),
             PAGE_24["synthese"].replace('\n', ' | '), 8, False, WH, PP_ALIGN.CENTER, "Inter")
    # Apports de cette double experience
    y_ap = bot + Cm(2.5)
    add_text(s, Cm(1.8), y_ap, Cm(17.4), Cm(0.55),
             "CE QUE CETTE DOUBLE EXPERIENCE M'APPORTE", 8, True, NV, PP_ALIGN.CENTER, "Montserrat")
    orange_bar(s, Cm(7.5), y_ap + Cm(0.6), Cm(6))
    apports = [
        (_MA, "Cote maitrise d'ouvrage — Maroc (2017-2022)",
         "Concevoir les marches, rediger les CPS/RC/BPDE, piloter la CAO, verifier les situations mensuelles, "
         "emettre les OS. Comprendre les attentes du maitre d'ouvrage en matiere de justification des prix "
         "et savoir analyser les offres avec pertinence — en identifiant les prix anormalement bas ou excessifs. "
         "7 marches publics, 100 M DH d'investissements pilotes."),
        (_FR, "Cote execution — France (2022-2024)",
         "Banches, ferraillage, betonnage, encadrement d'equipe, planning quotidien, controle qualite. "
         "Comprendre les couts reels de production (main-d'oeuvre, rendements, materiaux, matériel) "
         "et les contraintes d'execution : acces chantier, intemperies, approvisionnements. "
         "Cette connaissance terrain rend les estimations ancrees dans la realite — pas dans les baremes."),
        (TQ, "Synthese : la triple competence MOA + Execution + BIM",
         "Peu de professionnels combinent ces trois dimensions. Elle permet d'etablir des estimations "
         "realistes (ni trop basses, risque d'infructuosite — ni trop hautes, surcoût pour la collectivite), "
         "d'analyser les offres avec perspicacite, et d'exploiter les maquettes BIM pour des metres precis. "
         "C'est ce positionnement unique que traduit BIMCO."),
    ]
    y_ap += Cm(1.3)
    for i, (col, titre, texte) in enumerate(apports):
        add_rect(s, Cm(1.8), y_ap, Cm(17.4), Cm(2.6), WH, 0.08, shadow=True)
        add_rect(s, Cm(1.8), y_ap, Cm(0.15), Cm(2.6), col)
        add_circle(s, Cm(2.5), y_ap + Cm(0.4), Cm(0.9), col)
        add_text(s, Cm(2.5), y_ap + Cm(0.42), Cm(0.9), Cm(0.7), str(i+1), 11, True, WH, PP_ALIGN.CENTER, "Montserrat")
        add_text(s, Cm(3.8), y_ap + Cm(0.2), Cm(14.5), Cm(0.55), titre, 8, True, col, font_name="Montserrat")
        add_text(s, Cm(3.8), y_ap + Cm(0.85), Cm(14.5), Cm(1.6), texte, 7, False, RGBColor(0x44,0x44,0x44), PP_ALIGN.JUSTIFY, "Inter")
        y_ap += Cm(2.8)
    footer(s, 24)

    # ═══ S25 BILAN REFLEXIF ═══
    s = prs.slides.add_slide(blank)
    add_rect(s, Cm(0), Cm(0), W, Cm(0.4), OG)
    constructys_deco(s, OG)
    add_circle(s, Cm(1.2), Cm(0.8), Cm(1.2), OG)
    add_text(s, Cm(2.8), Cm(1), Cm(16), Cm(1), "BILAN REFLEXIF", 22, True, NV, font_name="Montserrat")
    orange_bar(s, Cm(1.8), Cm(2.8))
    add_text(s, Cm(1.8), Cm(3.2), Cm(17.4), Cm(0.7),
             "Avec le recul, chaque mission m'a appris quelque chose sur ma facon de travailler. "
             "Certains reflexes sont venus tardivement — le tableau de bord, la demande d'etudes "
             "complementaires, la traçabilite photographique. Ce bilan tire les lecons de 8 annees "
             "de pratique pour construire la suite.",
             7, False, RGBColor(0x55,0x55,0x55), PP_ALIGN.JUSTIFY, "Inter")
    bloc_colors = [(OG, LO), (TQ, RGBColor(0xE8,0xF8,0xF6)), (NV, RGBColor(0xED,0xF1,0xF5))]
    bloc_icons = ["A", "D", "M"]  # Appris, Differemment, MEC
    y = Cm(4)
    for i, (t, tx) in enumerate(PAGE_25["blocs"]):
        bc, bg = bloc_colors[i]
        add_rect(s, Cm(1.8), y, Cm(17.4), Cm(5.5), bg, 0.1, shadow=True)
        add_rect(s, Cm(1.8), y, Cm(0.15), Cm(5.5), bc)
        add_circle(s, Cm(2.3), y+Cm(0.4), Cm(1), bc)
        add_text(s, Cm(2.3), y+Cm(0.5), Cm(1), Cm(0.8), bloc_icons[i], 14, True, WH, PP_ALIGN.CENTER, "Montserrat")
        add_text(s, Cm(3.8), y+Cm(0.3), Cm(14), Cm(0.7), t, 11, True, bc, font_name="Montserrat")
        add_rect(s, Cm(3.8), y+Cm(1.1), Cm(3), Cm(0.06), bc)  # separateur accent
        add_text(s, Cm(2.5), y+Cm(1.3), Cm(16), Cm(4), tx, 8, False, RGBColor(0x55,0x55,0x55), PP_ALIGN.JUSTIFY, "Inter")
        y += Cm(6)
    footer(s, 25)

    # ═══ S26 BIM ═══
    s = prs.slides.add_slide(blank)
    add_rect(s, Cm(0), Cm(0), W, Cm(0.4), TQ)
    constructys_deco(s, TQ)
    add_circle(s, Cm(1.2), Cm(0.8), Cm(1.2), TQ)
    add_text(s, Cm(2.8), Cm(0.8), Cm(16), Cm(1.5), "PROTOCOLE DE\nCOLLABORATION BIM", 18, True, NV, font_name="Montserrat")
    orange_bar(s, Cm(1.8), Cm(3.2))
    add_rich_text(s, Cm(1.8), Cm(3.7), Cm(17.4), Cm(1.2),
             "Le BIM transforme la chaine metre - estimation - chiffrage en la rendant plus "
             "fiable et plus rapide. Ma formation a l'AFPA Colmar (8 mois, Technicien Modeleur BIM) "
             "m'a permis de modeliser un batiment R+2 complet et d'en extraire automatiquement "
             "**78 postes** de metres : surfaces de planchers, volumes beton, lineaires de murs "
             "et quantites d'acier. L'ecart avec le metre manuel traditionnel n'a ete que "
             "de **1,8%**. La detection de **12 clashs** structure/reseaux en amont du chantier "
             "illustre la valeur ajoutee concrete de l'approche BIM pour l'economiste de la construction.",
             8, TQ)
    # Convention left
    add_text(s, Cm(1.8), Cm(5.2), Cm(8), Cm(0.5), "Convention BIM appliquee", 10, True, OG, font_name="Montserrat")
    y = Cm(6)
    for k, v in PAGE_26["convention"]["items"]:
        add_text(s, Cm(2), y, Cm(3), Cm(0.6), k, 7, True, TQ, font_name="Inter")
        add_text(s, Cm(5.2), y, Cm(5), Cm(0.8), v, 7, False, DK, font_name="Inter")
        y += Cm(1.1)
    add_text(s, Cm(1.8), Cm(10.5), Cm(8), Cm(0.5), "Workflow BIM", 10, True, OG, font_name="Montserrat")
    y = Cm(11.3)
    for step in PAGE_26["workflow"]:
        add_circle(s, Cm(2), y+Cm(0.1), Cm(0.4), TQ)
        add_text(s, Cm(2.8), y, Cm(7), Cm(0.5), step, 8, False, RGBColor(0x55,0x55,0x55), font_name="Inter")
        y += Cm(0.9)
    # Cas concret right
    cas = PAGE_26["cas_concret"]
    add_rect(s, Cm(11), Cm(5.2), Cm(8), Cm(6), WH, 0.1, shadow=True)
    add_rect(s, Cm(11), Cm(5.2), Cm(0.15), Cm(6), OG)
    add_text(s, Cm(11.5), Cm(5.5), Cm(7), Cm(0.7), cas["titre"], 9, True, OG, font_name="Montserrat")
    add_text(s, Cm(11.5), Cm(6.5), Cm(7), Cm(4.5), cas["details"], 7, False, RGBColor(0x55,0x55,0x55), PP_ALIGN.JUSTIFY, "Inter")
    # Apport BIM
    add_rect(s, Cm(11), Cm(11.8), Cm(8), Cm(4.5), LG, 0.1, shadow=True)
    add_rect(s, Cm(11), Cm(11.8), Cm(0.15), Cm(4.5), TQ)
    add_text(s, Cm(11.5), Cm(12), Cm(7), Cm(0.6), "Apport du BIM pour le MEC", 9, True, TQ, font_name="Montserrat")
    add_text(s, Cm(11.5), Cm(12.8), Cm(7), Cm(3.2), PAGE_26["apport_mec"], 7, False, RGBColor(0x55,0x55,0x55), PP_ALIGN.JUSTIFY, "Inter")
    footer(s, 26)

    # ═══ S27 PROJET PRO ═══
    s = prs.slides.add_slide(blank)
    add_rect(s, Cm(0), Cm(0), W, Cm(0.4), OG)
    constructys_deco(s, OG)
    add_circle(s, Cm(1.2), Cm(0.8), Cm(1.2), OG)
    add_text(s, Cm(2.8), Cm(0.8), Cm(16), Cm(1), "MON PROJET PROFESSIONNEL", 18, True, NV, font_name="Montserrat")
    orange_bar(s, Cm(1.8), Cm(2.5))
    add_text(s, Cm(1.8), Cm(3), Cm(17.4), Cm(1.8),
             "Mon projet repose sur une conviction forte : les outils numeriques doivent etre "
             "au service de l'economiste de la construction — et non l'inverse. "
             "Le marche de l'ingenierie BIM est domine par des architectes et des ingenieurs ; "
             "tres peu d'economistes maitrisent a la fois le BIM, le developpement d'outils "
             "et la realite du terrain. BIMCO occupe ce creneau rare : la rigueur de l'economiste "
             "(metres, etudes de prix, suivi financier, analyse CAO) combinee aux outils du BIM "
             "(Revit, Dynamo, plugins C#, applications web React/Node.js). "
             "Ma double culture Maroc/France et ma formation BIM constituent les fondations de cette ambition.",
             8, False, RGBColor(0x55,0x55,0x55), PP_ALIGN.JUSTIFY, "Inter")
    hz_c = [OG, TQ, NV]
    hz_icons_text = ["CT", "MT", "LT"]
    for i, (label, date, desc) in enumerate(PAGE_27["horizons"]):
        x = Cm(1.5) + Cm(6.2)*i
        add_rect(s, x, Cm(5), Cm(5.8), Cm(10), WH, 0.1, shadow=True)
        add_rect(s, x, Cm(5), Cm(5.8), Cm(0.2), hz_c[i])
        add_circle(s, x+Cm(0.3), Cm(5.5), Cm(1), hz_c[i])
        add_text(s, x+Cm(1.5), Cm(5.5), Cm(4), Cm(0.7), label, 10, True, hz_c[i], font_name="Montserrat")
        add_text(s, x+Cm(0.5), Cm(6.5), Cm(5), Cm(0.5), date, 9, True, OG, font_name="Montserrat")
        add_text(s, x+Cm(0.5), Cm(7.2), Cm(5), Cm(6), desc.replace('\n', '\n'), 7, False, RGBColor(0x55,0x55,0x55), font_name="Inter")
    # Quote
    add_rect(s, Cm(1.8), Cm(16), Cm(17.4), Cm(2), LO, 0.1, shadow=True)
    add_rect(s, Cm(1.8), Cm(16), Cm(0.15), Cm(2), OG)
    add_text(s, Cm(2.5), Cm(16.3), Cm(16), Cm(1.5), PAGE_27["citation"].replace('\n', ' '), 9, True, NV, PP_ALIGN.CENTER, "Inter")
    # Double competence capsule
    add_rect(s, Cm(1.8), Cm(18.5), Cm(17.4), Cm(2.2), NV, 0.5, shadow=True)
    add_text(s, Cm(2.5), Cm(18.7), Cm(16), Cm(0.5), "MA DOUBLE COMPETENCE", 7, True, RGBColor(0xAA,0xAA,0xBB), PP_ALIGN.CENTER, "Montserrat")
    add_text(s, Cm(2.5), Cm(19.4), Cm(16), Cm(0.7), "Economiste terrain + Developpeur BIM", 12, True, WH, PP_ALIGN.CENTER, "Montserrat")
    add_text(s, Cm(2.5), Cm(20.2), Cm(16), Cm(0.4), "Un avantage differenciant rare sur le marche de l'ingenierie", 8, False, RGBColor(0xAA,0xAA,0xBB), PP_ALIGN.CENTER, "Inter")
    # Geometric shapes
    add_rect(s, Cm(18), Cm(22), Cm(3), Cm(2.5), TQ, 0.15)
    add_rect(s, Cm(18), Cm(24.5), Cm(2), Cm(1.5), NV, 0.1)
    footer(s, 27)

    # ═══ S28 CONCLUSION ═══
    s = prs.slides.add_slide(blank)
    constructys_header(s, "CONCLUSION", "CONCLUSION", NV, OG, tall=False)
    # Resume — paragraphe de synthese (ne repete pas les chiffres de la cover)
    add_text(s, Cm(1.8), Cm(4.8), Cm(17.4), Cm(1.5),
             PAGE_28["resume"], 8, False, RGBColor(0x44,0x44,0x44), PP_ALIGN.JUSTIFY, "Inter")
    # KPI row — 5 resultats cles, sur fond mini-cartes
    add_text(s, Cm(1.8), Cm(6.6), Cm(17.4), Cm(0.5), "EN CHIFFRES", 7, True, OG, PP_ALIGN.CENTER, "Montserrat")
    kpi_colors_28 = [OG, TQ, NV, TQ, OG]
    for i, (val, lab) in enumerate(PAGE_28["kpis"]):
        add_kpi_card(s, Cm(1.8) + Cm(3.6)*i, Cm(7.2), val, lab, kpi_colors_28[i], Cm(3.0))
    # Separator bar
    add_rect(s, Cm(1.8), Cm(9.9), Cm(17.4), Cm(0.06), RGBColor(0xEE,0xEE,0xEE))
    # Points — 4 enseignements (nouveau contenu, pas de répétition)
    pt_colors = [OG, TQ, NV, PR]
    y = Cm(10.2)
    for i, (t, d) in enumerate(PAGE_28["points"]):
        add_rect(s, Cm(1.8), y, Cm(0.2), Cm(2.4), pt_colors[i])
        add_circle(s, Cm(1.8), y+Cm(0.05), Cm(0.8), pt_colors[i])
        add_text(s, Cm(1.8), y+Cm(0.15), Cm(0.8), Cm(0.7), str(i+1), 10, True, WH, PP_ALIGN.CENTER, "Montserrat")
        add_text(s, Cm(2.8), y, Cm(16), Cm(0.7), t, 9, True, NV, font_name="Montserrat")
        add_text(s, Cm(2.8), y+Cm(0.75), Cm(16), Cm(1.5), d, 7, False, RGBColor(0x55,0x55,0x55), PP_ALIGN.JUSTIFY, "Inter")
        y += Cm(2.5)
    # Citation
    add_rect(s, Cm(1.8), y+Cm(0.3), Cm(17.4), Cm(3), LO, 0.1, shadow=True)
    add_rect(s, Cm(1.8), y+Cm(0.3), Cm(0.15), Cm(3), OG)
    add_text(s, Cm(2.5), y+Cm(0.6), Cm(16), Cm(2.5), PAGE_28["citation"].replace('\n', ' '), 8, True, NV, PP_ALIGN.CENTER, "Inter")
    # Pied capsule
    pied_c = add_rect(s, Cm(1.8), y+Cm(3.8), Cm(17.4), Cm(1), NV, 0.5, shadow=True)
    set_gradient(pied_c, NV, _rgb_shift(NV, -15), 0)
    add_text(s, Cm(2.5), y+Cm(4), Cm(16), Cm(0.7), PAGE_28["pied"], 7, False, WH, PP_ALIGN.CENTER, "Inter")
    footer(s, 28)

    # ═══ S29 ANNEXE 1 ═══
    s = prs.slides.add_slide(blank)
    add_rect(s, Cm(0), Cm(0), W, Cm(0.4), NV)
    constructys_deco(s, NV)
    add_text(s, Cm(1.8), Cm(1.5), Cm(17), Cm(1), "ANNEXE 1 : DOCUMENTS OFFICIELS", 20, True, NV, font_name="Montserrat")
    orange_bar(s, Cm(1.8), Cm(3.2))
    add_text(s, Cm(1.8), Cm(3.7), Cm(17.4), Cm(1.2),
             "Trois documents attestant de la realite des activites decrites dans ce rapport : "
             "convocation a la commission d'appel d'offres, notification de rejet d'une offre "
             "non conforme, et derniere page signee du cahier des prescriptions speciales du "
             "projet routier. Les originaux seront presentes lors de l'epreuve orale.",
             8, False, RGBColor(0x55,0x55,0x55), PP_ALIGN.JUSTIFY, "Inter")
    for i, (_, titre, desc) in enumerate(PAGE_29["documents"]):
        x = Cm(1.5) + Cm(6.2)*i
        add_rect(s, x, Cm(5.5), Cm(5.8), Cm(8), LG, 0.1, shadow=True)
        # Insert document image if available
        if i < len(DOC_IMGS):
            d_img = pic(s, DOC_IMGS[i], x+Cm(0.3), Cm(5.8), Cm(5.2), Cm(7))
            if not d_img:
                add_text(s, x, Cm(9), Cm(5.8), Cm(0.5), "Document", 8, False, GY, PP_ALIGN.CENTER, "Inter")
        add_text(s, x, Cm(14), Cm(5.8), Cm(0.7), titre, 9, True, OG, PP_ALIGN.CENTER, "Montserrat")
        add_text(s, x, Cm(14.8), Cm(5.8), Cm(1), desc.replace('\n', ' '), 7, False, GY, PP_ALIGN.CENTER, "Inter")
    add_rect(s, Cm(1.8), Cm(17), Cm(17.4), Cm(1), LG, 0.5)
    add_text(s, Cm(2.5), Cm(17.1), Cm(16), Cm(0.8),
             "Ces documents originaux seront presentes lors de l'epreuve orale.", 7, False, GY, PP_ALIGN.CENTER, "Inter")
    add_rect(s, Cm(18), Cm(22), Cm(3), Cm(4), TQ, 0.15)
    footer(s, 29)

    # ═══ S30 ANNEXE 2 ═══
    s = prs.slides.add_slide(blank)
    add_rect(s, Cm(0), Cm(0), W, Cm(0.4), NV)
    constructys_deco(s, NV)
    add_text(s, Cm(1.8), Cm(1.5), Cm(17), Cm(1), "ANNEXE 2 : PHOTOS DE CHANTIER", 20, True, NV, font_name="Montserrat")
    orange_bar(s, Cm(1.8), Cm(3.2))
    add_text(s, Cm(1.8), Cm(3.7), Cm(17.4), Cm(1.2),
             "Photographies prises sur les chantiers de la province de Khenifra : amenagement "
             "urbain des 4 communes (pose de bordures, tranchees d'assainissement, mise en "
             "oeuvre d'enrobes, eclairage public) et route Lehri-Kerrouchen (terrassement "
             "rocheux, ouvrages hydrauliques). Ces images temoignent des conditions reelles "
             "de travail : relief montagneux, sol rocheux, chantiers distants.",
             8, False, RGBColor(0x55,0x55,0x55), PP_ALIGN.JUSTIFY, "Inter")
    for i, (_, caption) in enumerate(PAGE_30["photos"]):
        col = i % 3
        row = i // 3
        x = Cm(1.5) + Cm(6.2)*col
        y = Cm(5) + Cm(5)*row
        add_rect(s, x, y, Cm(5.8), Cm(3.5), LG, 0.1, shadow=True)
        # Insert real chantier photo
        if i < len(CHANTIER_IMGS):
            ch_img = pic(s, CHANTIER_IMGS[i], x+Cm(0.1), y+Cm(0.1), Cm(5.6), Cm(3.3))
            if ch_img:
                add_shadow(ch_img)
            else:
                add_text(s, x, y+Cm(1.3), Cm(5.8), Cm(0.5), "Photo", 8, False, GY, PP_ALIGN.CENTER, "Inter")
        add_text(s, x, y+Cm(3.7), Cm(5.8), Cm(0.5), caption, 7, False, GY, PP_ALIGN.CENTER, "Inter")
    add_rect(s, Cm(0), Cm(27), Cm(1), Cm(2.7), OG, 0.15)
    footer(s, 30)

    return prs


def main():
    print("=" * 60)
    print("RAPPORT U62 - PPTX CANVA TEMPLATE")
    print("=" * 60)
    prs = build()
    os.makedirs(os.path.dirname(OUT), exist_ok=True)
    prs.save(OUT)
    sz = os.path.getsize(OUT) / 1024
    print(f"PPTX : {OUT}")
    print(f"Slides: {len(prs.slides)} | Taille: {sz:.0f} Ko")
    print("=" * 60)
    print(">> Importe ce fichier dans Canva pour le personnaliser !")


if __name__ == "__main__":
    main()
