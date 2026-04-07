"""Extract speaker notes from PPTX and create a printable HTML file."""
from pptx import Presentation
import html

prs = Presentation("SOUTENANCE_U62_v4.pptx")

notes_data = []
for i, slide in enumerate(prs.slides, 1):
    note_text = ""
    if slide.has_notes_slide:
        notes_slide = slide.notes_slide
        for para in notes_slide.notes_text_frame.paragraphs:
            note_text += para.text + "\n"
    notes_data.append((i, note_text.strip()))

with open("NOTES_ORAL_SOUTENANCE.html", "w", encoding="utf-8") as f:
    f.write("""<!DOCTYPE html>
<html lang="fr">
<head>
<meta charset="UTF-8">
<title>Notes de pr\u00e9sentation - Soutenance BTS MEC U62</title>
<style>
@media print { .page-break { page-break-before: always; } }
body { font-family: Calibri, Arial, sans-serif; font-size: 13pt; line-height: 1.6; max-width: 800px; margin: 0 auto; padding: 20px; color: #222; }
h1 { text-align: center; color: #007A7F; font-size: 20pt; border-bottom: 3px solid #007A7F; padding-bottom: 10px; }
h2 { color: #007A7F; font-size: 15pt; margin-top: 30px; padding: 8px 12px; background: #E0F4F5; border-left: 5px solid #007A7F; }
.timing { float: right; background: #007A7F; color: white; padding: 2px 10px; border-radius: 4px; font-size: 11pt; font-weight: bold; }
.note-body { padding: 5px 15px; white-space: pre-wrap; font-size: 12pt; }
.important { background: #FFF3CD; border-left: 4px solid #D97706; padding: 8px 12px; margin: 10px 0; }
.tip { color: #666; font-style: italic; font-size: 10pt; }
.footer { text-align: center; margin-top: 40px; padding-top: 15px; border-top: 2px solid #007A7F; color: #666; font-size: 10pt; }
.total { text-align: center; font-size: 16pt; font-weight: bold; color: #007A7F; margin: 20px 0; padding: 15px; background: #E0F4F5; border-radius: 8px; }
</style>
</head>
<body>
<h1>NOTES DE PR\u00c9SENTATION ORALE<br><span style="font-size:14pt;color:#666;">Soutenance BTS MEC U62 \u2013 Session 2026 \u2013 BAHAFID Mohamed</span></h1>

<div class="total">TIMING TOTAL VIS\u00c9 : 18\u201320 minutes sur 20 min allou\u00e9es</div>

<p class="tip">\u2139\ufe0f Ces notes sont \u00e0 utiliser pendant l'entra\u00eenement. Le jour J, vous ne devez PAS lire \u2013 parlez naturellement en vous appuyant sur les slides. Le mode Pr\u00e9sentateur (F5) affichera ces notes sur votre \u00e9cran.</p>
""")

    timings = ["1 min", "1 min 30", "1 min 30", "1 min", "1 min", "2 min", "2 min 30", "1 min", "1 min 30", "1 min", "1 min", "1 min 30", "2 min", "1 min 30", "1 min 30"]

    for i, (num, note) in enumerate(notes_data):
        if not note:
            continue
        timing = timings[i] if i < len(timings) else ""

        # Extract title from first line
        lines = note.split("\n")
        title = lines[0] if lines else f"SLIDE {num}"
        body = "\n".join(lines[1:]).strip()

        important = ""
        if "IMPORTANTE" in title.upper():
            important = ' class="important"'

        f.write(f'\n<h2>{html.escape(title)}<span class="timing">{timing}</span></h2>\n')
        f.write(f'<div class="note-body"{important}>{html.escape(body)}</div>\n')

        # Add page break after slide 7 (halfway)
        if num == 8:
            f.write('<div class="page-break"></div>\n')

    f.write("""
<div class="page-break"></div>
<h2>QUESTIONS PROBABLES DU JURY</h2>
<div class="note-body">
<strong>1. Pourquoi candidat libre et pas en alternance ?</strong>
\u2192 Mon exp\u00e9rience de 8 ans couvre d\u00e9j\u00e0 largement le r\u00e9f\u00e9rentiel. Le candidat libre me permet de valoriser ce parcours sans reprendre une formation classique. BIMCO est ma structure d'accueil.

<strong>2. Comment adaptez-vous la m\u00e9thode marocaine aux normes fran\u00e7aises ?</strong>
\u2192 Les fondamentaux sont identiques : transparence, mise en concurrence, offre \u00e9conomiquement la plus avantageuse. CPS = CCAP, BPDE = BPU/DQE. La logique de protection du MOA est la m\u00eame. Je dois adapter la terminologie et les r\u00e9f\u00e9rences normatives (DTU/Eurocodes vs normes marocaines).

<strong>3. Le BIM apporte quoi de plus qu'Excel pour les m\u00e9tr\u00e9s ?</strong>
\u2192 Tra\u00e7abilit\u00e9 : chaque quantit\u00e9 est li\u00e9e \u00e0 un objet 3D. Si le mod\u00e8le change, les quantit\u00e9s se recalculent. Sur mon cas AFPA : 78 postes en 2h vs 2 jours, \u00e9cart 1,8%. Plus fiable, plus rapide, plus tra\u00e7able.

<strong>4. L'estimation confidentielle, \u00e7a existe en France ?</strong>
\u2192 Oui, sous le nom d'estimation du ma\u00eetre d'ouvrage. Le principe est le m\u00eame : fixer un prix de r\u00e9f\u00e9rence avant la mise en concurrence. Au Maroc elle est confidentielle et constitue le prix plafond strict. En France elle sert de r\u00e9f\u00e9rence mais n'est pas toujours \u00e9liminatoire.

<strong>5. Comment g\u00e9rez-vous un client qui conteste vos m\u00e9tr\u00e9s BIM ?</strong>
\u2192 L'avantage du BIM : je peux montrer pr\u00e9cis\u00e9ment d'o\u00f9 vient chaque quantit\u00e9 dans la maquette. Chaque ligne du DPGF est li\u00e9e \u00e0 un objet. C'est une tra\u00e7abilit\u00e9 que le m\u00e9tr\u00e9 sur plan ne permet pas.

<strong>6. Votre tableau de bord, il ressemblait \u00e0 quoi concr\u00e8tement ?</strong>
\u2192 Excel consolid\u00e9, 3 indicateurs par poste : avancement physique (%), consommation budg\u00e9taire (DH), \u00e9cart pr\u00e9visionnel (\u0394%). Mis \u00e0 jour chaque semaine. Un code couleur : vert/orange/rouge selon la proximit\u00e9 du seuil d'avenant.
</div>

<h2>CHECKLIST JOUR J</h2>
<div class="note-body">
\u2610 Cl\u00e9 USB avec la pr\u00e9sentation (+ copie de secours)
\u2610 Rapport papier (le m\u00eame que celui d\u00e9pos\u00e9)
\u2610 Chrono/montre visible pendant la pr\u00e9sentation
\u2610 V\u00e9rifier le mode Pr\u00e9sentateur avant de commencer
\u2610 Tenue professionnelle
\u2610 Arriver 15 min en avance pour tester le mat\u00e9riel
</div>

<div class="footer">
BAHAFID Mohamed \u00b7 BTS MEC U62 \u00b7 Session 2026 \u00b7 Acad\u00e9mie de Lyon<br>
BIMCO | gestion.bimco-consulting.fr
</div>
</body>
</html>
""")

print("OK: NOTES_ORAL_SOUTENANCE.html generated")
