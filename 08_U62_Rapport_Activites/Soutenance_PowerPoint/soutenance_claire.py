import sys
sys.stdout.reconfigure(encoding='utf-8')

import os
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import ChartData
from pptx.oxml.ns import qn

# ── Dimensions 16:9 ──────────────────────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

# ── Colour palette ────────────────────────────────────────────────────────────
BG     = RGBColor(0xF5, 0xF8, 0xFA)
CARD   = RGBColor(0xE8, 0xF5, 0xF6)
CARD2  = RGBColor(0xF0, 0xF4, 0xF8)
TURQ   = RGBColor(0x00, 0x95, 0x9E)
TURQD  = RGBColor(0x00, 0x6E, 0x78)
TURQL  = RGBColor(0xCC, 0xEE, 0xF0)
NAVY   = RGBColor(0x1C, 0x33, 0x40)
GRIS   = RGBColor(0x6E, 0x8A, 0x96)
GRISL  = RGBColor(0xD4, 0xE1, 0xE5)
GOLD   = RGBColor(0xF5, 0xA1, 0x18)
ROUGE  = RGBColor(0xE5, 0x4E, 0x3C)
VERT   = RGBColor(0x27, 0xAE, 0x60)
ORANGE = RGBColor(0xF3, 0x9C, 0x12)
BLANC  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x0F, 0x2A, 0x33)

MEDIA = "d:/PREPA BTS MEC/08_U62_Rapport_Activites/Soutenance_PowerPoint/extracted_media"

# ── Low-level helpers ─────────────────────────────────────────────────────────
def blank_slide(prs):
    layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(layout)

def R(sl, x, y, w, h, fill_color=None):
    """Add a filled rectangle."""
    shape = sl.shapes.add_shape(1, x, y, w, h)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    return shape

def T(sl, text, x, y, w, h, sz=14, bold=False, italic=False,
      color=NAVY, align=PP_ALIGN.LEFT, wrap=True):
    """Add a text box."""
    tb = sl.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Calibri"
    run.font.size = Pt(sz)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb

def set_bg(sl, color):
    bg = sl.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_image(sl, path, x, y, w, h):
    if os.path.exists(path):
        return sl.shapes.add_picture(path, x, y, w, h)
    return None

def dark_overlay(sl, alpha=75):
    sh = R(sl, 0, 0, W, H, DARK)
    spPr = sh._element.spPr
    sf = spPr.find('.//' + qn('a:solidFill'))
    if sf is not None and len(sf):
        ae = etree.SubElement(sf[0], qn('a:alpha'))
        ae.set('val', str(alpha * 1000))

def footer(sl, num):
    R(sl, 0, H - Inches(0.28), W, Inches(0.28), GRISL)
    T(sl, "BAHAFID Mohamed . BTS MEC U62 . Session 2026",
      Inches(0.15), H - Inches(0.27), Inches(8), Inches(0.25),
      sz=9, color=GRIS)
    T(sl, str(num),
      W - Inches(0.6), H - Inches(0.27), Inches(0.45), Inches(0.25),
      sz=9, bold=True, color=TURQ, align=PP_ALIGN.RIGHT)

def header_band(sl, label, title, subtitle=""):
    R(sl, 0, 0, W, Inches(1.15), BLANC)
    R(sl, 0, 0, Inches(0.07), Inches(1.15), TURQ)
    T(sl, label, Inches(0.2), Inches(0.08), Inches(4), Inches(0.3),
      sz=10, bold=True, color=TURQ)
    T(sl, title, Inches(0.2), Inches(0.35), W - Inches(0.4), Inches(0.5),
      sz=22, bold=True, color=NAVY)
    if subtitle:
        T(sl, subtitle, Inches(0.2), Inches(0.82), W - Inches(0.4), Inches(0.3),
          sz=12, color=GRIS)

def bar_chart(sl, title, cats, vals, x, y, w, h, color="00959E", horizontal=False):
    cd = ChartData()
    cd.categories = cats
    cd.add_series("", vals)
    chart_type = XL_CHART_TYPE.BAR_CLUSTERED if horizontal else XL_CHART_TYPE.COLUMN_CLUSTERED
    shape = sl.shapes.add_chart(chart_type, x, y, w, h, cd)
    ch = shape.chart
    ch.has_title = True
    ch.chart_title.text_frame.text = title
    if ch.chart_title.text_frame.paragraphs[0].runs:
        r = ch.chart_title.text_frame.paragraphs[0].runs[0]
        r.font.size = Pt(11)
        r.font.bold = True
        r.font.color.rgb = NAVY
    plot = ch.plots[0]
    plot.gap_width = 55
    plot.has_data_labels = True
    plot.data_labels.show_value = True
    r2, g2, b2 = int(color[0:2], 16), int(color[2:4], 16), int(color[4:6], 16)
    plot.series[0].format.fill.solid()
    plot.series[0].format.fill.fore_color.rgb = RGBColor(r2, g2, b2)
    ch.has_legend = False
    return ch

def pie_chart(sl, cats, vals, colors_hex, x, y, w, h):
    cd = ChartData()
    cd.categories = cats
    cd.add_series("", vals)
    shape = sl.shapes.add_chart(XL_CHART_TYPE.PIE, x, y, w, h, cd)
    ch = shape.chart
    ch.has_title = False
    plot = ch.plots[0]
    plot.has_data_labels = True
    plot.data_labels.show_percentage = True
    plot.data_labels.show_category_name = False
    for i, hc in enumerate(colors_hex):
        r2, g2, b2 = int(hc[0:2], 16), int(hc[2:4], 16), int(hc[4:6], 16)
        plot.series[0].points[i].format.fill.solid()
        plot.series[0].points[i].format.fill.fore_color.rgb = RGBColor(r2, g2, b2)
    ch.has_legend = True
    ch.legend.position = XL_LEGEND_POSITION.BOTTOM
    return ch

def cpar_row(sl, ctx, pb, act, res, x, y, w, h=Inches(1.55)):
    bw = (w - Inches(0.32)) / 4
    configs = [
        ("C", "CONTEXTE", TURQ,  TURQL,                        ctx),
        ("P", "PROBLEME", ROUGE, RGBColor(0xFF, 0xEE, 0xEC),   pb),
        ("A", "ACTION",   VERT,  RGBColor(0xE8, 0xF8, 0xEE),   act),
        ("R", "RESULTAT", GOLD,  RGBColor(0xFF, 0xF5, 0xE0),   res),
    ]
    for i, (letter, label, accent, bg_col, content) in enumerate(configs):
        bx = x + i * (bw + Inches(0.1))
        R(sl, bx, y, bw, h, bg_col)
        R(sl, bx, y, bw, Inches(0.05), accent)
        c = sl.shapes.add_shape(9, bx + Inches(0.08), y + Inches(0.1),
                                Inches(0.4), Inches(0.4))
        c.fill.solid()
        c.fill.fore_color.rgb = accent
        c.line.fill.background()
        T(sl, letter, bx + Inches(0.08), y + Inches(0.1),
          Inches(0.4), Inches(0.4), sz=16, bold=True, color=BLANC,
          align=PP_ALIGN.CENTER)
        T(sl, label, bx + Inches(0.56), y + Inches(0.13),
          bw - Inches(0.65), Inches(0.3), sz=9, bold=True, color=accent)
        T(sl, content, bx + Inches(0.1), y + Inches(0.52),
          bw - Inches(0.18), h - Inches(0.6), sz=12, color=NAVY)

# ── SLIDE 1 – COVER ───────────────────────────────────────────────────────────
def slide1(prs):
    sl = blank_slide(prs)
    set_bg(sl, DARK)
    img1 = os.path.join(MEDIA, "image1.jpeg")
    add_image(sl, img1, 0, 0, W, H)
    dark_overlay(sl, 78)

    # Left turquoise bar
    R(sl, 0, 0, Inches(0.1), H, TURQ)
    # Gold top line
    R(sl, 0, 0, W, Inches(0.06), GOLD)

    T(sl, "Rapport d'Activites", Inches(0.5), Inches(0.8), Inches(10), Inches(1.0),
      sz=58, bold=True, color=BLANC)
    T(sl, "Professionnelles", Inches(0.5), Inches(1.7), Inches(10), Inches(0.9),
      sz=48, bold=True, color=TURQ)
    T(sl, "BAHAFID Mohamed", Inches(0.5), Inches(2.85), Inches(9), Inches(0.7),
      sz=38, bold=True, color=BLANC)
    T(sl, "N 02537399911 . Academie de Lyon", Inches(0.5), Inches(3.6),
      Inches(8), Inches(0.4), sz=15, color=GRIS)

    # Stat boxes
    boxes = [
        ("8 ans BTP", TURQ),
        ("2 pays Maroc + France", GOLD),
        ("7 marches publics", VERT),
    ]
    bw = Inches(3.5)
    for i, (txt, col) in enumerate(boxes):
        bx = Inches(0.5) + i * (bw + Inches(0.25))
        R(sl, bx, Inches(5.0), bw, Inches(0.65), col)
        T(sl, txt, bx + Inches(0.1), Inches(5.05), bw - Inches(0.2), Inches(0.55),
          sz=18, bold=True, color=BLANC, align=PP_ALIGN.CENTER)

    T(sl, "BTS MEC U62 . Session 2026 . BIMCO",
      Inches(0.5), H - Inches(0.65), Inches(10), Inches(0.35),
      sz=13, italic=True, color=TURQL)

# ── SLIDE 2 – PROFIL ─────────────────────────────────────────────────────────
def slide2(prs):
    sl = blank_slide(prs)
    set_bg(sl, BG)
    header_band(sl, "QUI SUIS-JE ?", "Profil Candidat - BAHAFID Mohamed")

    # Photo right
    img2 = os.path.join(MEDIA, "image2.jpeg")
    add_image(sl, img2, W - Inches(2.9), Inches(0.1), Inches(2.65), Inches(3.4))

    # Timeline
    ty = Inches(1.5)
    R(sl, Inches(0.5), ty + Inches(0.18), Inches(8.5), Inches(0.05), GRISL)
    timeline = [
        ("2016", TURQ,  Inches(0.5),  "Maroc / Conseil\nRegional BMK"),
        ("2019", GOLD,  Inches(3.0),  "France / Chef de\nchantier GO 5 ans"),
        ("2025", VERT,  Inches(5.5),  "BIM / Formation\nAFPA Colmar"),
        ("2026", TURQD, Inches(8.0),  "BIMCO /\nMicro-entreprise"),
    ]
    for yr, col, cx, lbl in timeline:
        c = sl.shapes.add_shape(9, cx, ty, Inches(0.42), Inches(0.42))
        c.fill.solid(); c.fill.fore_color.rgb = col; c.line.fill.background()
        T(sl, yr, cx - Inches(0.05), ty - Inches(0.3), Inches(0.6), Inches(0.28),
          sz=10, bold=True, color=col, align=PP_ALIGN.CENTER)
        T(sl, lbl, cx - Inches(0.15), ty + Inches(0.45), Inches(1.2), Inches(0.6),
          sz=9, color=NAVY)

    # KPI boxes
    kpis = [("8 ans", "Experience"), ("7 marches", "Publics"), ("82,5 M DH", "Budget total"), ("5 situations", "CPAR")]
    kw = Inches(2.1)
    ky = Inches(2.4)
    for i, (val, lbl) in enumerate(kpis):
        kx = Inches(0.3) + i * (kw + Inches(0.15))
        R(sl, kx, ky, kw, Inches(0.75), CARD)
        R(sl, kx, ky, kw, Inches(0.06), TURQ)
        T(sl, val, kx + Inches(0.1), ky + Inches(0.08), kw - Inches(0.2), Inches(0.35),
          sz=18, bold=True, color=TURQ, align=PP_ALIGN.CENTER)
        T(sl, lbl, kx + Inches(0.1), ky + Inches(0.44), kw - Inches(0.2), Inches(0.25),
          sz=9, color=GRIS, align=PP_ALIGN.CENTER)

    # Pie chart right
    pie_chart(sl,
              ["France BTP 5 ans", "Maroc MOA 3 ans", "BIM 8 mois"],
              [58, 35, 8],
              ["00959E", "F5A118", "27AE60"],
              W - Inches(3.5), Inches(3.4), Inches(3.0), Inches(2.5))

    # Info rows
    rows = [
        ("Structure :", "Agence Urbaine / Conseil Regional BMK (Maroc) . Puis chantiers France"),
        ("Poste :", "Technicien BTP / Conducteur de travaux"),
        ("Formation BIM :", "AFPA Colmar 2025 - Autodesk Revit + BIM360"),
        ("Activite actuelle :", "BIMCO micro-entreprise BIM"),
    ]
    ry = Inches(3.4)
    for lbl, txt in rows:
        R(sl, Inches(0.3), ry, Inches(9.3), Inches(0.38), BLANC)
        T(sl, lbl, Inches(0.4), ry + Inches(0.05), Inches(1.8), Inches(0.28),
          sz=10, bold=True, color=TURQ)
        T(sl, txt, Inches(2.1), ry + Inches(0.05), Inches(7.4), Inches(0.28),
          sz=10, color=NAVY)
        ry += Inches(0.42)

    footer(sl, 2)

# ── SLIDE 3 – VUE ENSEMBLE ────────────────────────────────────────────────────
def slide3(prs):
    sl = blank_slide(prs)
    set_bg(sl, BG)
    header_band(sl, "VUE D'ENSEMBLE",
                "MES ACTIVITES PROFESSIONNELLES",
                "2 Projets . 5 Situations CPAR")

    cw = (W - Inches(0.6)) / 2

    # Projet 1 card
    R(sl, Inches(0.2), Inches(1.25), cw, Inches(4.8), BLANC)
    R(sl, Inches(0.2), Inches(1.25), cw, Inches(0.08), TURQ)
    img31 = os.path.join(MEDIA, "image31.jpeg")
    add_image(sl, img31, Inches(0.2), Inches(1.33), cw, Inches(1.4))
    dark_overlay_rect(sl, Inches(0.2), Inches(1.33), cw, Inches(1.4), 60)
    T(sl, "PROJET 1 - Mise a Niveau - 4 Communes",
      Inches(0.35), Inches(1.38), cw - Inches(0.3), Inches(0.5),
      sz=14, bold=True, color=BLANC)

    p1_lines = [
        "Budget : 15,8 M DH HT",
        "Communes : Ouaoumana, El Hammam, Kerrouchen, Sebt Ait Rahou",
        "Nature : VRD - routes, trottoirs, eclairage, paysager",
        "Marche : Appel d'offres ouvert - Khenifra",
    ]
    for i, line in enumerate(p1_lines):
        T(sl, "• " + line, Inches(0.35), Inches(2.82) + i * Inches(0.33),
          cw - Inches(0.3), Inches(0.3), sz=10, color=NAVY)

    sits = [("Sit.1", TURQ), ("Sit.2", GOLD), ("Sit.3", VERT), ("Sit.4", ORANGE)]
    for i, (lbl, col) in enumerate(sits):
        cx = Inches(0.35) + i * Inches(1.1)
        c = sl.shapes.add_shape(9, cx, Inches(4.3), Inches(0.45), Inches(0.45))
        c.fill.solid(); c.fill.fore_color.rgb = col; c.line.fill.background()
        T(sl, lbl, cx, Inches(4.3), Inches(0.45), Inches(0.45),
          sz=9, bold=True, color=BLANC, align=PP_ALIGN.CENTER)

    # Projet 2 card
    px2 = Inches(0.2) + cw + Inches(0.2)
    R(sl, px2, Inches(1.25), cw, Inches(4.8), BLANC)
    R(sl, px2, Inches(1.25), cw, Inches(0.08), GOLD)
    img39 = os.path.join(MEDIA, "image39.jpeg")
    add_image(sl, img39, px2, Inches(1.33), cw, Inches(1.4))
    dark_overlay_rect(sl, px2, Inches(1.33), cw, Inches(1.4), 60)
    T(sl, "PROJET 2 - Route Lehri-Kerrouchen - 25 km",
      px2 + Inches(0.15), Inches(1.38), cw - Inches(0.3), Inches(0.5),
      sz=14, bold=True, color=BLANC)

    p2_lines = [
        "Budget : 29 M DH HT",
        "Longueur : 25 km zone montagneuse",
        "Nature : Route principale - terrassements + chaussee",
        "Marche : Appel d'offres ouvert regional",
    ]
    for i, line in enumerate(p2_lines):
        T(sl, "• " + line, px2 + Inches(0.15), Inches(2.82) + i * Inches(0.33),
          cw - Inches(0.3), Inches(0.3), sz=10, color=NAVY)

    c5 = sl.shapes.add_shape(9, px2 + Inches(0.15), Inches(4.3), Inches(0.45), Inches(0.45))
    c5.fill.solid(); c5.fill.fore_color.rgb = ROUGE; c5.line.fill.background()
    T(sl, "Sit.5", px2 + Inches(0.15), Inches(4.3), Inches(0.45), Inches(0.45),
      sz=9, bold=True, color=BLANC, align=PP_ALIGN.CENTER)

    # Pie chart bottom right
    pie_chart(sl,
              ["Projet 1", "Projet 2", "Autres"],
              [53.5, 29, 25],
              ["00959E", "F5A118", "27AE60"],
              W - Inches(3.2), Inches(4.8), Inches(3.0), Inches(2.3))

    footer(sl, 3)

def dark_overlay_rect(sl, x, y, w, h, alpha=60):
    sh = sl.shapes.add_shape(1, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = DARK
    sh.line.fill.background()
    spPr = sh._element.spPr
    sf = spPr.find('.//' + qn('a:solidFill'))
    if sf is not None and len(sf):
        ae = etree.SubElement(sf[0], qn('a:alpha'))
        ae.set('val', str(alpha * 1000))

# ── SLIDES 4-8 – SITUATIONS ───────────────────────────────────────────────────
SIT_DATA = [
    {
        "num": 1, "label": "SITUATION 1",
        "title": "Estimation confidentielle - Mise a niveau Ouaoumana",
        "subtitle": "Competence C18 - Metrer et estimer",
        "accent": TURQ, "metric": "3,2%",
        "metric_label": "Ecart estimation / offre retenue",
        "metric_sub": "objectif < 5%",
        "ctx": "Estimation confidentielle Ouaoumana. 15,8 M DH HT - 35% budget. Prix plafond avant AO. Delai 3 semaines.",
        "pb":  "Plans incoherents. Mercuriale 2014 obsolete : prix enrobes sous-evalues 15 a 20%.",
        "act": "112 lignes metres AutoCAD + 4 visites terrain. Prix croises mercuriale + 5 marches similaires + devis fournisseurs.",
        "res": "Ecart 3,2% vs 5-10% habituel. Livre en 3 semaines. Methode adoptee standard Agence.",
        "chart_type": "column",
        "chart_cats": ["Estimation confidentielle", "Offre retenue", "Ecart"],
        "chart_vals": (15.8, 15.29, 0.51),
        "chart_color": "00959E",
        "chart_title": "Estimation vs offre (M DH HT)",
    },
    {
        "num": 2, "label": "SITUATION 2",
        "title": "Analyse des offres - Commission AO Ouaoumana",
        "subtitle": "Competence C18 - Analyser les offres",
        "accent": GOLD, "metric": "94/100",
        "metric_label": "Note entreprise retenue",
        "metric_sub": "0 recours depose",
        "ctx": "Membre technique commission. 3 entreprises soumissionnaires. Analyse comparative conformite + prix.",
        "pb":  "Dossier A : 7 erreurs arithmetiques. Dossier B : prix bas anormaux 42% du montant, ecarts 25-33%.",
        "act": "Correction arithmetique. Demande explication ecrite entreprise B. Rapport analyse + PV commission.",
        "res": "Attribution 15 jours. 0 recours. 94/100. Rapport valide sans reserve par Directeur.",
        "chart_type": "column",
        "chart_cats": ["Entreprise A (retenue)", "Entreprise B (prix bas)", "Entreprise C (non conf.)"],
        "chart_vals": (94, 72, 0),
        "chart_color": "F5A118",
        "chart_title": "Notes attribuees /100",
    },
    {
        "num": 3, "label": "SITUATION 3",
        "title": "Suivi financier - Route Lehri-Kerrouchen",
        "subtitle": "Competence C18 - Suivre le budget",
        "accent": VERT, "metric": "+0,8%",
        "metric_label": "Depassement final maitrise",
        "metric_sub": "Avenant evite",
        "ctx": "Kerrouchen 7,3 M DH 18 mois. Suivi mensuel situations travaux et quantites reelles.",
        "pb":  "Mi-parcours : chaussee +12% terrain rocheux, murs +15%. Depassement projete +4,8% = 349 000 DH. Seuil avenant 5%.",
        "act": "Tableau de bord hebdomadaire. Attachements contradictoires. Compensation Paysager -44k + Mobilier -12k.",
        "res": "Depassement final +0,8% (292 000 DH sous seuil). Aucun avenant. Tableau adopte 3 autres communes.",
        "chart_type": "column",
        "chart_cats": ["Chaussee", "Murs", "Trottoirs", "Eclairage", "Paysager*", "Mobilier*", "GLOBAL"],
        "chart_vals": (12, 15, 2.5, 3.2, -8, -10, 0.8),
        "chart_color": "27AE60",
        "chart_title": "Derive budgetaire par poste (%)",
    },
    {
        "num": 4, "label": "SITUATION 4",
        "title": "Coordination multi-sites - Province Khenifra",
        "subtitle": "Competence C19 - Coordonner et communiquer",
        "accent": ORANGE, "metric": "48h",
        "metric_label": "3 crises resolues",
        "metric_sub": "4 sites 20-80 km",
        "ctx": "Relais unique Directeur pour 4 chantiers VRD province Khenifra. Interface entreprise/BET/labo/hierarchie.",
        "pb":  "Semaine 23 : retard enrobes Ouaoumana, alerte meteo Kerrouchen, litige bordures +15% Sebt Ait Rahou.",
        "act": "48h : note Directeur + planning rattrapage. OS arret meteo + photos datees. Re-mesurage contradictoire : 15% vers 2,3%.",
        "res": "Retard rattrape 2 semaines. Litige resolu paiement debloque. Confie redaction tous OS et CR.",
        "chart_type": "column",
        "chart_cats": ["El Hammam", "Kerrouchen", "Ouaoumana", "Sebt Ait Rahou"],
        "chart_vals": (98, 95, 100, 97),
        "chart_color": "F3921A",
        "chart_title": "Avancement chantiers Projet 1 (%)",
    },
    {
        "num": 5, "label": "SITUATION 5",
        "title": "Controle des cubatures - Route Lehri-Kerrouchen",
        "subtitle": "Competence C18 - Verifier les quantites",
        "accent": ROUGE, "metric": "-8,3%",
        "metric_label": "Ecart cubatures corrige",
        "metric_sub": "Paiement ajuste",
        "ctx": "Route Lehri-Kerrouchen 25 km zone montagneuse. Denivele 400 m pentes 12%. Calcaire fracture imprevue.",
        "pb":  "Etude geotechnique non representative. Calcaire fracture 85 DH/m3 vs 28. Quantites surestimees +8,3%.",
        "act": "Profils en travers tous 25 m. Releves GPS contradictoires tous 500 ml. Verification avant visa situations travaux.",
        "res": "Situations corrigees paiement ajuste quantites reelles. Alea documente. Tracabilite complete validee.",
        "chart_type": "bar",
        "chart_cats": ["Corps chaussee 30%", "Enrobes 16%", "Terrassements 18%", "Ouv. hydr. 14%", "Assainissement 12%", "Divers 10%"],
        "chart_vals": (7.3, 3.9, 4.4, 3.4, 2.9, 2.3),
        "chart_color": "E54E3C",
        "chart_title": "Repartition budget route (M DH HT)",
    },
]

def situation_slide(prs, data, slide_num):
    sl = blank_slide(prs)
    set_bg(sl, BG)
    header_band(sl, data["label"], data["title"], data["subtitle"])

    left_w = Inches(3.3)
    left_h = Inches(2.5)
    ly = Inches(1.25)

    # Left metric block
    R(sl, 0, ly, left_w, left_h, data["accent"])
    T(sl, data["metric"], Inches(0.1), ly + Inches(0.25), left_w - Inches(0.2), Inches(1.0),
      sz=46, bold=True, color=BLANC, align=PP_ALIGN.CENTER)
    T(sl, data["metric_label"], Inches(0.1), ly + Inches(1.25), left_w - Inches(0.2), Inches(0.6),
      sz=12, bold=True, color=BLANC, align=PP_ALIGN.CENTER, wrap=True)
    T(sl, data["metric_sub"], Inches(0.1), ly + Inches(1.9), left_w - Inches(0.2), Inches(0.35),
      sz=10, italic=True, color=BLANC, align=PP_ALIGN.CENTER)

    # CPAR row
    cpar_x = left_w + Inches(0.1)
    cpar_w = W - cpar_x - Inches(0.15)
    cpar_row(sl, data["ctx"], data["pb"], data["act"], data["res"],
             cpar_x, ly, cpar_w, Inches(2.5))

    # Chart bottom
    chart_y = ly + Inches(2.6)
    chart_h = H - chart_y - Inches(0.35)
    horizontal = (data["chart_type"] == "bar")
    bar_chart(sl, data["chart_title"],
              data["chart_cats"], data["chart_vals"],
              cpar_x, chart_y, cpar_w, chart_h,
              color=data["chart_color"], horizontal=horizontal)

    footer(sl, slide_num)

# ── SLIDE 9 – SYNTHESE ────────────────────────────────────────────────────────
def slide9(prs):
    sl = blank_slide(prs)
    set_bg(sl, BG)
    header_band(sl, "BILAN", "Synthese des competences - C18 et C19")

    # Table
    table_data = [
        ["Sit.1", "Estimation confidentielle Ouaoumana", "C18 - Metres et estimation", "Ecart 3,2% - methode adoptee"],
        ["Sit.2", "Analyse offres commission AO", "C18 - Analyse des offres", "94/100 - 0 recours"],
        ["Sit.3", "Suivi financier Kerrouchen", "C18 - Suivi budgetaire", "+0,8% - avenant evite"],
        ["Sit.4", "Coordination multi-sites", "C19 - Coordination", "3 crises - 48h resolution"],
        ["Sit.5", "Controle cubatures route", "C18 - Controle quantites", "Paiement ajuste - tracabilite"],
    ]
    headers = ["Sit", "Activite", "Competence", "Resultat cle"]
    col_ws = [Inches(0.7), Inches(3.3), Inches(3.5), Inches(3.5)]
    tx = Inches(0.3)
    ty = Inches(1.3)
    rh = Inches(0.38)

    # Header row
    cx = tx
    for ci, hdr in enumerate(headers):
        R(sl, cx, ty, col_ws[ci], rh, TURQ)
        T(sl, hdr, cx + Inches(0.05), ty + Inches(0.05),
          col_ws[ci] - Inches(0.1), rh - Inches(0.1),
          sz=11, bold=True, color=BLANC)
        cx += col_ws[ci]

    for ri, row in enumerate(table_data):
        bg = CARD if ri % 2 == 0 else BLANC
        cx = tx
        for ci, cell in enumerate(row):
            R(sl, cx, ty + (ri + 1) * rh, col_ws[ci], rh, bg)
            col_color = VERT if ci == 3 else NAVY
            bold = (ci == 3)
            T(sl, cell, cx + Inches(0.05),
              ty + (ri + 1) * rh + Inches(0.05),
              col_ws[ci] - Inches(0.1), rh - Inches(0.1),
              sz=10, bold=bold, color=col_color)
            cx += col_ws[ci]

    # Bar chart left bottom
    chart_y = Inches(4.1)
    chart_h = H - chart_y - Inches(0.35)
    bar_chart(sl, "Auto-evaluation /5",
              ["Metres", "Analyse offres", "Suivi finan.", "Coordination", "Controle quant.", "Marches pub."],
              (4.5, 4.5, 4.8, 4.2, 4.3, 4.6),
              Inches(0.3), chart_y, Inches(6.0), chart_h,
              color="00959E")

    # CARD2 box right bottom
    bx = Inches(6.7)
    bw2 = W - bx - Inches(0.3)
    R(sl, bx, chart_y, bw2, chart_h, CARD2)
    R(sl, bx, chart_y, bw2, Inches(0.06), GOLD)
    T(sl, "Comparaison Maroc / France", bx + Inches(0.15), chart_y + Inches(0.1),
      bw2 - Inches(0.3), Inches(0.35), sz=13, bold=True, color=NAVY)
    bullets = [
        "Maroc : MOA publique, vision globale projet, AO, commission",
        "France : terrain operationnel, planning, sous-traitants, PV",
        "BIM : pont entre conception et execution - efficacite metres",
        "BIMCO : valorisation des deux cultures professionnelles",
    ]
    for i, b in enumerate(bullets):
        T(sl, "• " + b, bx + Inches(0.15),
          chart_y + Inches(0.55) + i * Inches(0.45),
          bw2 - Inches(0.3), Inches(0.4), sz=10, color=NAVY)

    footer(sl, 9)

# ── SLIDE 10 – BIM + PROJET PRO ───────────────────────────────────────────────
def slide10(prs):
    sl = blank_slide(prs)
    set_bg(sl, BG)
    header_band(sl, "BIMCO ET PERSPECTIVES",
                "Protocole BIM - Mon projet professionnel")

    half = (W - Inches(0.6)) / 2

    # Left: BIM
    lx = Inches(0.2)
    ly = Inches(1.3)
    lh = H - ly - Inches(0.35)
    R(sl, lx, ly, half, lh, BLANC)
    R(sl, lx, ly, half, Inches(0.07), TURQ)
    T(sl, "Protocole BIM - BIMCO", lx + Inches(0.15), ly + Inches(0.1),
      half - Inches(0.3), Inches(0.35), sz=14, bold=True, color=TURQ)

    bim_items = [
        ("Convention BIM", "Charte projet, niveaux de developpement LOD"),
        ("Format echange", "IFC / RVT Revit - interoperabilite"),
        ("Extraction metres", "Quantitatifs automatises depuis maquette"),
        ("Gain de temps", "Estimation 3x plus rapide vs AutoCAD"),
        ("Outils", "Revit, BIM360, Navisworks - formation AFPA"),
    ]
    for i, (k, v) in enumerate(bim_items):
        iy = ly + Inches(0.55) + i * Inches(0.42)
        T(sl, k + " : ", lx + Inches(0.15), iy, Inches(1.7), Inches(0.35),
          sz=10, bold=True, color=TURQ)
        T(sl, v, lx + Inches(1.85), iy, half - Inches(2.05), Inches(0.35),
          sz=10, color=NAVY)

    chart_y2 = ly + Inches(2.7)
    chart_h2 = lh - Inches(2.8)
    bar_chart(sl, "Gain BIM vs traditionnel base 100",
              ["Metres", "Fiabilite", "Temps"],
              (60, 115, 65),
              lx, chart_y2, half, chart_h2,
              color="00959E")

    # Right: Projet pro
    rx = lx + half + Inches(0.2)
    R(sl, rx, ly, half, lh, CARD2)
    R(sl, rx, ly, half, Inches(0.07), GOLD)
    T(sl, "Mon projet professionnel", rx + Inches(0.15), ly + Inches(0.1),
      half - Inches(0.3), Inches(0.35), sz=14, bold=True, color=GOLD)

    phases = [
        ("2026 - Court terme", TURQ, [
            "Obtenir BTS MEC - validation acquis",
            "Developper BIMCO - premiers clients",
            "Maquettes BIM chantiers VRD",
        ]),
        ("2027-2028 - Moyen terme", GOLD, [
            "Licence pro gestion de projet BTP",
            "Partenariat BET et maitres d'oeuvre",
            "Logiciel estimation BIM integree",
        ]),
        ("Long terme", TURQD, [
            "Expert BIM reconnu MOE / AMO",
            "Formation jeunes techniciens BTP",
            "Pont Maroc - France - expertise BIM",
        ]),
    ]
    py = ly + Inches(0.55)
    for phase_label, col, items in phases:
        R(sl, rx + Inches(0.1), py, half - Inches(0.2), Inches(0.3), col)
        T(sl, phase_label, rx + Inches(0.15), py + Inches(0.02),
          half - Inches(0.3), Inches(0.26), sz=11, bold=True, color=BLANC)
        py += Inches(0.35)
        for item in items:
            T(sl, "  • " + item, rx + Inches(0.15), py,
              half - Inches(0.3), Inches(0.3), sz=10, color=NAVY)
            py += Inches(0.32)
        py += Inches(0.15)

    footer(sl, 10)

# ── SLIDE 11 – BILAN REFLEXIF ─────────────────────────────────────────────────
def slide11(prs):
    sl = blank_slide(prs)
    set_bg(sl, BG)
    header_band(sl, "BILAN REFLEXIF", "Ce que le terrain m'a appris")

    cards = [
        ("Ce que j'ai appris", TURQ, "C",
         "Semaine 23 : gerer 3 crises simultanees a 80 km d'ecart m'a oblige a prioriser, formaliser et communiquer avec clarte. L'urgence structure la rigueur."),
        ("Ce que j'aurais fait autrement", ROUGE, "A",
         "Mettre en place le tableau de bord des le debut, pas a mi-parcours. Anticiper les derives avant qu'elles deviennent critiques."),
        ("Ce que je maitrise", VERT, "M",
         "Le lien technique / financier : metres, situations de travaux, quantites reelles. Passage du terrain au document contractuel."),
        ("Ce que ca m'a apporte", GOLD, "P",
         "BIMCO est ne de ce parcours. Le BIM est l'outil qui manquait pour industrialiser la methode metres + estimation que j'ai developpe sur le terrain."),
    ]

    cw2 = (W - Inches(0.5)) / 4
    for i, (lbl, col, letter, txt) in enumerate(cards):
        cx = Inches(0.2) + i * (cw2 + Inches(0.1))
        cy = Inches(1.25)
        ch2 = Inches(4.3)
        R(sl, cx, cy, cw2, ch2, BLANC)
        R(sl, cx, cy, cw2, Inches(0.07), col)

        c = sl.shapes.add_shape(9, cx + Inches(0.1), cy + Inches(0.15),
                                Inches(0.4), Inches(0.4))
        c.fill.solid(); c.fill.fore_color.rgb = col; c.line.fill.background()
        T(sl, letter, cx + Inches(0.1), cy + Inches(0.15),
          Inches(0.4), Inches(0.4), sz=16, bold=True, color=BLANC,
          align=PP_ALIGN.CENTER)

        T(sl, lbl, cx + Inches(0.6), cy + Inches(0.18),
          cw2 - Inches(0.7), Inches(0.4), sz=11, bold=True, color=col)

        R(sl, cx + Inches(0.1), cy + Inches(0.65), cw2 - Inches(0.2),
          Inches(0.02), GRISL)

        T(sl, txt, cx + Inches(0.1), cy + Inches(0.75),
          cw2 - Inches(0.2), ch2 - Inches(0.85), sz=10, color=NAVY, wrap=True)

    # Bar chart bottom
    chart_y = Inches(5.65)
    chart_h = H - chart_y - Inches(0.35)
    bar_chart(sl, "Progression percue /10",
              ["Sit.1", "Sit.2", "Sit.3", "Sit.4", "Sit.5"],
              (7, 7.5, 8.5, 8, 8.5),
              Inches(0.2), chart_y, W - Inches(0.4), chart_h,
              color="00959E")

    footer(sl, 11)

# ── SLIDE 12 – CONCLUSION ─────────────────────────────────────────────────────
def slide12(prs):
    sl = blank_slide(prs)
    set_bg(sl, DARK)
    img1 = os.path.join(MEDIA, "image1.jpeg")
    add_image(sl, img1, 0, 0, W, H)
    dark_overlay(sl, 80)

    R(sl, 0, 0, Inches(0.1), H, TURQ)
    R(sl, 0, 0, W, Inches(0.06), GOLD)

    T(sl, "En conclusion...", Inches(0.4), Inches(0.3), Inches(10), Inches(0.35),
      sz=14, italic=True, color=TURQL)
    T(sl, "5 competences. 2 projets. 82 millions de dirhams.",
      Inches(0.4), Inches(0.7), Inches(12), Inches(0.8),
      sz=44, bold=True, color=BLANC)
    R(sl, Inches(0.4), Inches(1.55), Inches(5), Inches(0.05), GOLD)

    bullets = [
        (TURQ,   "Estimer, controler, verifier",
                 "Metres 15,8 M DH . Ecart 3,2% . Cubatures corrigees"),
        (GOLD,   "Analyser, decider, formaliser",
                 "Commission AO . 94/100 . 0 recours . 15 jours"),
        (VERT,   "Suivre, anticiper, corriger",
                 "Tableau de bord . Avenant evite . +0,8% final"),
        (ORANGE, "Coordonner, communiquer",
                 "4 sites . 3 crises . 48h . OS + CR systematiques"),
    ]

    for i, (col, main_txt, sub_txt) in enumerate(bullets):
        by = Inches(1.75) + i * Inches(0.85)
        c = sl.shapes.add_shape(9, Inches(0.4), by + Inches(0.05),
                                Inches(0.4), Inches(0.4))
        c.fill.solid(); c.fill.fore_color.rgb = col; c.line.fill.background()
        T(sl, main_txt, Inches(0.95), by + Inches(0.08),
          Inches(9), Inches(0.35), sz=16, bold=True, color=BLANC)
        T(sl, sub_txt, Inches(0.95), by + Inches(0.44),
          Inches(9), Inches(0.3), sz=11, color=TURQL)

    R(sl, Inches(0.4), Inches(5.3), W - Inches(0.8), Inches(0.04), TURQ)

    T(sl, "Le BIM est le prolongement naturel de ce parcours.",
      Inches(0.4), Inches(5.42), Inches(12), Inches(0.35),
      sz=16, italic=True, color=TURQL)
    T(sl, "Merci pour votre attention.",
      Inches(0.4), Inches(5.9), Inches(12), Inches(0.55),
      sz=26, bold=True, color=BLANC)
    T(sl, "Je suis disponible pour repondre a vos questions.",
      Inches(0.4), Inches(6.5), Inches(12), Inches(0.35),
      sz=15, color=TURQL)

# ── BUILD ─────────────────────────────────────────────────────────────────────
def build():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    print("Slide 1 - Cover...")
    slide1(prs)

    print("Slide 2 - Profil...")
    slide2(prs)

    print("Slide 3 - Vue ensemble...")
    slide3(prs)

    for idx, data in enumerate(SIT_DATA):
        sn = 4 + idx
        print(f"Slide {sn} - Situation {data['num']}...")
        situation_slide(prs, data, sn)

    print("Slide 9 - Synthese...")
    slide9(prs)

    print("Slide 10 - BIM + Projet pro...")
    slide10(prs)

    print("Slide 11 - Bilan reflexif...")
    slide11(prs)

    print("Slide 12 - Conclusion...")
    slide12(prs)

    out = os.path.join(os.path.dirname(__file__),
                       "SOUTENANCE_16x9_GRAND_ECRAN.pptx")
    prs.save(out)
    print(f"\nSaved: {out}")
    print(f"Slides: {len(prs.slides)}")

if __name__ == "__main__":
    build()
