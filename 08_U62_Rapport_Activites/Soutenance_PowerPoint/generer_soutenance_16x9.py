"""
Génère la présentation soutenance BTS MEC U62 — 16:9 grand écran
Fond turquoise + graphiques professionnels
"""
import sys, os
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from lxml import etree
import copy

# ─── PALETTE TURQUOISE ────────────────────────────────────────────────────────
BG          = RGBColor(0x00, 0x4D, 0x52)   # turquoise foncé (fond principal)
BG_CARD     = RGBColor(0x00, 0x38, 0x3D)   # cartes sombres
BG_DARK     = RGBColor(0x00, 0x28, 0x2C)   # très sombre
TURQ        = RGBColor(0x00, 0xC8, 0xD0)   # turquoise vif (accent principal)
TURQ_MID    = RGBColor(0x40, 0xD8, 0xDE)   # turquoise moyen
TURQ_LIGHT  = RGBColor(0xA0, 0xEC, 0xF0)   # turquoise clair
TURQ_PALE   = RGBColor(0xD0, 0xF4, 0xF6)   # très clair
BLANC       = RGBColor(0xFF, 0xFF, 0xFF)
GRIS        = RGBColor(0xA0, 0xC4, 0xC6)
GOLD        = RGBColor(0xF5, 0xC5, 0x18)
ROUGE       = RGBColor(0xFF, 0x55, 0x44)
VERT        = RGBColor(0x44, 0xCC, 0x88)
ORANGE      = RGBColor(0xFF, 0x99, 0x33)

W = Inches(13.333)
H = Inches(7.5)
MEDIA = "d:/PREPA BTS MEC/08_U62_Rapport_Activites/Soutenance_PowerPoint/extracted_media"

# ─── UTILITAIRES ─────────────────────────────────────────────────────────────
def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def add_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def bg(slide, color=BG):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color

def rect(slide, x, y, w, h, color, alpha_pct=None):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    if alpha_pct is not None:
        spPr = s._element.spPr
        sf = spPr.find('.//' + qn('a:solidFill'))
        if sf is not None:
            ch = sf[0] if len(sf) else None
            if ch is not None:
                ae = etree.SubElement(ch, qn('a:alpha'))
                ae.set('val', str(int(alpha_pct * 1000)))
    return s

def txt(slide, text, x, y, w, h, sz=18, bold=False, color=BLANC,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size   = Pt(sz)
    r.font.bold   = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name   = "Calibri"
    return tb

def txt_lines(slide, lines, x, y, w, h, spacing=4):
    """lines = [(text, size, bold, color, align)]"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for (t, sz, bold, col, al) in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = al
        p.space_before = Pt(spacing)
        r = p.add_run()
        r.text  = t
        r.font.size  = Pt(sz)
        r.font.bold  = bold
        r.font.color.rgb = col
        r.font.name  = "Calibri"
    return tb

def img(slide, path, x, y, w, h):
    if os.path.exists(path):
        return slide.shapes.add_picture(path, x, y, w, h)

def overlay(slide, alpha=60):
    """Overlay sombre sur la dernière image insérée."""
    s = rect(slide, 0, 0, W, H, BG_DARK, alpha_pct=alpha)
    return s

def footer(slide, num, total=15):
    fy = H - Inches(0.34)
    rect(slide, 0, fy, W, Inches(0.34), BG_DARK)
    txt(slide, "BAHAFID Mohamed  │  Rapport U62  │  BTS MEC 2026",
        Inches(0.3), fy + Pt(4), Inches(10), Inches(0.28), sz=10, color=GRIS)
    txt(slide, f"{num} / {total}", W - Inches(1.2), fy + Pt(4),
        Inches(1.0), Inches(0.28), sz=10, color=TURQ_LIGHT, align=PP_ALIGN.RIGHT)

def top_bar(slide, color=TURQ, h=Inches(0.07)):
    rect(slide, 0, 0, W, h, color)

def left_bar(slide, color=TURQ, w=Inches(0.07)):
    rect(slide, 0, 0, w, H, color)

def kpi_box(slide, val, label, x, y, w, h, val_color=GOLD):
    rect(slide, x, y, w, h, BG_CARD)
    rect(slide, x, y, w, Inches(0.04), TURQ)
    txt(slide, val,   x, y + Inches(0.1), w, Inches(0.55),
        sz=32, bold=True, color=val_color, align=PP_ALIGN.CENTER)
    txt(slide, label, x, y + Inches(0.6), w, Inches(0.38),
        sz=11, color=TURQ_LIGHT, align=PP_ALIGN.CENTER)

def cpar_col(slide, label, content, x, y, w, h, label_color=TURQ):
    BG_MAP = {
        "CONTEXTE":    RGBColor(0x00, 0x3A, 0x50),
        "PROBLÈME":    RGBColor(0x50, 0x18, 0x18),
        "ACTION":      RGBColor(0x18, 0x45, 0x28),
        "RÉSULTAT":    RGBColor(0x1A, 0x35, 0x10),
    }
    bg_col = BG_MAP.get(label, BG_CARD)
    rect(slide, x, y, w, h, bg_col)
    rect(slide, x, y, w, Inches(0.04), label_color)
    txt(slide, label, x + Inches(0.1), y + Inches(0.07),
        w - Inches(0.2), Inches(0.32), sz=12, bold=True, color=label_color)
    txt(slide, content, x + Inches(0.1), y + Inches(0.42),
        w - Inches(0.2), h - Inches(0.52), sz=12.5, color=BLANC)

# ─── GRAPHIQUES ──────────────────────────────────────────────────────────────
def add_pie_chart(slide, title, categories, values, colors_hex, x, y, w, h):
    """Camembert coloré avec légende intégrée."""
    cd = ChartData()
    cd.categories = categories
    cd.add_series("", values)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, x, y, w, h, cd
    ).chart
    chart.has_title = True
    chart.chart_title.text_frame.text = title
    chart.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(13)
    chart.chart_title.text_frame.paragraphs[0].runs[0].font.bold = True
    chart.chart_title.text_frame.paragraphs[0].runs[0].font.color.rgb = BLANC

    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.show_percentage = True
    plot.data_labels.show_category_name = False
    plot.data_labels.number_format = '0%'
    for lbl in plot.data_labels._element.iter(qn('a:solidFill')):
        srgb = lbl.find(qn('a:srgbClr'))
        if srgb is not None:
            srgb.set('val', 'FFFFFF')

    # Couleurs des tranches
    for i, hex_col in enumerate(colors_hex):
        r2, g2, b2 = int(hex_col[0:2],16), int(hex_col[2:4],16), int(hex_col[4:6],16)
        pt = plot.series[0].points[i]
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = RGBColor(r2, g2, b2)

    # Fond du graphique transparent


    chart.has_legend = True
    chart.legend.position = -4152  # RIGHT
    chart.legend.include_in_layout = False
    # Couleur légende
    leg_elem = chart.legend._element
    for txPr in leg_elem.iter(qn('a:solidFill')):
        sc = txPr.find(qn('a:srgbClr'))
        if sc is not None:
            sc.set('val', 'A0C4C6')
    return chart

def add_bar_chart(slide, title, categories, series_data, x, y, w, h,
                  bar_colors=None, show_val=True, gap=80):
    """Histogramme horizontal ou vertical."""
    cd = ChartData()
    cd.categories = categories
    for name, vals in series_data:
        cd.add_series(name, vals)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, w, h, cd
    ).chart

    chart.has_title = True
    chart.chart_title.text_frame.text = title
    tf_para = chart.chart_title.text_frame.paragraphs[0]
    if tf_para.runs:
        tf_para.runs[0].font.size = Pt(12)
        tf_para.runs[0].font.bold = True
        tf_para.runs[0].font.color.rgb = BLANC

    plot = chart.plots[0]
    plot.gap_width = gap
    plot.has_data_labels = show_val
    if show_val:
        plot.data_labels.show_value = True
        plot.data_labels.number_format = '0.0"%"' if any('%' in str(v) for _, vals in series_data for v in vals) else 'General'

    if bar_colors:
        for si, (_, _vals) in enumerate(series_data):
            ser = plot.series[si]
            if si < len(bar_colors):
                hc = bar_colors[si]
                r2,g2,b2 = int(hc[0:2],16),int(hc[2:4],16),int(hc[4:6],16)
                ser.format.fill.solid()
                ser.format.fill.fore_color.rgb = RGBColor(r2,g2,b2)



    chart.has_legend = (len(series_data) > 1)

    # Axe couleur
    va = chart.value_axis
    va.tick_labels.font.size = Pt(10)
    va.tick_labels.font.color.rgb = TURQ_LIGHT
    ca = chart.category_axis
    ca.tick_labels.font.size = Pt(10)
    ca.tick_labels.font.color.rgb = TURQ_LIGHT
    return chart

def add_bar_horizontal(slide, title, categories, values, x, y, w, h,
                        bar_color="00C8D0"):
    cd = ChartData()
    cd.categories = categories
    cd.add_series("", values)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, x, y, w, h, cd
    ).chart
    chart.has_title = True
    chart.chart_title.text_frame.text = title
    if chart.chart_title.text_frame.paragraphs[0].runs:
        r = chart.chart_title.text_frame.paragraphs[0].runs[0]
        r.font.size = Pt(12)
        r.font.bold = True
        r.font.color.rgb = BLANC

    plot = chart.plots[0]
    plot.gap_width = 60
    plot.has_data_labels = True
    plot.data_labels.show_value = True

    ser = plot.series[0]
    r2,g2,b2 = int(bar_color[0:2],16),int(bar_color[2:4],16),int(bar_color[4:6],16)
    ser.format.fill.solid()
    ser.format.fill.fore_color.rgb = RGBColor(r2,g2,b2)



    chart.has_legend = False

    va = chart.value_axis
    va.tick_labels.font.color.rgb = TURQ_LIGHT
    va.tick_labels.font.size = Pt(9)
    ca = chart.category_axis
    ca.tick_labels.font.color.rgb = TURQ_LIGHT
    ca.tick_labels.font.size = Pt(10)
    return chart

# ─── SLIDES SITUATION CPAR ───────────────────────────────────────────────────
def situation_slide(prs, num_sit, title, metric, metric_label,
                    contexte, problematique, action, resultat,
                    competence, slide_num, total=15):
    sl = add_slide(prs)
    bg(sl, BG)
    top_bar(sl, TURQ, Inches(0.07))

    # Badge numéro
    rect(sl, Inches(0.1), Inches(0.1), Inches(0.72), Inches(0.72), TURQ)
    txt(sl, str(num_sit), Inches(0.1), Inches(0.12), Inches(0.72), Inches(0.62),
        sz=30, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)

    txt(sl, f"SITUATION {num_sit}", Inches(0.94), Inches(0.1),
        Inches(8), Inches(0.35), sz=11, bold=True, color=TURQ_MID)
    txt(sl, title, Inches(0.94), Inches(0.4), Inches(9.0), Inches(0.55),
        sz=25, bold=True, color=BLANC)

    # Métrique (haut droite)
    rect(sl, W - Inches(3.1), Inches(0.08), Inches(2.95), Inches(0.95), BG_DARK)
    rect(sl, W - Inches(3.1), Inches(0.08), Inches(2.95), Inches(0.05), GOLD)
    txt(sl, metric, W - Inches(3.05), Inches(0.12), Inches(2.85), Inches(0.52),
        sz=36, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    txt(sl, metric_label, W - Inches(3.05), Inches(0.64), Inches(2.85), Inches(0.32),
        sz=10.5, color=TURQ_LIGHT, align=PP_ALIGN.CENTER)

    # Bande compétence
    rect(sl, 0, Inches(1.05), W, Inches(0.33), RGBColor(0x00, 0x38, 0x40))
    txt(sl, f"▸  {competence}", Inches(0.25), Inches(1.08),
        W - Inches(0.4), Inches(0.28), sz=12, bold=True, color=TURQ_PALE)

    # 4 blocs CPAR
    bw = (W - Inches(0.55)) / 4 - Inches(0.1)
    bh = H - Inches(1.58)
    by = Inches(1.46)
    labels   = ["CONTEXTE", "PROBLÈME", "ACTION", "RÉSULTAT"]
    contents = [contexte, problematique, action, resultat]
    lcolors  = [TURQ_MID, ROUGE, VERT, GOLD]
    for i, (lbl, cnt, lc) in enumerate(zip(labels, contents, lcolors)):
        bx = Inches(0.15) + i * (bw + Inches(0.12))
        cpar_col(sl, lbl, cnt, bx, by, bw, bh, label_color=lc)

    footer(sl, slide_num, total)
    return sl

# ─── CONSTRUCTION COMPLÈTE ───────────────────────────────────────────────────
def build():
    prs = new_prs()
    TOTAL = 15

    # ═══════════════════════════════════════════════════════════════════════
    # 1. COVER
    # ═══════════════════════════════════════════════════════════════════════
    sl = add_slide(prs)
    bg(sl, BG)
    im1 = os.path.join(MEDIA, "image1.jpeg")
    if os.path.exists(im1):
        img(sl, im1, 0, 0, W, H)
        overlay(sl, alpha=72)

    left_bar(sl, GOLD, Inches(0.1))
    top_bar(sl, TURQ, Inches(0.07))

    txt(sl, "ACADÉMIE DE LYON", Inches(0.3), Inches(0.3),
        Inches(8), Inches(0.4), sz=13, color=TURQ_MID)
    txt(sl, "RAPPORT D'ACTIVITÉS", Inches(0.3), Inches(1.0),
        Inches(11), Inches(1.1), sz=60, bold=True, color=BLANC)
    txt(sl, "PROFESSIONNELLES", Inches(0.3), Inches(2.05),
        Inches(11), Inches(0.85), sz=50, bold=True, color=TURQ)

    rect(sl, Inches(0.3), Inches(3.05), Inches(5.8), Inches(0.05), TURQ)

    txt(sl, "BTS Management Économique de la Construction",
        Inches(0.3), Inches(3.22), Inches(9), Inches(0.45), sz=20, color=GRIS)
    txt(sl, "SESSION 2026",
        Inches(0.3), Inches(3.72), Inches(5), Inches(0.45), sz=22, bold=True, color=TURQ_MID)
    txt(sl, "BAHAFID Mohamed",
        Inches(0.3), Inches(4.55), Inches(9), Inches(0.7), sz=40, bold=True, color=BLANC)
    txt(sl, "N° Candidat  02537399911  │  Académie de Lyon",
        Inches(0.3), Inches(5.28), Inches(9), Inches(0.38), sz=16, color=GRIS)
    rect(sl, Inches(0.3), Inches(5.8), Inches(5.8), Inches(0.04), TURQ)
    txt(sl, "BIMCO — Projeteur BIM & Économiste de la Construction",
        Inches(0.3), Inches(5.95), Inches(9), Inches(0.38), sz=14, color=TURQ_LIGHT, italic=True)

    footer(sl, 1, TOTAL)

    # ═══════════════════════════════════════════════════════════════════════
    # 2. PROFIL CANDIDAT + graphique donut parcours
    # ═══════════════════════════════════════════════════════════════════════
    sl = add_slide(prs)
    bg(sl, BG)
    top_bar(sl, TURQ)
    left_bar(sl, TURQ)

    txt(sl, "PROFIL CANDIDAT", Inches(0.3), Inches(0.1),
        Inches(8), Inches(0.45), sz=13, bold=True, color=TURQ_MID)
    txt(sl, "BAHAFID Mohamed", Inches(0.3), Inches(0.52),
        Inches(9), Inches(0.65), sz=36, bold=True, color=BLANC)

    im2 = os.path.join(MEDIA, "image2.jpeg")
    if os.path.exists(im2):
        img(sl, im2, W - Inches(2.9), Inches(0.45), Inches(2.6), Inches(3.3))

    infos = [
        ("N° Candidat",       "02537399911"),
        ("Académie",          "Lyon"),
        ("Structure d'accueil", "Conseil Régional — Béni Mellal-Khénifra (Maroc)"),
        ("Poste occupé",      "Technicien Études et Suivi des Travaux"),
        ("Expérience BTP",    "8 ans  (3 ans Maroc + 5 ans France)"),
        ("Formation BIM",     "Technicien Modeleur BIM — AFPA Colmar (8 mois)"),
        ("Activité actuelle", "BIMCO — Projeteur BIM / Économiste de la construction"),
        ("SIREN / APE",       "999580053 / 7112B — Ingénierie, études techniques"),
    ]
    y0 = Inches(1.38)
    for label, val in infos:
        txt(sl, label, Inches(0.25), y0, Inches(2.9), Inches(0.28),
            sz=10, bold=True, color=TURQ_MID)
        txt(sl, val,   Inches(3.25), y0, Inches(6.2), Inches(0.28),
            sz=12.5, color=BLANC)
        rect(sl, Inches(0.25), y0 + Inches(0.27), Inches(9.15), Inches(0.01),
             RGBColor(0x00, 0x60, 0x65))
        y0 += Inches(0.36)

    # 4 KPI boxes
    kpis = [("8 ans", "Expérience BTP"), ("2 pays", "Maroc + France"),
            ("7 marchés", "Publics en MOA"), ("82,5 M DH", "Budget géré")]
    kw = (W - Inches(0.5)) / 4 - Inches(0.1)
    ky = H - Inches(1.52)
    for i, (v, l) in enumerate(kpis):
        kx = Inches(0.15) + i * (kw + Inches(0.12))
        kpi_box(sl, v, l, kx, ky, kw, Inches(1.18))

    # Graphique Donut — Répartition expérience (années)
    cd = ChartData()
    cd.categories = ["France (BTP)", "Maroc (MOA)", "Formation BIM"]
    cd.add_series("Années", (5, 3, 0.67))
    chart = sl.shapes.add_chart(
        XL_CHART_TYPE.DOUGHNUT, W - Inches(5.2), Inches(3.85), Inches(2.8), Inches(2.6), cd
    ).chart
    chart.has_title = False
    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.show_percentage = True
    plot.data_labels.show_category_name = False
    seg_colors = ["00C8D0", "F5C518", "44CC88"]
    for j, hc in enumerate(seg_colors):
        r2,g2,b2 = int(hc[0:2],16),int(hc[2:4],16),int(hc[4:6],16)
        plot.series[0].points[j].format.fill.solid()
        plot.series[0].points[j].format.fill.fore_color.rgb = RGBColor(r2,g2,b2)


    chart.has_legend = True
    chart.legend.position = -4152  # RIGHT

    txt(sl, "Répartition\nde l'expérience", W - Inches(5.2), Inches(3.72),
        Inches(2.8), Inches(0.42), sz=11, bold=True, color=TURQ_MID, align=PP_ALIGN.CENTER)

    footer(sl, 2, TOTAL)

    # ═══════════════════════════════════════════════════════════════════════
    # 3. CADRE PROFESSIONNEL
    # ═══════════════════════════════════════════════════════════════════════
    sl = add_slide(prs)
    bg(sl, BG)
    top_bar(sl, TURQ)

    txt(sl, "01  CADRE PROFESSIONNEL", Inches(0.3), Inches(0.1),
        Inches(10), Inches(0.4), sz=13, bold=True, color=TURQ_MID)
    txt(sl, "CONSEIL RÉGIONAL DE BÉNI MELLAL-KHÉNIFRA", Inches(0.3), Inches(0.47),
        Inches(9.5), Inches(0.58), sz=26, bold=True, color=BLANC)

    im7 = os.path.join(MEDIA, "image7.jpeg")
    if os.path.exists(im7):
        img(sl, im7, W - Inches(5.2), Inches(0.1), Inches(5.0), Inches(3.5))
        rect(sl, W - Inches(5.2), Inches(0.1), Inches(5.0), Inches(3.5),
             BG, alpha_pct=35)

    desc = ("Collectivité territoriale — 5 provinces — 2,5 millions d'habitants\n"
            "L'Agence d'Exécution des Projets assure la MOA des infrastructures régionales :\n"
            "routes, VRD, adduction d'eau potable, équipements publics.")
    txt(sl, desc, Inches(0.3), Inches(1.15), Inches(7.8), Inches(1.2), sz=14, color=GRIS)

    # Organigramme
    org_y = Inches(2.5)
    txt(sl, "MON POSTE dans l'organigramme", Inches(0.3), org_y,
        Inches(7.5), Inches(0.35), sz=12, bold=True, color=TURQ_MID)
    org = [
        ("Président du Conseil Régional",             RGBColor(0x00, 0x3A, 0x50)),
        ("Directeur Agence — M. A. DOGHMANI",         RGBColor(0x00, 0x4A, 0x58)),
        ("Service Marchés  │  Études  │  Suivi Travaux", RGBColor(0x00, 0x40, 0x48)),
        ("► Technicien Études & Suivi des Travaux  [MON POSTE]", RGBColor(0x00, 0x5C, 0x20)),
    ]
    oy = org_y + Inches(0.38)
    for k, (t, c) in enumerate(org):
        ox = Inches(0.3 + k * 0.5)
        ow = Inches(7.5) - Inches(k * 0.5)
        rect(sl, ox, oy, ow, Inches(0.38), c)
        txt(sl, t, ox + Inches(0.1), oy + Pt(2), ow - Inches(0.2), Inches(0.34),
            sz=12.5 if k < 3 else 13,
            bold=(k == 3), color=BLANC if k < 3 else GOLD)
        oy += Inches(0.44)

    # Cadre réglementaire
    rect(sl, Inches(8.1), Inches(2.5), Inches(5.0), Inches(2.9), BG_CARD)
    rect(sl, Inches(8.1), Inches(2.5), Inches(5.0), Inches(0.05), TURQ)
    txt(sl, "Cadre réglementaire — Marchés publics", Inches(8.22), Inches(2.6),
        Inches(4.8), Inches(0.35), sz=11, bold=True, color=TURQ_MID)
    rg = ("• Appel d'offres ouvert / restreint / négocié\n"
          "• CPS — Cahier des Prescriptions Spéciales\n"
          "• RC — Règlement de Consultation\n"
          "• BPDE — Bordereau des Prix / Détail Estimatif\n"
          "• Estimation confidentielle de l'administration")
    txt(sl, rg, Inches(8.22), Inches(3.0), Inches(4.75), Inches(2.2), sz=12.5, color=BLANC)

    # BIMCO + outils
    rect(sl, Inches(0.3), Inches(5.7), Inches(12.7), Inches(1.45), BG_DARK)
    rect(sl, Inches(0.3), Inches(5.7), Inches(12.7), Inches(0.05), GOLD)
    txt(sl, "BIMCO — Mon activité indépendante  (SIREN 999580053)", Inches(0.45), Inches(5.78),
        Inches(7), Inches(0.35), sz=12, bold=True, color=GOLD)
    txt(sl, "Revit · Navisworks · Dynamo · AutoCAD · MS Project · Excel avancé · Python · API Revit · React · Docker",
        Inches(0.45), Inches(6.18), Inches(12.4), Inches(0.35), sz=12, color=TURQ_LIGHT)
    txt(sl, "Métrés BIM · DPGF · Études de prix · Plugins Revit · Modélisation 3D · Coordination maquettes",
        Inches(0.45), Inches(6.57), Inches(12.4), Inches(0.35), sz=12, color=BLANC)

    footer(sl, 3, TOTAL)

    # ═══════════════════════════════════════════════════════════════════════
    # 4. PROJET 1 — INTRO + GRAPHIQUE budget par commune
    # ═══════════════════════════════════════════════════════════════════════
    sl = add_slide(prs)
    bg(sl, BG)
    top_bar(sl, TURQ)

    im31 = os.path.join(MEDIA, "image31.jpeg")
    if os.path.exists(im31):
        img(sl, im31, W - Inches(5.8), Inches(0.07), Inches(5.8), H - Inches(0.42))
        rect(sl, W - Inches(5.8), Inches(0.07), Inches(5.8), H - Inches(0.42),
             BG_DARK, alpha_pct=55)

    txt(sl, "02", Inches(0.2), Inches(0.08), Inches(2), Inches(1.0),
        sz=72, bold=True, color=TURQ)
    txt(sl, "PROJET 1", Inches(0.2), Inches(1.08), Inches(7), Inches(0.45),
        sz=14, bold=True, color=TURQ_MID)
    txt(sl, "MISE À NIVEAU DES CENTRES\nDE 4 COMMUNES", Inches(0.2), Inches(1.48),
        Inches(7), Inches(1.3), sz=30, bold=True, color=BLANC)

    fiche = [
        ("Marché",        "n°38-RBK-2017 (Lot 4) — Appel d'offres ouvert"),
        ("MOA",           "Conseil Régional de Béni Mellal-Khénifra"),
        ("Localisation",  "Province de Khénifra — 4 communes — 20 à 80 km"),
        ("Nature",        "Aménagement urbain — VRD — 8 corps d'état"),
        ("Montant",       "53,5 M DH TTC  (≈ 4,86 M€) | 44,6 M DH HT"),
    ]
    fy2 = Inches(3.0)
    for lbl, val in fiche:
        txt(sl, lbl,  Inches(0.22), fy2, Inches(2.0), Inches(0.3),
            sz=10, bold=True, color=TURQ_MID)
        txt(sl, val,  Inches(2.3),  fy2, Inches(5.3), Inches(0.3),
            sz=12.5, color=BLANC)
        rect(sl, Inches(0.22), fy2 + Inches(0.28), Inches(7.35), Inches(0.01),
             RGBColor(0x00, 0x6A, 0x70))
        fy2 += Inches(0.36)

    # ── Graphique camembert : Répartition budget par commune ──
    add_pie_chart(sl,
        title="Répartition budget 44,6 M DH HT par commune",
        categories=["Ouaoumana (35,5%)", "Sebt Ait Rahou (33,2%)", "El Hammam (18,5%)", "Kerrouchen (12,8%)"],
        values=(15.83, 14.82, 8.25, 5.70),
        colors_hex=["00C8D0", "F5C518", "44CC88", "FF9933"],
        x=W - Inches(5.5), y=Inches(0.1), w=Inches(5.2), h=Inches(4.0))

    # Bar chart : 8 corps d'état (montant approximatif)
    add_bar_chart(sl,
        title="Répartition par corps d'état (M DH HT)",
        categories=["Assainisse.", "Chaussée", "Trottoirs", "Signalis.", "Éclairage", "Murs", "Paysager", "Mobilier"],
        series_data=[("M DH", (9.8, 12.2, 6.4, 2.1, 5.8, 4.6, 2.5, 1.2))],
        x=W - Inches(5.5), y=Inches(4.1), w=Inches(5.2), h=Inches(3.05),
        bar_colors=["00C8D0"], gap=60)

    footer(sl, 4, TOTAL)

    # ═══════════════════════════════════════════════════════════════════════
    # 5. SITUATION 1 — ESTIMATION
    # ═══════════════════════════════════════════════════════════════════════
    situation_slide(prs,
        num_sit=1,
        title="Estimation confidentielle de l'administration — Ouaoumana",
        metric="3,2%", metric_label="Écart estimation / offre retenue",
        contexte=("J'ai établi l'estimation confidentielle (15,8 M DH HT — 35% du budget global). "
                  "Ce document fixe le prix plafond avant l'AO. Délai : 3 semaines. "
                  "Périmètre : 8 corps d'état, 4 communes, plans incomplets."),
        problematique=("Plans d'assainissement incohérents (profils non calés, regards mal positionnés). "
                       "Mercuriale 2014 obsolète : prix enrobés, canalisations et acier sous-évalués "
                       "de 15 à 20%."),
        action=("Phase 1 : 112 lignes métrés AutoCAD + 4 visites terrain (altimétries manquantes). "
                "Phase 2 : prix unitaires croisés — mercuriale actualisée + 5 marchés similaires + "
                "devis fournisseurs. Phase 3 : vérification contradictoire chef de projet."),
        resultat=("Livraison dans le délai. Écart avec l'offre retenue : 3,2% "
                  "(vs 5–10% habituellement). Méthode adoptée comme standard "
                  "par l'Agence pour les 3 autres communes du programme."),
        competence="C18 — RÉALISER DES MÉTRÉS TOUS CORPS D'ÉTAT ET ESTIMER UN OUVRAGE",
        slide_num=5, total=TOTAL)

    # ═══════════════════════════════════════════════════════════════════════
    # 6. SITUATION 2 — COMMISSION AO
    # ═══════════════════════════════════════════════════════════════════════
    situation_slide(prs,
        num_sit=2,
        title="Analyse des offres et Commission d'Appel d'Offres",
        metric="94/100", metric_label="Note attribuée — 0 recours déposé",
        contexte=("Membre technique de la commission — 3 soumissionnaires. "
                  "Responsable : vérification conformité administrative, cohérence des prix, "
                  "rédaction du rapport d'analyse fondant la décision d'attribution."),
        problematique=("Dossier A : 7 erreurs arithmétiques dans les sous-détails de prix. "
                       "Dossier B : prix anormalement bas sur 3 postes majeurs (42% du montant) "
                       "— écarts de 25 à 33% vs marché. Risque de sous-traitance non déclarée."),
        action=("Vérification conformité des 3 dossiers. Correction arithmétique "
                "(règle : prix unitaire prime sur total). Demande d'explication écrite "
                "à l'entreprise aux prix bas. Rapport d'analyse + PV de commission."),
        resultat=("Attribution en 15 jours — 0 recours. Entreprise notée 94/100 retenue. "
                  "Rapport et PV validés sans réserve par le Directeur. "
                  "Leçon : l'analyse rigoureuse en amont protège tout le processus."),
        competence="C19 — ANALYSER LES OFFRES ET PRÉPARER LA DÉCISION EN COMMISSION",
        slide_num=6, total=TOTAL)

    # ═══════════════════════════════════════════════════════════════════════
    # 7. SITUATION 3 — SUIVI FINANCIER + graphique dérive budgétaire
    # ═══════════════════════════════════════════════════════════════════════
    sl = add_slide(prs)
    bg(sl, BG)
    top_bar(sl, TURQ)

    rect(sl, Inches(0.1), Inches(0.1), Inches(0.72), Inches(0.72), TURQ)
    txt(sl, "3", Inches(0.1), Inches(0.12), Inches(0.72), Inches(0.62),
        sz=30, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)

    txt(sl, "SITUATION 3", Inches(0.94), Inches(0.1),
        Inches(8), Inches(0.35), sz=11, bold=True, color=TURQ_MID)
    txt(sl, "Suivi financier et tableau de bord — Kerrouchen",
        Inches(0.94), Inches(0.4), Inches(9.0), Inches(0.55), sz=25, bold=True, color=BLANC)

    rect(sl, W - Inches(3.1), Inches(0.08), Inches(2.95), Inches(0.95), BG_DARK)
    rect(sl, W - Inches(3.1), Inches(0.08), Inches(2.95), Inches(0.05), GOLD)
    txt(sl, "+0,8%", W - Inches(3.05), Inches(0.12), Inches(2.85), Inches(0.52),
        sz=36, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    txt(sl, "Dépassement final — avenant évité",
        W - Inches(3.05), Inches(0.64), Inches(2.85), Inches(0.32),
        sz=10.5, color=TURQ_LIGHT, align=PP_ALIGN.CENTER)

    rect(sl, 0, Inches(1.05), W, Inches(0.33), RGBColor(0x00, 0x38, 0x40))
    txt(sl, "▸  C19 — SUIVRE L'EXÉCUTION FINANCIÈRE ET ANTICIPER LES DÉRIVES",
        Inches(0.25), Inches(1.08), W - Inches(0.4), Inches(0.28),
        sz=12, bold=True, color=TURQ_PALE)

    # 2 blocs CPAR gauche (Contexte + Problème)
    bh2 = Inches(2.8)
    bw2 = Inches(3.7)
    cpar_col(sl, "CONTEXTE",
        "Chantier Kerrouchen : 7,3 M DH sur 18 mois. Suivi mensuel : "
        "vérification situations de travaux, contrôle quantités réelles, "
        "décomptes transmis au Directeur. Découverte terrain rocheux imprévu.",
        Inches(0.15), Inches(1.45), bw2, bh2, TURQ_MID)
    cpar_col(sl, "PROBLÈME",
        "À mi-parcours : chaussée +12% (terrain rocheux), murs +15% "
        "(fondations plus profondes). Dépassement global projeté : "
        "+4,8% = 349 000 DH. Seuil avenant = 5% — procédure de 3 à 6 mois.",
        Inches(3.95), Inches(1.45), bw2, bh2, ROUGE)

    # Graphique dérive budgétaire
    add_bar_chart(sl,
        title="Dérive budgétaire par poste (%) — Kerrouchen",
        categories=["Chaussée", "Murs", "Trottoirs", "Éclairage", "Paysager", "Mobilier", "GLOBAL"],
        series_data=[
            ("Dérive (%)", (12.0, 15.0, 2.5, 3.2, -8.0, -10.0, 0.8)),
        ],
        x=Inches(7.75), y=Inches(1.42), w=Inches(5.35), h=Inches(2.85),
        bar_colors=["FF5544"], gap=70)

    cpar_col(sl, "ACTION",
        "Tableau de bord hebdomadaire : avancement physique + consommation + écart prévisionnel. "
        "Attachements contradictoires terrain. Compensation inter-postes proposée au Directeur : "
        "Paysager –44 000 DH + Mobilier –12 000 DH = 56 000 DH dégagés.",
        Inches(0.15), Inches(4.33), bw2, Inches(2.82), VERT)
    cpar_col(sl, "RÉSULTAT",
        "Dépassement final : +0,8% (292 000 DH — sous le seuil). "
        "Chantier terminé sans avenant. Tableau de bord répliqué sur "
        "les 3 autres communes et adopté comme référence par l'Agence.",
        Inches(3.95), Inches(4.33), bw2, Inches(2.82), GOLD)

    # Schéma compensation inter-postes
    rect(sl, Inches(7.75), Inches(4.33), Inches(5.35), Inches(2.82), BG_CARD)
    rect(sl, Inches(7.75), Inches(4.33), Inches(5.35), Inches(0.05), GOLD)
    txt(sl, "Mécanisme de compensation inter-postes", Inches(7.9), Inches(4.4),
        Inches(5.1), Inches(0.35), sz=11, bold=True, color=GOLD)
    comp_data = [
        ("Chaussée (dépassement)", "+87 600 DH", ROUGE),
        ("Murs (dépassement)",     "+48 000 DH", ROUGE),
        ("Paysager (réduit)",      "–44 000 DH", VERT),
        ("Mobilier (réduit)",      "–12 000 DH", VERT),
        ("Dépassement net final",  "+0,8%",       GOLD),
    ]
    cy = Inches(4.82)
    for label, val, col in comp_data:
        txt(sl, label, Inches(7.9),  cy, Inches(3.5), Inches(0.36), sz=12, color=BLANC)
        txt(sl, val,   Inches(11.5), cy, Inches(1.5), Inches(0.36), sz=13, bold=True,
            color=col, align=PP_ALIGN.RIGHT)
        rect(sl, Inches(7.9), cy + Inches(0.35), Inches(5.1), Inches(0.01),
             RGBColor(0x00, 0x6A, 0x70))
        cy += Inches(0.41)

    footer(sl, 7, TOTAL)

    # ═══════════════════════════════════════════════════════════════════════
    # 8. SITUATION 4 — COORDINATION
    # ═══════════════════════════════════════════════════════════════════════
    situation_slide(prs,
        num_sit=4,
        title="Communication et coordination — 4 chantiers simultanés",
        metric="48h", metric_label="3 crises résolues en simultané",
        contexte=("Relais unique du Directeur pour 4 chantiers VRD (20 à 80 km). "
                  "Interface : entreprise, BET, laboratoire, hiérarchie. "
                  "Chaque décision tracée par écrit — OS, CR, notes."),
        problematique=("Semaine 23 — 3 crises simultanées : retard livraison enrobés (Ouaoumana), "
                       "alerte météo → arrêt chantier (Kerrouchen), litige quantités bordures "
                       "+15% (Sebt Ait Rahou)."),
        action=("48h : note factuelle Directeur + planning rattrapage enrobés. "
                "OS d'arrêt météo avec photos datées géolocalisées. "
                "Re-mesurage contradictoire bordures sur place → 15% → 2,3%. "
                "Point téléphonique quotidien 8h avec les 4 chefs de chantier."),
        resultat=("Retard rattrapé en 2 semaines. Arrêt formalisé, reprise ordonnée. "
                  "Écart bordures : 15% → 2,3% — paiement débloqué. "
                  "Confié la rédaction systématique de tous les OS et CR de chantier."),
        competence="C18 — COMMUNIQUER ET COORDONNER EN CONTEXTE MULTI-SITES",
        slide_num=8, total=TOTAL)

    # ═══════════════════════════════════════════════════════════════════════
    # 9. PROJET 2 — ROUTE + graphique budget
    # ═══════════════════════════════════════════════════════════════════════
    sl = add_slide(prs)
    bg(sl, BG)
    top_bar(sl, TURQ)

    im39 = os.path.join(MEDIA, "image39.jpeg")
    if os.path.exists(im39):
        img(sl, im39, W - Inches(5.6), Inches(0.07), Inches(5.6), H - Inches(0.42))
        rect(sl, W - Inches(5.6), Inches(0.07), Inches(5.6), H - Inches(0.42),
             BG_DARK, alpha_pct=50)

    txt(sl, "03", Inches(0.2), Inches(0.08), Inches(2), Inches(1.0),
        sz=72, bold=True, color=TURQ)
    txt(sl, "PROJET 2", Inches(0.2), Inches(1.08), Inches(7), Inches(0.45),
        sz=14, bold=True, color=TURQ_MID)
    txt(sl, "ROUTE LEHRI-KERROUCHEN\n25 KM — Zone montagneuse", Inches(0.2), Inches(1.48),
        Inches(7), Inches(1.3), sz=28, bold=True, color=BLANC)

    fiche2 = [
        ("Marché",        "n°46-RBK-2017 — Programme PRR3"),
        ("Nature",        "Route rurale — Moyen Atlas — Dénivelé 400 m — Pentes 12%"),
        ("Montant",       "29 M DH TTC  (≈ 2,6 M€) | 24,2 M DH HT + TVA 20%"),
        ("MOA",           "Conseil Régional de Béni Mellal-Khénifra"),
    ]
    fy3 = Inches(3.0)
    for lbl, val in fiche2:
        txt(sl, lbl,  Inches(0.22), fy3, Inches(2.0), Inches(0.3),
            sz=10, bold=True, color=TURQ_MID)
        txt(sl, val,  Inches(2.3),  fy3, Inches(5.3), Inches(0.3),
            sz=12.5, color=BLANC)
        rect(sl, Inches(0.22), fy3 + Inches(0.28), Inches(7.35), Inches(0.01),
             RGBColor(0x00, 0x6A, 0x70))
        fy3 += Inches(0.36)

    # KPI route
    route_kpis = [("25 km", "Longueur totale"), ("400 m", "Dénivelé cumulé"),
                  ("12%", "Pente max en lacets"), ("24,2 M DH", "Budget HT")]
    rkw = Inches(7.3) / 4 - Inches(0.06)
    rky = Inches(4.55)
    for i, (v, l) in enumerate(route_kpis):
        rkx = Inches(0.15) + i * (rkw + Inches(0.08))
        kpi_box(sl, v, l, rkx, rky, rkw, Inches(1.08))

    # Graphique horizontal répartition budget route
    add_bar_horizontal(sl,
        title="Répartition budget Projet 2 (M DH HT)",
        categories=["Corps de chaussée", "Enrobés", "Terrassements", "Ouv. hydrauliques",
                    "Assainissement", "Divers"],
        values=(7.3, 3.9, 4.4, 3.4, 2.9, 2.3),
        x=W - Inches(5.4), y=Inches(3.5), w=Inches(5.2), h=Inches(3.6),
        bar_color="00C8D0")

    footer(sl, 9, TOTAL)

    # ═══════════════════════════════════════════════════════════════════════
    # 10. SITUATION 5 — CUBATURES
    # ═══════════════════════════════════════════════════════════════════════
    situation_slide(prs,
        num_sit=5,
        title="Contrôle des cubatures de terrassement — Zone montagneuse",
        metric="-8,3%", metric_label="Écart cubatures corrigé → paiement juste",
        contexte=("Route Lehri-Kerrouchen — 25 km — Moyen Atlas. "
                  "Contrôle des cubatures de terrassement sur terrain accidenté : "
                  "dénivelé 400 m, pentes 12%, calcaire fracturé imprévu."),
        problematique=("Étude géotechnique non représentative. Dès PK 4+200 : calcaire fracturé "
                       "(85 DH/m³ vs 28 DH/m³ ordinaire). Quantités déclarées par l'entreprise "
                       "surestimées de +8,3% — surcoût non justifié."),
        action=("Profils en travers tous les 25 m. Relevés GPS contradictoires "
                "tous les 500 ml avec le conducteur de travaux. "
                "Vérification systématique avant visa de chaque situation de travaux. "
                "Attachements datés et géolocalisés photographiés."),
        resultat=("Situations corrigées — paiement ajusté aux quantités réelles exécutées. "
                  "Aléas géotechniques documentés. Traçabilité complète. "
                  "Méthode de contrôle adaptée aux reliefs montagneux validée par le Directeur."),
        competence="C18 — CONTRÔLER LES QUANTITÉS ET VÉRIFIER LES SITUATIONS DE TRAVAUX",
        slide_num=10, total=TOTAL)

    # ═══════════════════════════════════════════════════════════════════════
    # 11. ACTIVITÉS COMPLÉMENTAIRES
    # ═══════════════════════════════════════════════════════════════════════
    sl = add_slide(prs)
    bg(sl, BG)
    top_bar(sl, TURQ)

    txt(sl, "ACTIVITÉS COMPLÉMENTAIRES", Inches(0.3), Inches(0.1),
        Inches(10), Inches(0.42), sz=13, bold=True, color=TURQ_MID)
    txt(sl, "5 Autres marchés publics — Cycle complet de gestion MOA",
        Inches(0.3), Inches(0.48), Inches(12.5), Inches(0.55), sz=26, bold=True, color=BLANC)

    # Tableau marchés
    cols = ["N° Marché", "Objet des travaux", "Montant", "Mission principale"]
    marches = [
        ("M01", "Route rurale — Province de Khenifra", "18 M DH", "Métrés + suivi financier"),
        ("M02", "Piste rurale — Zone Ait Sghir",       "6 M DH",  "Réception provisoire + réserves"),
        ("M03", "AEP rurale — Commune Ouaoumana",      "12 M DH", "Suivi technique et décomptes"),
        ("M04", "École — Province Béni Mellal",         "9 M DH",  "Vérification situations travaux"),
        ("M05", "Terrain de sport — Kerrouchen",        "4 M DH",  "Contrôle réception travaux"),
    ]
    cw2 = [Inches(1.5), Inches(5.3), Inches(1.8), Inches(4.6)]
    hx2 = Inches(0.2)
    hy2 = Inches(1.18)
    for j, (h2, cw) in enumerate(zip(cols, cw2)):
        rect(sl, hx2, hy2, cw, Inches(0.38), RGBColor(0x00, 0x60, 0x65))
        txt(sl, h2, hx2 + Inches(0.06), hy2 + Pt(3), cw - Inches(0.1), Inches(0.32),
            sz=11, bold=True, color=BLANC, align=PP_ALIGN.CENTER)
        hx2 += cw + Inches(0.06)
    ry2 = hy2 + Inches(0.42)
    for k, (nm, ob, mt, mi) in enumerate(marches):
        rbg2 = BG_CARD if k % 2 == 0 else RGBColor(0x00, 0x3A, 0x3F)
        rx2 = Inches(0.2)
        for j, (cell, cw) in enumerate(zip([nm, ob, mt, mi], cw2)):
            rect(sl, rx2, ry2, cw, Inches(0.46), rbg2)
            txt(sl, cell, rx2 + Inches(0.06), ry2 + Pt(3), cw - Inches(0.1), Inches(0.4),
                sz=12.5, color=GOLD if j == 0 else BLANC, bold=(j == 0))
            rx2 += cw + Inches(0.06)
        ry2 += Inches(0.49)

    # Graphique total des marchés (bar chart)
    add_bar_chart(sl,
        title="Total marchés gérés en MOA (M DH HT)",
        categories=["Mise à niveau\n4 communes", "Route\nLehri-Kerr.", "Route\nrurale M01",
                    "AEP\nM03", "École\nM04", "Piste\nM02", "Sport\nM05"],
        series_data=[("M DH HT", (44.6, 24.2, 18.0, 12.0, 9.0, 6.0, 4.0))],
        x=Inches(0.2), y=Inches(4.5), w=Inches(7.5), h=Inches(2.65),
        bar_colors=["00C8D0"], gap=50)

    # Résumé
    rect(sl, Inches(7.85), Inches(4.5), Inches(5.3), Inches(2.65), BG_DARK)
    rect(sl, Inches(7.85), Inches(4.5), Inches(5.3), Inches(0.05), GOLD)
    txt(sl, "BILAN GLOBAL", Inches(8.0), Inches(4.58), Inches(5.0), Inches(0.35),
        sz=13, bold=True, color=GOLD)
    bilan = [
        ("7 marchés publics gérés", BLANC),
        ("Cycle complet : AO → Réception", BLANC),
        ("+82 M DH d'investissements", TURQ_MID),
        ("Estimation · Commission · Suivi · Contrôle", GRIS),
        ("Compétences C18 + C19 mobilisées", TURQ_LIGHT),
    ]
    by4 = Inches(5.05)
    for bt, bc in bilan:
        txt(sl, f"▸ {bt}", Inches(8.0), by4, Inches(5.1), Inches(0.36), sz=13, color=bc)
        by4 += Inches(0.4)

    footer(sl, 11, TOTAL)

    # ═══════════════════════════════════════════════════════════════════════
    # 12. SYNTHÈSE COMPÉTENCES + radar / graphique
    # ═══════════════════════════════════════════════════════════════════════
    sl = add_slide(prs)
    bg(sl, BG)
    top_bar(sl, TURQ)

    txt(sl, "04  BILAN ET ANALYSE", Inches(0.3), Inches(0.1),
        Inches(10), Inches(0.4), sz=13, bold=True, color=TURQ_MID)
    txt(sl, "SYNTHÈSE DES COMPÉTENCES DÉVELOPPÉES",
        Inches(0.3), Inches(0.47), Inches(12.5), Inches(0.58), sz=26, bold=True, color=BLANC)

    # Tableau synthèse
    cols3 = ["Activité réalisée", "Compétence", "Situation", "Niveau"]
    rows3 = [
        ("Estimation confidentielle — 15,8 M DH HT",  "C18 — Métrés & estimation",   "Sit. 1", "Maîtrise"),
        ("Commission AO — analyse 3 offres",            "C19 — Analyse des offres",    "Sit. 2", "Maîtrise"),
        ("Tableau de bord financier — Kerrouchen",      "C19 — Suivi financier",       "Sit. 3", "Maîtrise"),
        ("Coordination 4 chantiers — 3 crises 48h",    "C18 — Communication chantier","Sit. 4", "Maîtrise"),
        ("Cubatures — Route Lehri-Kerrouchen",          "C18 — Contrôle quantités",    "Sit. 5", "Maîtrise"),
        ("5 marchés complémentaires",                   "C18+C19 — Cycle MPC complet", "Compl.", "Application"),
    ]
    cw3 = [Inches(4.1), Inches(3.5), Inches(1.4), Inches(1.8)]
    hx3 = Inches(0.2)
    hy3 = Inches(1.2)
    for j, (h3, cw) in enumerate(zip(cols3, cw3)):
        rect(sl, hx3, hy3, cw, Inches(0.38), RGBColor(0x00, 0x60, 0x65))
        txt(sl, h3, hx3 + Inches(0.06), hy3 + Pt(3), cw - Inches(0.1), Inches(0.32),
            sz=11, bold=True, color=BLANC, align=PP_ALIGN.CENTER)
        hx3 += cw + Inches(0.05)

    ry3 = hy3 + Inches(0.42)
    for k, row in enumerate(rows3):
        rbg3 = BG_CARD if k % 2 == 0 else RGBColor(0x00, 0x3A, 0x3F)
        rx3 = Inches(0.2)
        for j, (cell, cw) in enumerate(zip(row, cw3)):
            rect(sl, rx3, ry3, cw, Inches(0.48), rbg3)
            txt(sl, cell, rx3 + Inches(0.06), ry3 + Pt(3), cw - Inches(0.1), Inches(0.42),
                sz=12 if j < 2 else 13, color=GOLD if j == 3 else BLANC, bold=(j == 3))
            rx3 += cw + Inches(0.05)
        ry3 += Inches(0.52)

    # Graphique niveau de compétences
    add_bar_chart(sl,
        title="Niveau de maîtrise des compétences (auto-évaluation /5)",
        categories=["Métrés &\nestimation", "Analyse\noffres", "Suivi\nfinancier",
                    "Coordination\nchantier", "Contrôle\nquantités", "Procédures\nMarchés pub."],
        series_data=[("Niveau /5", (4.5, 4.5, 4.8, 4.2, 4.3, 4.6))],
        x=Inches(10.9), y=Inches(1.15), w=Inches(2.25), h=Inches(3.8),
        bar_colors=["00C8D0"], gap=60)

    # Analyse comparative
    rect(sl, Inches(0.2), ry3 + Inches(0.1), W - Inches(3.4), Inches(1.35), BG_DARK)
    rect(sl, Inches(0.2), ry3 + Inches(0.1), W - Inches(3.4), Inches(0.05), TURQ)
    txt(sl, "ANALYSE COMPARATIVE — MAROC vs FRANCE", Inches(0.35), ry3 + Inches(0.18),
        Inches(7), Inches(0.35), sz=12, bold=True, color=TURQ_MID)
    compare = ("Maroc (Décret 2013) : AO ouvert/restreint/négocié · CPS + RC + BPDE + Estimation confidentielle\n"
               "France (Code commande publique) : CCAG Travaux · DPGF · BIM Level 2 progressif\n"
               "Commun : transparence · égalité de traitement · offre économiquement la plus avantageuse")
    txt(sl, compare, Inches(0.35), ry3 + Inches(0.57), W - Inches(3.6), Inches(0.75), sz=11.5, color=BLANC)

    footer(sl, 12, TOTAL)

    # ═══════════════════════════════════════════════════════════════════════
    # 13. PROTOCOLE BIM + PROJET PRO
    # ═══════════════════════════════════════════════════════════════════════
    sl = add_slide(prs)
    bg(sl, BG)
    top_bar(sl, TURQ)

    txt(sl, "PROTOCOLE BIM & PROJET PROFESSIONNEL",
        Inches(0.3), Inches(0.1), Inches(12), Inches(0.5), sz=26, bold=True, color=BLANC)

    # Colonne BIM
    rect(sl, Inches(0.2), Inches(0.78), Inches(6.3), Inches(6.35), BG_CARD)
    rect(sl, Inches(0.2), Inches(0.78), Inches(6.3), Inches(0.05), TURQ)
    txt(sl, "PROTOCOLE BIM — BIMCO", Inches(0.35), Inches(0.86),
        Inches(6.0), Inches(0.38), sz=14, bold=True, color=GOLD)
    bim_items = [
        ("Convention BIM", "LOD 300 (chiffrage) / LOI 3 (données matériaux)"),
        ("Format d'échange", "IFC 2x3 — export Revit → coordination Navisworks"),
        ("Extraction métrés", "Schedules Revit paramétrés → Excel → DPGF automatisé"),
        ("Gain vs méthode trad.", "–40% temps métrés | +15% fiabilité quantités"),
        ("Outils", "Revit · Navisworks · Dynamo · AutoCAD · Python · API Revit"),
    ]
    by5 = Inches(1.38)
    for lbl5, val5 in bim_items:
        txt(sl, lbl5, Inches(0.35), by5, Inches(6.0), Inches(0.28),
            sz=11, bold=True, color=TURQ_MID)
        txt(sl, val5, Inches(0.35), by5 + Inches(0.28), Inches(6.0), Inches(0.36),
            sz=12, color=BLANC)
        by5 += Inches(0.7)

    # Graphique gain BIM
    add_bar_chart(sl,
        title="Gain BIM vs méthode traditionnelle",
        categories=["Métrés", "Fiabilité\nquantités", "Temps\nestimation"],
        series_data=[
            ("Traditionnel (%)", (100, 100, 100)),
            ("BIM (%)", (60, 115, 65)),
        ],
        x=Inches(0.25), y=Inches(5.05), w=Inches(6.2), h=Inches(2.1),
        bar_colors=["005A60", "00C8D0"], gap=30)

    im41 = os.path.join(MEDIA, "image41.png")
    if os.path.exists(im41):
        img(sl, im41, Inches(0.2), Inches(4.8), Inches(3.0), Inches(0.2))

    # Colonne Projet pro
    rect(sl, Inches(6.7), Inches(0.78), Inches(6.4), Inches(6.35), BG_DARK)
    rect(sl, Inches(6.7), Inches(0.78), Inches(6.4), Inches(0.05), GOLD)
    txt(sl, "MON PROJET PROFESSIONNEL", Inches(6.85), Inches(0.86),
        Inches(6.1), Inches(0.38), sz=14, bold=True, color=GOLD)
    proj = [
        ("2026 — Court terme",
         "• Obtenir le BTS MEC — validation officielle du parcours terrain\n"
         "• BIMCO : métrés BIM, DPGF, études de prix pour TPE du BTP\n"
         "• Développer API Revit pour extraction de métrés automatisée"),
        ("2027-2028 — Moyen terme",
         "• Certification BIM Manager\n"
         "• Partenariats avec MOE et promoteurs immobiliers\n"
         "• Déployer outils BIMCO aux TPE du secteur BTP"),
        ("Long terme",
         "• Référence régionale BIM + économie de la construction\n"
         "• Formation professionnelle sur les outils BIM numériques"),
    ]
    py3 = Inches(1.38)
    for lbl6, val6 in proj:
        rect(sl, Inches(6.85), py3, Inches(6.1), Inches(0.3), RGBColor(0x00, 0x50, 0x55))
        txt(sl, lbl6, Inches(6.92), py3 + Pt(2), Inches(5.85), Inches(0.27),
            sz=11, bold=True, color=TURQ_LIGHT)
        txt(sl, val6, Inches(6.92), py3 + Inches(0.33), Inches(5.95), Inches(1.0),
            sz=12.5, color=BLANC)
        py3 += Inches(1.45)

    footer(sl, 13, TOTAL)

    # ═══════════════════════════════════════════════════════════════════════
    # 14. BILAN RÉFLEXIF
    # ═══════════════════════════════════════════════════════════════════════
    sl = add_slide(prs)
    bg(sl, BG)
    top_bar(sl, TURQ)

    txt(sl, "BILAN RÉFLEXIF", Inches(0.3), Inches(0.1),
        Inches(10), Inches(0.42), sz=13, bold=True, color=TURQ_MID)
    txt(sl, "Ce que je retiens — Ce que j'aurais fait différemment",
        Inches(0.3), Inches(0.48), Inches(12.5), Inches(0.58), sz=24, bold=True, color=BLANC)

    reflexions = [
        ("A — Ce que j'ai appris", TURQ_MID,
         "La semaine 23 : 3 crises simultanées sur 4 chantiers à 80 km. "
         "L'anticipation documentée vaut mieux que la réaction rapide. "
         "Un OS émis 24h après l'événement, avec photos datées, "
         "protège tout le monde."),
        ("D — Ce que j'aurais fait autrement", ROUGE,
         "Mettre en place le tableau de bord dès le premier mois, pas à mi-chantier. "
         "Le réflexe de traçabilité systématique (même pour les décisions mineures) "
         "est venu tardivement. Aujourd'hui c'est un automatisme."),
        ("C — Ce que je maîtrise maintenant", VERT,
         "L'articulation entre le technique et le financier : une dérive de "
         "terrassement, c'est aussi un seuil d'avenant, une compensation inter-postes "
         "et des délais administratifs. Penser les deux simultanément."),
        ("E — Ce que ça m'a apporté", GOLD,
         "BIMCO : automatiser les métrés par extraction BIM pour libérer du temps "
         "sur ce qui demande du jugement — analyse, coordination, décision. "
         "La formation BIM et le BTS MEC se renforcent mutuellement."),
    ]
    rw2 = (W - Inches(0.5)) / 2 - Inches(0.1)
    rh2 = Inches(2.5)
    positions2 = [
        (Inches(0.15), Inches(1.18)), (W/2 + Inches(0.15), Inches(1.18)),
        (Inches(0.15), Inches(3.78)), (W/2 + Inches(0.15), Inches(3.78)),
    ]
    bgs2 = [RGBColor(0x00, 0x3A, 0x50), RGBColor(0x50, 0x12, 0x12),
             RGBColor(0x12, 0x40, 0x20), RGBColor(0x30, 0x28, 0x00)]
    for (rx4, ry4), (lbl7, lc7, t7), rbc2 in zip(positions2, reflexions, bgs2):
        rect(sl, rx4, ry4, rw2, rh2, rbc2)
        rect(sl, rx4, ry4, rw2, Inches(0.05), lc7)
        txt(sl, lbl7, rx4 + Inches(0.1), ry4 + Inches(0.1),
            rw2 - Inches(0.2), Inches(0.35), sz=12, bold=True, color=lc7)
        txt(sl, t7, rx4 + Inches(0.1), ry4 + Inches(0.48),
            rw2 - Inches(0.2), rh2 - Inches(0.58), sz=13, color=BLANC)

    footer(sl, 14, TOTAL)

    # ═══════════════════════════════════════════════════════════════════════
    # 15. CONCLUSION
    # ═══════════════════════════════════════════════════════════════════════
    sl = add_slide(prs)
    bg(sl, BG)
    im1b = os.path.join(MEDIA, "image1.jpeg")
    if os.path.exists(im1b):
        img(sl, im1b, 0, 0, W, H)
        overlay(sl, alpha=78)
    top_bar(sl, TURQ)
    left_bar(sl, GOLD, Inches(0.1))

    txt(sl, "CONCLUSION", Inches(0.3), Inches(0.25),
        Inches(10), Inches(0.52), sz=38, bold=True, color=TURQ)
    txt(sl, ("Ce rapport est avant tout une démonstration — non pas de ce que je sais, "
             "mais de comment je résous un problème sous contrainte."),
        Inches(0.3), Inches(0.88), Inches(11), Inches(0.62), sz=17, color=GRIS)

    points2 = [
        ("1", "5 compétences construites sur le terrain",
         "Estimer · analyser les offres · suivre financièrement · coordonner en crise · contrôler les quantités"),
        ("2", "Une lecture double — Maroc et France",
         "Deux cadres réglementaires, les mêmes exigences : rigueur, traçabilité, transparence"),
        ("3", "Le BIM comme prolongement naturel",
         "BIMCO : appliquer le numérique à l'économie de la construction — pas un projet parallèle"),
        ("4", "Prêt pour la prochaine étape",
         "BTS MEC + Certification BIM Manager + BIMCO — un parcours cohérent et délibéré"),
    ]
    py4 = Inches(1.68)
    for n2, tit2, txt2 in points2:
        rect(sl, Inches(0.35), py4, Inches(0.55), Inches(0.55), TURQ)
        txt(sl, n2, Inches(0.35), py4 + Pt(2), Inches(0.55), Inches(0.5),
            sz=24, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)
        txt(sl, tit2, Inches(1.02), py4, Inches(12.0), Inches(0.38),
            sz=17, bold=True, color=BLANC)
        txt(sl, txt2, Inches(1.02), py4 + Inches(0.38), Inches(12.0), Inches(0.48),
            sz=13.5, color=GRIS)
        py4 += Inches(1.12)

    rect(sl, Inches(0.35), py4 + Inches(0.15), W - Inches(0.6), Inches(0.06), TURQ)
    txt(sl, "Merci pour votre attention — Je suis disponible pour vos questions.",
        Inches(0.35), py4 + Inches(0.3), W - Inches(0.6), Inches(0.45),
        sz=20, bold=True, color=TURQ_MID, align=PP_ALIGN.CENTER)

    footer(sl, 15, TOTAL)

    # ─── SAUVEGARDE ──────────────────────────────────────────────────────────
    out = ("d:/PREPA BTS MEC/08_U62_Rapport_Activites/Soutenance_PowerPoint/"
           "SOUTENANCE_16x9_GRAND_ECRAN.pptx")
    prs.save(out)
    print(f"Présentation générée : {out}")
    print(f"{len(prs.slides)} slides  |  16:9 ({W.inches:.2f}\" x {H.inches:.2f}\")")


if __name__ == "__main__":
    build()
