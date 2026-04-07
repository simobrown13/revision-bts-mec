"""
Génération d'un PowerPoint 30 pages pour le Rapport U62 - BAHAFID Mohamed
Style : Minimaliste Moderne Rouge (inspiré du template Canva)
BTS Management Économique de la Construction - Session 2026

POUR MODIFIER LE TEXTE : éditez le fichier contenu.py (pas ce fichier)
"""
import sys, os
sys.stdout.reconfigure(encoding="utf-8", errors="replace")

from contenu import (
    CANDIDAT, PAGE_01, PAGE_02, PAGE_03, PAGE_04, PAGE_05, PAGE_06, PAGE_07,
    PAGE_08, PAGE_09, PAGE_10, PAGE_11, PAGE_12, PAGE_13, PAGE_14, PAGE_15,
    PAGE_16, PAGE_17, PAGE_18, PAGE_19, PAGE_20, PAGE_21, PAGE_22, PAGE_23,
    PAGE_24, PAGE_25, PAGE_26, PAGE_27, PAGE_28, PAGE_29, PAGE_30,
)

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData
from pptx.oxml.ns import qn
from lxml import etree
import math

# =============================================================================
# CONFIGURATION
# =============================================================================
BASE_DIR = r"D:\PREPA BTS MEC\08_U62_Rapport_Activites"
PHOTOS_DIR = os.path.join(BASE_DIR, "Documents_Entreprise", "Photos_Chantier")
ANNEXES_DIR = os.path.join(BASE_DIR, "Annexes", "Images")
OUTPUT = os.path.join(BASE_DIR, "Rapport_Redaction", "RAPPORT_U62_BAHAFID_30pages_v8.pptx")

# Couleurs
ROUGE = RGBColor(0xC0, 0x39, 0x2B)
ROUGE_CLAIR = RGBColor(0xE7, 0x4C, 0x3C)
ROUGE_FONCE = RGBColor(0x8B, 0x00, 0x00)
NOIR = RGBColor(0x33, 0x33, 0x33)
BLANC = RGBColor(0xFF, 0xFF, 0xFF)
GRIS_CLAIR = RGBColor(0xF5, 0xF5, 0xF5)
GRIS = RGBColor(0x99, 0x99, 0x99)
GRIS_FONCE = RGBColor(0x55, 0x55, 0x55)
VERT = RGBColor(0x27, 0xAE, 0x60)
ORANGE = RGBColor(0xF3, 0x9C, 0x12)

# Dimensions A4 portrait en EMU
A4_W = Cm(21)
A4_H = Cm(29.7)

# Photos sélectionnées
PHOTO = {
    "cover": os.path.join(PHOTOS_DIR, "tom-shamberger--KAoVUNv9W0-unsplash.jpg"),
    "bahafid": os.path.join(ANNEXES_DIR, "photo_bahafid.jpg"),
    "maroc": os.path.join(PHOTOS_DIR, "20180330_153333.jpg"),
    "france": os.path.join(PHOTOS_DIR, "IMG_20250718_155345[1].jpg"),
    "bimco_tech": os.path.join(PHOTOS_DIR, "d-c-NFVIq89r_8Q-unsplash.jpg"),
    "conseil": os.path.join(PHOTOS_DIR, "20181206_125854.jpg"),
    "chantier1": os.path.join(PHOTOS_DIR, "20180523_102152.jpg"),
    "chantier2": os.path.join(PHOTOS_DIR, "20180605_122447.jpg"),
    "chantier3": os.path.join(PHOTOS_DIR, "20180604_133353.jpg"),
    "chantier4": os.path.join(PHOTOS_DIR, "20180511_155634.jpg"),
    "chantier5": os.path.join(PHOTOS_DIR, "20180503_164033.jpg"),
    "chantier6": os.path.join(PHOTOS_DIR, "20180403_152746.jpg"),
    "chantier7": os.path.join(PHOTOS_DIR, "20180405_115900.jpg"),
    "projet1": os.path.join(PHOTOS_DIR, "20180403_123714.jpg"),
    "projet2": os.path.join(PHOTOS_DIR, "20181206_125636.jpg"),
    "route": os.path.join(PHOTOS_DIR, "20181206_130504.jpg"),
    "route2": os.path.join(PHOTOS_DIR, "20181206_141917.jpg"),
    "bilan": os.path.join(PHOTOS_DIR, "pexels-efeburakbaydar-35846752.jpg"),
    "conclusion": os.path.join(PHOTOS_DIR, "d-c-zEjnjuA_KBY-unsplash.jpg"),
    "gros_oeuvre": os.path.join(PHOTOS_DIR, "20180228_144045.jpg"),
    "vrd": os.path.join(PHOTOS_DIR, "20180330_154009.jpg"),
    "terrassement": os.path.join(PHOTOS_DIR, "20181206_133336.jpg"),
    "logo": os.path.join(ANNEXES_DIR, "logo_bimco.png"),
    "banniere": os.path.join(ANNEXES_DIR, "banniere_bimco.png"),
    "cao": os.path.join(ANNEXES_DIR, "convocation_cao.jpg"),
    "rejet": os.path.join(ANNEXES_DIR, "notification_rejet.jpg"),
    "cps": os.path.join(ANNEXES_DIR, "cps_marche46_signe.jpg"),
    "chantier_urban": os.path.join(PHOTOS_DIR, "20180503_105806.jpg"),
    "montagne": os.path.join(PHOTOS_DIR, "20181113_154208.jpg"),
    "france2": os.path.join(PHOTOS_DIR, "IMG_20250910_163813.jpg"),
}

# =============================================================================
# HELPERS AVANCÉS (transparence, ombres, dégradés)
# =============================================================================
def img(name):
    """Retourne le chemin d'une photo si elle existe."""
    p = PHOTO.get(name, "")
    return p if os.path.exists(p) else None

def set_transparency(shape, opacity_pct):
    """Appliquer une transparence à un shape (0=invisible, 100=opaque)."""
    fill_elem = shape._element.spPr.find(qn('a:solidFill'))
    if fill_elem is None:
        return
    srgb = fill_elem.find(qn('a:srgbClr'))
    if srgb is None:
        return
    for a in srgb.findall(qn('a:alpha')):
        srgb.remove(a)
    alpha = srgb.makeelement(qn('a:alpha'), {'val': str(int(opacity_pct * 1000))})
    srgb.append(alpha)

def add_shadow(shape, blur=50800, dist=38100, direction=5400000, alpha_pct=35):
    """Ajouter une ombre portée à un shape."""
    spPr = shape._element.spPr
    effectLst = spPr.find(qn('a:effectLst'))
    if effectLst is None:
        effectLst = spPr.makeelement(qn('a:effectLst'), {})
        spPr.append(effectLst)
    outerShdw = effectLst.makeelement(qn('a:outerShdw'), {
        'blurRad': str(blur), 'dist': str(dist),
        'dir': str(direction), 'algn': 'bl', 'rotWithShape': '0'
    })
    srgb = outerShdw.makeelement(qn('a:srgbClr'), {'val': '000000'})
    alpha = srgb.makeelement(qn('a:alpha'), {'val': str(int(alpha_pct * 1000))})
    srgb.append(alpha)
    outerShdw.append(srgb)
    effectLst.append(outerShdw)

def add_gradient_fill(shape, color1, color2, angle_deg=270):
    """Remplir un shape avec un dégradé."""
    spPr = shape._element.spPr
    # Supprimer les fills existants
    for tag in ['a:solidFill', 'a:noFill', 'a:gradFill']:
        for e in spPr.findall(qn(tag)):
            spPr.remove(e)
    gradFill = spPr.makeelement(qn('a:gradFill'), {'rotWithShape': '1'})
    gsLst = gradFill.makeelement(qn('a:gsLst'), {})
    for pos, color in [('0', color1), ('100000', color2)]:
        gs = gsLst.makeelement(qn('a:gs'), {'pos': pos})
        srgb = gs.makeelement(qn('a:srgbClr'), {
            'val': f'{color.red:02X}{color.green:02X}{color.blue:02X}'
        })
        gs.append(srgb)
        gsLst.append(gs)
    gradFill.append(gsLst)
    lin = gradFill.makeelement(qn('a:lin'), {
        'ang': str(angle_deg * 60000), 'scaled': '1'
    })
    gradFill.append(lin)
    spPr.insert(0, gradFill)

def add_gradient_3stops(shape, c1, c2, c3, angle_deg=270):
    """Dégradé à 3 couleurs pour plus de profondeur."""
    spPr = shape._element.spPr
    for tag in ['a:solidFill', 'a:noFill', 'a:gradFill']:
        for e in spPr.findall(qn(tag)):
            spPr.remove(e)
    gradFill = spPr.makeelement(qn('a:gradFill'), {'rotWithShape': '1'})
    gsLst = gradFill.makeelement(qn('a:gsLst'), {})
    for pos, color in [('0', c1), ('50000', c2), ('100000', c3)]:
        gs = gsLst.makeelement(qn('a:gs'), {'pos': pos})
        srgb = gs.makeelement(qn('a:srgbClr'), {
            'val': f'{color.red:02X}{color.green:02X}{color.blue:02X}'
        })
        gs.append(srgb)
        gsLst.append(gs)
    gradFill.append(gsLst)
    lin = gradFill.makeelement(qn('a:lin'), {'ang': str(angle_deg * 60000), 'scaled': '1'})
    gradFill.append(lin)
    spPr.insert(0, gradFill)

def set_text_vertical_center(textbox):
    """Centrer verticalement le texte dans une textbox."""
    txBody = textbox._element.find(qn('p:txBody'))
    if txBody is None:
        return
    bodyPr = txBody.find(qn('a:bodyPr'))
    if bodyPr is not None:
        bodyPr.set('anchor', 'ctr')

def add_bg_rect(slide, color, left=0, top=0, w=None, h=None, opacity=100):
    """Ajouter un rectangle de fond coloré avec transparence optionnelle."""
    w = w or A4_W
    h = h or A4_H
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if opacity < 100:
        set_transparency(shape, opacity)
    return shape

def add_image_safe(slide, path, left, top, width=None, height=None):
    """Ajouter une image en toute sécurité."""
    if path and os.path.exists(path):
        try:
            if width and height:
                return slide.shapes.add_picture(path, left, top, width, height)
            elif width:
                return slide.shapes.add_picture(path, left, top, width=width)
            elif height:
                return slide.shapes.add_picture(path, left, top, height=height)
            else:
                return slide.shapes.add_picture(path, left, top)
        except Exception as e:
            print(f"  Erreur image {path}: {e}")
    return None

def add_text_box(slide, text, left, top, width, height, font_size=12,
                 color=NOIR, bold=False, align=PP_ALIGN.LEFT, font_name="Calibri"):
    """Ajouter une zone de texte."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_multi_text(slide, lines, left, top, width, height, font_size=11,
                   color=NOIR, spacing=1.2, bold_first=False):
    """Ajouter du texte multi-lignes."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(font_size * spacing * 0.3)
        if bold_first and i == 0:
            p.font.bold = True
    return txBox

def add_circle(slide, left, top, size, color):
    """Ajouter un cercle coloré."""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_card(slide, left, top, w, h, title, text, icon_text="",
             bg_color=BLANC, title_color=ROUGE, text_color=NOIR, shadow=True):
    """Ajouter une carte stylisée avec ombre portée."""
    # Fond
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    rect.fill.solid()
    rect.fill.fore_color.rgb = bg_color
    rect.line.fill.background()
    if shadow:
        add_shadow(rect)
    # Icone texte en haut
    if icon_text:
        add_text_box(slide, icon_text, left + Cm(0.3), top + Cm(0.2),
                     w - Cm(0.6), Cm(1.2), font_size=18, color=title_color,
                     bold=True, align=PP_ALIGN.CENTER)
    # Titre
    add_text_box(slide, title, left + Cm(0.3), top + Cm(1.3),
                 w - Cm(0.6), Cm(1), font_size=9, color=title_color,
                 bold=True, align=PP_ALIGN.CENTER)
    # Texte
    add_text_box(slide, text, left + Cm(0.3), top + Cm(2.1),
                 w - Cm(0.6), h - Cm(2.4), font_size=7, color=text_color,
                 align=PP_ALIGN.CENTER)

def add_bar_h(slide, left, top, max_w, h, pct, color, label="", value=""):
    """Ajouter une barre horizontale de progression."""
    # Fond gris
    bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, max_w, h)
    bg.fill.solid()
    bg.fill.fore_color.rgb = GRIS_CLAIR
    bg.line.fill.background()
    # Barre de progression
    bar_w = int(max_w * pct / 100)
    if bar_w > 0:
        bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, bar_w, h)
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
    # Label
    if label:
        add_text_box(slide, label, left - Cm(5), top, Cm(4.8), h,
                     font_size=8, color=NOIR, align=PP_ALIGN.RIGHT)
    if value:
        add_text_box(slide, value, left + max_w + Cm(0.2), top, Cm(2), h,
                     font_size=8, color=GRIS_FONCE)

def add_page_number(slide, num):
    """Ajouter le numéro de page avec barre rouge et nom candidat."""
    # Petite frise islamique en bas (remplace la simple barre rouge)
    add_islamic_border(slide, A4_H - Cm(1.5), Cm(0.35), ROUGE, opacity=25)
    # Nom à gauche
    add_text_box(slide, "BAHAFID Mohamed | U62 | BTS MEC 2026",
                 Cm(1), A4_H - Cm(1.2), Cm(12), Cm(0.8), font_size=7, color=GRIS)
    # Numéro à droite
    add_text_box(slide, str(num).zfill(2), A4_W - Cm(2), A4_H - Cm(1.2),
                 Cm(1.5), Cm(0.8), font_size=8, color=GRIS, align=PP_ALIGN.RIGHT)

def add_red_accent(slide, left=0, top=0, w=Cm(0.4), h=None):
    """Bande rouge décorative."""
    h = h or A4_H
    return add_bg_rect(slide, ROUGE, left, top, w, h)

def add_section_separator(slide, num, title, subtitle, photo_path):
    """Page de séparation de section avec photo pleine page et overlay semi-transparent."""
    # Photo de fond
    add_image_safe(slide, photo_path, 0, 0, A4_W, A4_H)
    # Overlay dégradé semi-transparent (transparent en haut, opaque en bas)
    overlay = add_bg_rect(slide, ROUGE_FONCE, 0, 0, A4_W, A4_H, opacity=70)
    # --- Motifs islamiques ---
    add_islamic_border(slide, Cm(3.5), Cm(0.8), BLANC, opacity=12)
    add_islamic_star(slide, Cm(3), Cm(4), Cm(2.5), BLANC, opacity=10)
    add_islamic_star(slide, int(A4_W) - Cm(3), Cm(4), Cm(2.5), BLANC, opacity=10)
    add_islamic_border(slide, A4_H - Cm(3), Cm(0.8), BLANC, opacity=12)
    # Grand numéro semi-transparent
    num_box = add_text_box(slide, num, Cm(2), Cm(5), Cm(17), Cm(10),
                           font_size=160, color=BLANC, bold=True, align=PP_ALIGN.CENTER)
    # Ligne décorative remplacée par frise islamique
    add_islamic_border(slide, Cm(15.8), Cm(0.5), ROUGE_CLAIR, opacity=30)
    # Titre
    add_text_box(slide, title, Cm(2), Cm(17), Cm(17), Cm(3),
                 font_size=28, color=BLANC, bold=True, align=PP_ALIGN.CENTER)
    # Sous-titre
    add_text_box(slide, subtitle, Cm(2), Cm(20.5), Cm(17), Cm(2),
                 font_size=14, color=BLANC, align=PP_ALIGN.CENTER)

def add_fullpage_photo_overlay(slide, photo_path, overlay_color=ROUGE_FONCE, opacity=65):
    """Photo pleine page avec overlay semi-transparent."""
    add_image_safe(slide, photo_path, 0, 0, A4_W, A4_H)
    add_bg_rect(slide, overlay_color, 0, 0, A4_W, A4_H, opacity=opacity)

# =============================================================================
# MOTIFS GÉOMÉTRIQUES ISLAMIQUES (vectoriels)
# =============================================================================

def add_islamic_star(slide, cx, cy, size, color=None, opacity=25):
    """Étoile islamique à 8 branches (Rub el Hizb) en vectoriel.
    cx, cy : centre absolu sur la slide (EMU)
    size   : diamètre global (EMU)
    """
    c = color or ROUGE
    r = size / 2
    inner_r = r * 0.38
    # 16 sommets en coordonnées locales (centre = (r, r))
    verts = []
    for i in range(16):
        angle = math.radians(i * 22.5 - 90)
        radius = r if i % 2 == 0 else inner_r
        verts.append((int(r + radius * math.cos(angle)),
                      int(r + radius * math.sin(angle))))
    try:
        builder = slide.shapes.build_freeform(verts[0][0], verts[0][1])
        for v in verts[1:]:
            builder.add_line_to(v[0], v[1])
        builder.add_line_to(verts[0][0], verts[0][1])
        shape = builder.convert_to_shape(int(cx - r), int(cy - r))
        shape.fill.solid()
        shape.fill.fore_color.rgb = c
        shape.line.fill.background()
        if opacity < 100:
            set_transparency(shape, opacity)
        return shape
    except Exception:
        # Fallback : losange simple
        d = slide.shapes.add_shape(MSO_SHAPE.DIAMOND,
                                   int(cx - r), int(cy - r), int(size), int(size))
        d.fill.solid()
        d.fill.fore_color.rgb = c
        d.line.fill.background()
        if opacity < 100:
            set_transparency(d, opacity)
        return d


def add_islamic_border(slide, y, band_h=None, color=None, opacity=18):
    """Frise horizontale de losanges entrelacés (motif islamique)."""
    c = color or ROUGE
    h = band_h or Cm(0.6)
    cell = int(h)
    x = 0
    while x < int(A4_W) + cell:
        d = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, x, int(y), cell, cell)
        d.fill.solid()
        d.fill.fore_color.rgb = c
        d.line.fill.background()
        set_transparency(d, opacity)
        x += cell


def add_islamic_corner(slide, corner, size=None, color=None, opacity=18):
    """Ornement géométrique en coin : étoile + losanges satellites.
    corner: 'tl', 'tr', 'bl', 'br'
    """
    sz = size or Cm(3)
    c = color or ROUGE_CLAIR
    half = int(sz) // 2
    corners_pos = {
        'tl': (0, 0),
        'tr': (int(A4_W), 0),
        'bl': (0, int(A4_H)),
        'br': (int(A4_W), int(A4_H)),
    }
    cx, cy = corners_pos[corner]
    add_islamic_star(slide, cx, cy, int(sz), c, opacity)
    # Losanges satellites
    small = int(sz * 0.35)
    for dx, dy in [(half, 0), (0, half), (-half, 0), (0, -half)]:
        nx, ny = cx + dx, cy + dy
        d = slide.shapes.add_shape(MSO_SHAPE.DIAMOND,
                                   nx - small // 2, ny - small // 2, small, small)
        d.fill.solid()
        d.fill.fore_color.rgb = c
        d.line.fill.background()
        set_transparency(d, max(opacity - 5, 5))


def add_islamic_star_outline(slide, cx, cy, size, color=None, opacity=30, line_w=0.75):
    """Étoile en filigrane (contour uniquement, pas de remplissage)."""
    c = color or ROUGE_CLAIR
    half = int(size) // 2
    for rot in [0, 45.0]:
        sq = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    cx - half, cy - half, int(size), int(size))
        sq.fill.background()
        sq.line.color.rgb = c
        sq.line.width = Pt(line_w)
        sq.rotation = rot
        # Transparence de la ligne via XML
        ln = sq._element.spPr.find(qn('a:ln'))
        if ln is not None:
            srgb = ln.find(qn('a:solidFill'))
            if srgb is None:
                srgb = ln.find('./' + qn('a:solidFill'))
            sf = ln.find(qn('a:solidFill'))
            if sf is not None:
                clr = sf.find(qn('a:srgbClr'))
                if clr is not None:
                    alpha = clr.makeelement(qn('a:alpha'),
                                            {'val': str(int(opacity * 1000))})
                    clr.append(alpha)


def add_islamic_lattice(slide, left, top, width, height, cell_size=None,
                        color=None, opacity=10):
    """Treillis losangé subtil couvrant une zone rectangulaire."""
    c = color or ROUGE
    cs = cell_size or Cm(1.2)
    cell = int(cs)
    half = cell // 2
    y = int(top)
    row = 0
    while y < int(top + height):
        x = int(left) + (half if row % 2 else 0)
        while x < int(left + width):
            d = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, x, y, cell, cell)
            d.fill.solid()
            d.fill.fore_color.rgb = c
            d.line.fill.background()
            set_transparency(d, opacity)
            x += cell
        y += half
        row += 1


# =============================================================================
# PAGES
# =============================================================================

def page_01_couverture(prs):
    """Page 1 - Couverture avec overlay semi-transparent et dégradé."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    # Photo de fond pleine page
    add_image_safe(slide, img("cover"), 0, 0, A4_W, A4_H)
    # Overlay dégradé semi-transparent (visible la photo en dessous)
    overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, A4_W, A4_H)
    overlay.line.fill.background()
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = RGBColor(0x2A, 0x00, 0x00)
    set_transparency(overlay, 65)
    # --- Motifs islamiques décoratifs ---
    # Étoile coin haut-droit (semi-transparent, derrière le texte)
    add_islamic_star(slide, int(A4_W) - Cm(2), Cm(2), Cm(4), ROUGE_CLAIR, opacity=15)
    add_islamic_star(slide, int(A4_W) - Cm(0.5), Cm(5), Cm(1.8), ROUGE_CLAIR, opacity=12)
    # Étoile coin bas-gauche
    add_islamic_star(slide, Cm(4), int(A4_H) - Cm(3), Cm(3), ROUGE_CLAIR, opacity=12)
    # Frise de losanges en bas de page
    add_islamic_border(slide, A4_H - Cm(1.5), Cm(0.7), ROUGE_CLAIR, opacity=15)
    # Bande verticale gauche avec dégradé
    bande = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Cm(1.5), A4_H)
    bande.line.fill.background()
    bande.fill.solid()
    bande.fill.fore_color.rgb = ROUGE
    # U62 vertical dans la bande
    add_text_box(slide, "U62", Cm(0.1), Cm(12), Cm(1.3), Cm(4),
                 font_size=18, color=BLANC, bold=True, align=PP_ALIGN.CENTER)
    # Ligne décorative horizontale fine
    deco = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(3), Cm(2.5), Cm(5), Pt(3))
    deco.fill.solid()
    deco.fill.fore_color.rgb = ROUGE_CLAIR
    deco.line.fill.background()
    # Titre principal - plus grand, plus impactant
    add_text_box(slide, PAGE_01["titre_1"], Cm(3), Cm(3.5), Cm(16), Cm(3),
                 font_size=48, color=BLANC, bold=True, align=PP_ALIGN.LEFT)
    add_text_box(slide, PAGE_01["titre_2"], Cm(3), Cm(6.5), Cm(16), Cm(3),
                 font_size=48, color=BLANC, bold=True, align=PP_ALIGN.LEFT)
    add_text_box(slide, PAGE_01["titre_3"], Cm(3), Cm(9.5), Cm(16), Cm(3),
                 font_size=36, color=ROUGE_CLAIR, bold=True, align=PP_ALIGN.LEFT)
    # BTS MEC avec fond semi-transparent
    mec_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(3), Cm(13.5), Cm(13), Cm(2.5))
    mec_bg.fill.solid()
    mec_bg.fill.fore_color.rgb = ROUGE
    mec_bg.line.fill.background()
    set_transparency(mec_bg, 80)
    add_text_box(slide, PAGE_01["bts"], Cm(3.5), Cm(13.8), Cm(12), Cm(2),
                 font_size=16, color=BLANC, bold=True, align=PP_ALIGN.LEFT)
    add_text_box(slide, PAGE_01["session"], Cm(3), Cm(17), Cm(10), Cm(1.5),
                 font_size=20, color=ROUGE_CLAIR, bold=True)
    # Ligne séparatrice épaisse
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(3), Cm(19), Cm(8), Pt(4))
    line.fill.solid()
    line.fill.fore_color.rgb = ROUGE_CLAIR
    line.line.fill.background()
    # Nom candidat - plus visible
    add_text_box(slide, CANDIDAT["nom"], Cm(3), Cm(20), Cm(15), Cm(2),
                 font_size=28, color=BLANC, bold=True)
    add_text_box(slide, PAGE_01["candidat_info"], Cm(3), Cm(22.5), Cm(15), Cm(1),
                 font_size=12, color=RGBColor(0xDD, 0xDD, 0xDD))
    # Logo BIMCO avec fond
    logo_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(14.5), Cm(24.5), Cm(5.5), Cm(4))
    logo_bg.fill.solid()
    logo_bg.fill.fore_color.rgb = BLANC
    logo_bg.line.fill.background()
    set_transparency(logo_bg, 85)
    add_shadow(logo_bg, alpha_pct=20)
    add_image_safe(slide, img("logo"), Cm(15), Cm(25), Cm(4.5))
    print("  Page 1 : Couverture OK")

def page_02_fiche_candidat(prs):
    """Page 2 - Fiche d'identité du candidat."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_rect(slide, BLANC)
    # Bandes décoratives remplacées par frises islamiques
    add_islamic_border(slide, Cm(-0.1), Cm(0.5), ROUGE, opacity=35)
    add_islamic_border(slide, A4_H - Cm(0.4), Cm(0.5), ROUGE, opacity=35)
    # Titre
    add_text_box(slide, PAGE_02["titre"], Cm(1), Cm(1), Cm(19), Cm(1.5),
                 font_size=22, color=ROUGE, bold=True, align=PP_ALIGN.CENTER)
    # Photo portrait dans un cercle (simulé par un cadre)
    circle = add_circle(slide, Cm(2), Cm(4), Cm(5), ROUGE)
    add_image_safe(slide, img("bahafid"), Cm(2.2), Cm(4.2), Cm(4.6), Cm(4.6))
    # Informations à droite
    infos = PAGE_02["champs"]
    y = Cm(3.5)
    for label, value in infos:
        add_text_box(slide, label, Cm(9), y, Cm(5), Cm(0.6),
                     font_size=7, color=GRIS, bold=True)
        add_text_box(slide, value, Cm(9), y + Cm(0.4), Cm(10), Cm(0.8),
                     font_size=10, color=NOIR, bold=False)
        y += Cm(1.7)
    # Bannière BIMCO en bas
    add_image_safe(slide, img("banniere"), Cm(3), Cm(24), Cm(15))
    # Site web
    add_text_box(slide, PAGE_02["pied"],
                 Cm(1), Cm(27), Cm(19), Cm(1), font_size=9, color=ROUGE, align=PP_ALIGN.CENTER)
    add_page_number(slide, 2)
    print("  Page 2 : Fiche candidat OK")

def page_03_sommaire(prs):
    """Page 3 - Sommaire."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_rect(slide, BLANC)
    add_red_accent(slide, 0, 0, Cm(1), A4_H)
    add_text_box(slide, "SOMMAIRE", Cm(2), Cm(1.5), Cm(17), Cm(2),
                 font_size=30, color=ROUGE, bold=True)
    # Frise islamique sous le titre
    add_islamic_border(slide, Cm(3.5), Cm(0.4), ROUGE, opacity=12)
    # Étoile décorative en bas à droite (subtile)
    add_islamic_star(slide, int(A4_W) - Cm(2.5), int(A4_H) - Cm(3.5), Cm(3), ROUGE_CLAIR, opacity=10)
    sections = [
        ("01", "Introduction", "Parcours, contexte, objectifs du rapport", "p. 4"),
        ("02", "Cadre professionnel", "Mon parcours | Conseil Régional | BIMCO | Outils", "p. 5-12"),
        ("03", "Projet 1 : Mise à niveau 4 communes", "53,5 M DH TTC – 8 corps d'état – Province de Khénifra", "p. 13-21"),
        ("04", "Projet 2 : Route Lehri-Kerrouchen", "29 M DH TTC – 25 km en zone montagneuse", "p. 22-25"),
        ("05", "Activités complémentaires", "5 autres marchés au Maroc + Expérience terrain France", "p. 26"),
        ("06", "Compétences et analyse", "Référentiel BTS MEC + Comparaison Maroc / France", "p. 27-28"),
        ("07", "Projet professionnel", "Court, moyen et long terme – Vision BIMCO", "p. 29"),
        ("08", "Conclusion", "Bilan et perspectives", "p. 30"),
    ]
    y = Cm(4.5)
    for num, title, subtitle, page in sections:
        # Cercle numéro
        c = add_circle(slide, Cm(2.5), y, Cm(1.2), ROUGE)
        add_text_box(slide, num, Cm(2.5), y + Cm(0.15), Cm(1.2), Cm(1),
                     font_size=11, color=BLANC, bold=True, align=PP_ALIGN.CENTER)
        # Ligne
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(4.2), y + Cm(0.55), Cm(12), Pt(1))
        line.fill.solid()
        line.fill.fore_color.rgb = GRIS_CLAIR
        line.line.fill.background()
        # Titre section
        add_text_box(slide, title, Cm(4.2), y, Cm(12), Cm(1),
                     font_size=13, color=NOIR, bold=True)
        # Sous-titre
        add_text_box(slide, subtitle, Cm(4.2), y + Cm(1), Cm(12), Cm(0.8),
                     font_size=8, color=GRIS_FONCE)
        # Page
        add_text_box(slide, page, Cm(16.5), y + Cm(0.1), Cm(3), Cm(1),
                     font_size=10, color=GRIS)
        y += Cm(3)
    add_page_number(slide, 3)
    print("  Page 3 : Sommaire OK")

def page_04_introduction(prs):
    """Page 4 - Introduction."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_rect(slide, BLANC)
    # Étoile islamique décorative en haut à droite
    add_islamic_star(slide, int(A4_W) - Cm(2), Cm(2), Cm(2.5), ROUGE_CLAIR, opacity=8)
    add_text_box(slide, "INTRODUCTION", Cm(1.5), Cm(1), Cm(18), Cm(1.5),
                 font_size=28, color=ROUGE, bold=True)
    # Points clés avec descriptions enrichies
    points = [
        ("8 ans d'expérience dans le BTP",
         "Parcours complet depuis la formation au Maroc jusqu'à la création de BIMCO en France"),
        ("Maîtrise d'ouvrage publique au Maroc (4,5 ans)",
         "7 marchés publics suivis au Conseil Régional de Béni Mellal-Khénifra, +100 M DH d'investissements"),
        ("Chef de chantier gros œuvre en France (5 ans)",
         "Coffrage, ferraillage, bétonnage chez Ergalis (Feurs) et Minssieux et Fils (Mornant)"),
        ("Formation BIM Modeleur - AFPA Colmar (8 mois)",
         "Revit, Navisworks, Dynamo, Python, C#/API Revit, formats IFC"),
        ("Création de BIMCO - Indépendant (janv. 2026)",
         "Projeteur BIM et Économiste de la construction - SIREN 999580053"),
    ]
    y = Cm(3.5)
    for title, desc in points:
        add_circle(slide, Cm(2), y + Cm(0.15), Cm(0.4), ROUGE)
        add_text_box(slide, title, Cm(3), y, Cm(16), Cm(0.8),
                     font_size=12, color=NOIR, bold=True)
        add_text_box(slide, desc, Cm(3), y + Cm(0.8), Cm(16), Cm(0.8),
                     font_size=10, color=GRIS_FONCE)
        y += Cm(2)
    # Paragraphe sur les deux projets
    add_text_box(slide, PAGE_04["projets_titre"],
                 Cm(2), Cm(14), Cm(17), Cm(0.8), font_size=11, color=NOIR, bold=True)
    add_multi_text(slide, PAGE_04["projets"],
                   Cm(2.5), Cm(15), Cm(16), Cm(2), font_size=10, color=NOIR)
    # Citation en encadré avec ombre et barre latérale rouge
    quote_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(2), Cm(17.5), Cm(17), Cm(2.5))
    quote_bg.fill.solid()
    quote_bg.fill.fore_color.rgb = GRIS_CLAIR
    quote_bg.line.fill.background()
    add_shadow(quote_bg, alpha_pct=15)
    quote_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(2), Cm(17.5), Cm(0.3), Cm(2.5))
    quote_bar.fill.solid()
    quote_bar.fill.fore_color.rgb = ROUGE
    quote_bar.line.fill.background()
    add_text_box(slide, PAGE_04["citation"],
                 Cm(3), Cm(17.7), Cm(15), Cm(2.2), font_size=12, color=ROUGE,
                 bold=True, align=PP_ALIGN.CENTER)
    # Infographie parcours en bas
    etapes = ["Formation\nMaroc", "MOA\nMaroc", "Chantier\nFrance", "BIM\nColmar", "BIMCO"]
    x = Cm(1.5)
    for i, etape in enumerate(etapes):
        c = add_circle(slide, x, Cm(20.5), Cm(2.5), ROUGE if i < 4 else ROUGE_FONCE)
        add_text_box(slide, etape, x, Cm(20.8), Cm(2.5), Cm(2),
                     font_size=7, color=BLANC, bold=True, align=PP_ALIGN.CENTER)
        dates = ["2014-17", "2017-22", "2022-24", "2023-24", "2026"]
        add_text_box(slide, dates[i], x, Cm(23.2), Cm(2.5), Cm(0.8),
                     font_size=7, color=GRIS_FONCE, align=PP_ALIGN.CENTER)
        if i < 4:
            # Flèche entre les cercles
            arrow = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Cm(2.5), Cm(21.5), Cm(1.3), Pt(3))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = ROUGE_CLAIR
            arrow.line.fill.background()
        x += Cm(3.8)
    add_page_number(slide, 4)
    print("  Page 4 : Introduction OK")

def page_05_section1(prs):
    """Page 5 - Séparateur Partie 1."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_separator(slide, "01", "CADRE PROFESSIONNEL",
                          "Mon parcours en 5 phases | Le Conseil Régional de Béni Mellal-Khénifra\nBIMCO – Mon activité indépendante | Outils et méthodes",
                          img("conseil"))
    print("  Page 5 : Séparateur Partie 1 OK")

def page_06_parcours(prs):
    """Page 6 - Timeline du parcours."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_rect(slide, BLANC)
    add_text_box(slide, "MON PARCOURS EN 5 PHASES", Cm(1), Cm(0.8), Cm(19), Cm(1),
                 font_size=22, color=ROUGE, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, "Plus de 10 années d'expérience, de la formation initiale au Maroc à la création "
                 "de BIMCO en France, en passant par la maîtrise d'ouvrage publique, le chantier et le BIM.",
                 Cm(1), Cm(2), Cm(19), Cm(1), font_size=8, color=GRIS_FONCE, align=PP_ALIGN.CENTER)
    # Timeline centrale
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(10.2), Cm(3.5), Cm(0.3), Cm(22))
    line.fill.solid()
    line.fill.fore_color.rgb = ROUGE
    line.line.fill.background()
    phases = [
        ("2014-2017", "Formation BTP au Maroc", "Technicien Chef de Chantier BTP\nDiplôme TS Gros Œuvre à l'ISBTP\nPlanification, organisation, métrés", True),
        ("2017-2022", "Chargé d'affaires – MOA", "Conseil Régional Béni Mellal-Khénifra\n7 marchés publics, +100 M DH\nEstimations, DCE, analyse offres, suivi", False),
        ("2022-2024", "Chef de chantier France", "Ergalis (Feurs, 42) puis Minssieux (69)\nCoffrage, banches, ferraillage, béton\nGestion d'équipe, contrôle qualité", True),
        ("2023-2024", "Formation BIM - AFPA", "8 mois à Colmar (68)\nRevit, Navisworks, Dynamo, C#/API\nIFC, extraction quantités, maquettes", False),
        ("Depuis 2026", "Création BIMCO", "Micro-entreprise, APE 7112B\nProjeteur BIM + Économiste construction\nApp « Gestion Chantiers » développée", True),
    ]
    y = Cm(3.8)
    for i, (date, title, desc, is_left) in enumerate(phases):
        # Point sur la timeline
        add_circle(slide, Cm(9.8), y + Cm(0.5), Cm(1.1), ROUGE_FONCE)
        add_text_box(slide, str(i + 1), Cm(9.9), y + Cm(0.6), Cm(0.9), Cm(0.9),
                     font_size=10, color=BLANC, bold=True, align=PP_ALIGN.CENTER)
        if is_left:
            # Carte à gauche avec ombre
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(1), y, Cm(8), Cm(3.5))
            card.fill.solid()
            card.fill.fore_color.rgb = BLANC
            card.line.fill.background()
            add_shadow(card, alpha_pct=15)
            # Barre rouge latérale
            bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(1), y, Cm(0.3), Cm(3.5))
            bar.fill.solid()
            bar.fill.fore_color.rgb = ROUGE
            bar.line.fill.background()
            add_text_box(slide, date, Cm(1.6), y + Cm(0.2), Cm(7.2), Cm(0.7),
                         font_size=9, color=ROUGE, bold=True)
            add_text_box(slide, title, Cm(1.6), y + Cm(0.9), Cm(7.2), Cm(0.7),
                         font_size=11, color=NOIR, bold=True)
            add_text_box(slide, desc, Cm(1.6), y + Cm(1.7), Cm(7.2), Cm(1.5),
                         font_size=9, color=GRIS_FONCE)
        else:
            # Carte à droite avec ombre
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(11.8), y, Cm(8), Cm(3.5))
            card.fill.solid()
            card.fill.fore_color.rgb = BLANC
            card.line.fill.background()
            add_shadow(card, alpha_pct=15)
            # Barre rouge latérale droite
            bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(19.5), y, Cm(0.3), Cm(3.5))
            bar.fill.solid()
            bar.fill.fore_color.rgb = ROUGE
            bar.line.fill.background()
            add_text_box(slide, date, Cm(12.1), y + Cm(0.2), Cm(7.2), Cm(0.7),
                         font_size=9, color=ROUGE, bold=True)
            add_text_box(slide, title, Cm(12.1), y + Cm(0.9), Cm(7.2), Cm(0.7),
                         font_size=11, color=NOIR, bold=True)
            add_text_box(slide, desc, Cm(12.1), y + Cm(1.7), Cm(7.2), Cm(1.5),
                         font_size=9, color=GRIS_FONCE)
        y += Cm(4.3)
    add_page_number(slide, 6)
    print("  Page 6 : Parcours timeline OK")

def page_07_chiffres_cles(prs):
    """Page 7 - Mon parcours en chiffres."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_rect(slide, GRIS_CLAIR)
    # Treillis islamique très subtil en fond
    add_islamic_lattice(slide, 0, Cm(3), A4_W, A4_H - Cm(5), Cm(1.5), ROUGE, opacity=5)
    add_text_box(slide, "MON PARCOURS EN CHIFFRES", Cm(1), Cm(1), Cm(19), Cm(1.5),
                 font_size=22, color=ROUGE, bold=True, align=PP_ALIGN.CENTER)
    chiffres = [
        ("8", "ANS", "d'expérience BTP", "Depuis la formation initiale\nau Maroc en 2014"),
        ("2", "PAYS", "Maroc + France", "MOA publique à Khénifra\npuis chantier en Rhône-Alpes"),
        ("7", "MARCHÉS", "publics suivis", "Routes, pistes, VRD, AEP,\naménagement urbain"),
        ("+100M", "DH", "d'investissements", "Volume global des marchés\ngérés (~9 M€)"),
        ("5", "ANS", "sur chantier", "Chef d'équipe et chef de\nchantier gros œuvre"),
        ("8", "MOIS", "de formation BIM", "AFPA Colmar : Revit,\nNavisworks, Dynamo, C#"),
        ("10+", "LOGICIELS", "maîtrisés", "BIM, DAO, bureautique,\ndéveloppement web"),
        ("1", "ENTREPRISE", "créée (BIMCO)", "Micro-entreprise depuis\njanvier 2026, APE 7112B"),
    ]
    positions = [(0, 0), (1, 0), (2, 0), (3, 0), (0, 1), (1, 1), (2, 1), (3, 1)]
    for idx, (num, unit, desc, detail) in enumerate(chiffres):
        col, row = positions[idx]
        x = Cm(1.2) + col * Cm(4.8)
        y = Cm(4) + row * Cm(12)
        # Carte avec ombre
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Cm(4.3), Cm(10))
        card.fill.solid()
        card.fill.fore_color.rgb = BLANC
        card.line.fill.background()
        add_shadow(card, alpha_pct=20)
        # Accent rouge en haut de la carte
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Cm(0.5), y + Cm(0.3), Cm(3.3), Pt(4))
        accent.fill.solid()
        accent.fill.fore_color.rgb = ROUGE
        accent.line.fill.background()
        # Grand chiffre
        add_text_box(slide, num, x, y + Cm(1), Cm(4.3), Cm(3),
                     font_size=42, color=ROUGE, bold=True, align=PP_ALIGN.CENTER)
        # Unité
        add_text_box(slide, unit, x, y + Cm(4), Cm(4.3), Cm(1.2),
                     font_size=14, color=NOIR, bold=True, align=PP_ALIGN.CENTER)
        # Description
        add_text_box(slide, desc, x, y + Cm(5.2), Cm(4.3), Cm(1.2),
                     font_size=9, color=GRIS_FONCE, align=PP_ALIGN.CENTER)
        # Ligne séparatrice
        sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Cm(1), y + Cm(6.5), Cm(2.3), Pt(1))
        sep.fill.solid()
        sep.fill.fore_color.rgb = GRIS_CLAIR
        sep.line.fill.background()
        # Détail supplémentaire
        add_text_box(slide, detail, x, y + Cm(7), Cm(4.3), Cm(2.5),
                     font_size=8, color=GRIS, align=PP_ALIGN.CENTER)
    add_page_number(slide, 7)
    print("  Page 7 : Chiffres clés OK")

def page_08_conseil_regional(prs):
    """Page 8 - Le Conseil Régional."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_rect(slide, BLANC)
    # Étoile islamique discrète (institution marocaine)
    add_islamic_star(slide, int(A4_W) - Cm(2.5), Cm(2), Cm(3), ROUGE_CLAIR, opacity=8)
    add_islamic_star(slide, Cm(2), Cm(2), Cm(2), ROUGE_CLAIR, opacity=6)
    add_text_box(slide, "STRUCTURE D'ACCUEIL", Cm(1), Cm(0.8), Cm(19), Cm(1),
                 font_size=14, color=GRIS, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, "Conseil Régional de\nBéni Mellal-Khénifra", Cm(1), Cm(2), Cm(19), Cm(3),
                 font_size=24, color=ROUGE, bold=True, align=PP_ALIGN.CENTER)
    # Photo à gauche
    add_image_safe(slide, img("conseil"), Cm(0.5), Cm(6), Cm(9.5), Cm(12))
    # Infos à droite
    facts = [
        ("5 Provinces", "Béni Mellal, Azilal, Fquih Ben Salah,\nKhénifra, Khouribga"),
        ("28 374 km²", "Superficie de la région"),
        ("2,5 millions", "Habitants"),
        ("Compétences", "Routes régionales et communales,\naménagement urbain et rural,\nadduction d'eau potable,\néquipements publics"),
        ("Budget invest.", "Centaines de millions de DH/an\nconsacrés aux infrastructures"),
    ]
    y = Cm(6.5)
    for title, desc in facts:
        add_text_box(slide, title, Cm(11), y, Cm(9), Cm(1),
                     font_size=14, color=ROUGE, bold=True)
        add_text_box(slide, desc, Cm(11), y + Cm(1), Cm(9), Cm(1.5),
                     font_size=10, color=NOIR)
        y += Cm(3.2)
    # Encadré réglementation
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(11), Cm(23), Cm(9), Cm(4))
    rect.fill.solid()
    rect.fill.fore_color.rgb = GRIS_CLAIR
    rect.line.fill.background()
    add_text_box(slide, "Décret n°2-12-349\ndu 20 mars 2013", Cm(11.5), Cm(23.3), Cm(8), Cm(1.5),
                 font_size=11, color=ROUGE_FONCE, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, "Cadre réglementaire marchés publics\nAppel d'offres ouvert, restreint\net marché négocié", Cm(11.5), Cm(24.8), Cm(8), Cm(2),
                 font_size=9, color=GRIS_FONCE, align=PP_ALIGN.CENTER)
    # Paragraphe explicatif MOA
    add_text_box(slide, "En tant que maître d'ouvrage public, le Conseil Régional est le donneur d'ordre qui "
                 "définit les besoins, finance les projets et contrôle leur bonne exécution. Mon rôle au sein "
                 "de cette structure m'a permis de comprendre toute la chaîne de la commande publique : de la "
                 "programmation budgétaire à la réception des ouvrages, en passant par la passation des marchés, "
                 "l'analyse des offres et le suivi financier. Cette expérience côté MOA est rare pour un technicien "
                 "et constitue un atout majeur pour comprendre les attentes du maître d'ouvrage.",
                 Cm(0.5), Cm(27), Cm(10), Cm(2.5), font_size=7, color=GRIS_FONCE)
    add_page_number(slide, 8)
    print("  Page 8 : Conseil Régional OK")

def page_09_organigramme(prs):
    """Page 9 - Organigramme."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_rect(slide, BLANC)
    add_text_box(slide, "MON POSTE AU SEIN DE L'AGENCE", Cm(1), Cm(1), Cm(19), Cm(1.5),
                 font_size=22, color=ROUGE, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, "Agence d'Exécution des Projets", Cm(1), Cm(2.5), Cm(19), Cm(1),
                 font_size=14, color=GRIS_FONCE, align=PP_ALIGN.CENTER)
    # Organigramme simplifié
    def add_org_box(text, x, y, w, h, is_me=False):
        rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        rect.fill.solid()
        rect.fill.fore_color.rgb = ROUGE if is_me else BLANC
        if not is_me:
            rect.line.color.rgb = GRIS
            rect.line.width = Pt(1)
        else:
            rect.line.fill.background()
        add_text_box(slide, text, x + Cm(0.2), y + Cm(0.2), w - Cm(0.4), h - Cm(0.4),
                     font_size=9, color=BLANC if is_me else NOIR, bold=True, align=PP_ALIGN.CENTER)

    # Président
    add_org_box("Président du\nConseil Régional", Cm(6.5), Cm(5), Cm(8), Cm(2))
    # Ligne vers directeur
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(10.35), Cm(7), Cm(0.15), Cm(1.5)).fill.solid()
    # Directeur
    add_org_box("Directeur de l'Agence\nM. DOGHMANI", Cm(6.5), Cm(8.5), Cm(8), Cm(2))
    # Lignes vers services
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(10.35), Cm(10.5), Cm(0.15), Cm(1.5)).fill.solid()
    # 3 services
    add_org_box("Service Études\n& Programmation", Cm(0.5), Cm(13), Cm(6), Cm(2.5))
    add_org_box("Service Marchés\n& Contrats", Cm(7.5), Cm(13), Cm(6), Cm(2.5))
    add_org_box("Service Suivi\ndes Travaux", Cm(14.5), Cm(13), Cm(6), Cm(2.5))
    # Lignes horizontales
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(3.5), Cm(12), Cm(14), Pt(1.5)).fill.solid()
    # Ligne vers mon poste
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(17.35), Cm(15.5), Cm(0.15), Cm(2)).fill.solid()
    # Mon poste (en rouge)
    add_org_box("Technicien de Suivi\nBAHAFID Mohamed", Cm(13.5), Cm(17.5), Cm(7.5), Cm(3), is_me=True)
    # Flèche "VOUS ÊTES ICI"
    add_text_box(slide, "MON POSTE", Cm(13.5), Cm(20.8), Cm(7.5), Cm(1),
                 font_size=11, color=ROUGE, bold=True, align=PP_ALIGN.CENTER)
    # Contexte du poste
    add_text_box(slide, "Au sein de l'Agence d'Exécution des Projets, j'intervenais en appui technique "
                 "au Directeur pour le pilotage de 7 marchés publics représentant plus de 100 M DH "
                 "d'investissement en infrastructures.",
                 Cm(1), Cm(22), Cm(19), Cm(1.5), font_size=9, color=GRIS_FONCE, align=PP_ALIGN.CENTER)
    # Missions en bas - enrichies
    add_multi_text(slide, [
        "Missions principales :",
        "- Métrés détaillés et estimations confidentielles de l'administration (8 parties techniques × 4 communes)",
        "- Préparation des DCE : rédaction CPS, vérification cohérence CPS/BPDE, plans et RC",
        "- Analyse technique et financière des offres, participation aux CAO comme membre technique",
        "- Suivi financier : vérification des situations mensuelles, décomptes provisoires, gestion avenants",
        "- Visites de chantier : contrôle de conformité, attachements contradictoires, PV de réception",
        "- Rédaction des ordres de service (OS) d'arrêt et de reprise des travaux",
    ], Cm(1), Cm(23.5), Cm(19), Cm(5), font_size=8, color=GRIS_FONCE)
    add_page_number(slide, 9)
    print("  Page 9 : Organigramme OK")

def page_10_bimco(prs):
    """Page 10 - BIMCO Présentation."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_rect(slide, BLANC)
    # Logo
    add_image_safe(slide, img("logo"), Cm(7), Cm(1), Cm(7))
    add_text_box(slide, "BIMCO", Cm(1), Cm(5.5), Cm(19), Cm(2),
                 font_size=36, color=ROUGE, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, "Expert BIM & Économie de la Construction", Cm(1), Cm(7.5), Cm(19), Cm(1),
                 font_size=16, color=GRIS_FONCE, align=PP_ALIGN.CENTER)
    # Fiche d'identité
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(3), Cm(9.5), Cm(15), Cm(6.5))
    card.fill.solid()
    card.fill.fore_color.rgb = GRIS_CLAIR
    card.line.fill.background()
    infos = [
        "Créé le 9 janvier 2026 | Micro-entreprise",
        "SIREN : 999580053 | SIRET : 99958005300018",
        "Code APE : 7112B – Ingénierie, études techniques",
        "44 rue de la République, 42510 Bussières (Loire)",
        "BIM + COnstruction = BIMCO",
    ]
    y = Cm(10)
    for info in infos:
        add_text_box(slide, info, Cm(4), y, Cm(13), Cm(0.8),
                     font_size=11, color=NOIR, align=PP_ALIGN.CENTER)
        y += Cm(1.1)
    # Mission / value proposition
    add_text_box(slide, "Pourquoi créer BIMCO ? Après 8 ans dans le BTP (MOA publique, chantier, BIM), j'ai constaté "
                 "que les économistes de la construction manquent d'outils numériques adaptés à leurs besoins réels. "
                 "BIMCO allie ma connaissance métier (7 marchés publics, 5 ans de chantier) et ma maîtrise du "
                 "développement logiciel (Python, C#, React) pour concevoir des outils qui automatisent les tâches "
                 "répétitives (métrés, chiffrages, situations de travaux) et laissent au technicien le temps de "
                 "se concentrer sur l'analyse et la prise de décision. Le statut de micro-entreprise me permet "
                 "de tester rapidement mes solutions auprès de professionnels du BTP.",
                 Cm(2), Cm(16.5), Cm(17), Cm(3), font_size=8, color=GRIS_FONCE, align=PP_ALIGN.CENTER)
    # Domaines en cercles avec sous-descriptions
    domaines = [
        ("Projeteur\nBIM", "Revit, Navisworks\nIFC, maquettes 3D"),
        ("Métrés\nQuantitatifs", "BIM + traditionnel\nExtraction automatique"),
        ("Études\nde prix", "Estimations, DPGF\nBordereaux, DQE"),
        ("Outils\nnumériques", "Scripts Python, C#\nPlugins Revit, Apps"),
        ("Suivi\néconomique", "Budgets, décomptes\nSituations, avenants"),
    ]
    x = Cm(1)
    for dom, sub in domaines:
        c = add_circle(slide, x, Cm(19.5), Cm(3.2), ROUGE)
        add_text_box(slide, dom, x, Cm(20.1), Cm(3.2), Cm(1.8),
                     font_size=8, color=BLANC, bold=True, align=PP_ALIGN.CENTER)
        add_text_box(slide, sub, x, Cm(23), Cm(3.2), Cm(1.2),
                     font_size=7, color=GRIS_FONCE, align=PP_ALIGN.CENTER)
        x += Cm(3.8)
    # Bannière
    add_image_safe(slide, img("banniere"), Cm(3), Cm(24.5), Cm(15))
    add_page_number(slide, 10)
    print("  Page 10 : BIMCO OK")

def page_11_outils(prs):
    """Page 11 - Outils et compétences techniques."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Fond sombre
    add_bg_rect(slide, RGBColor(0x1A, 0x1A, 0x2E))
    add_text_box(slide, "NOS OUTILS ET MÉTHODES", Cm(1), Cm(1), Cm(19), Cm(1.5),
                 font_size=22, color=BLANC, bold=True, align=PP_ALIGN.CENTER)
    # Grille de logiciels avec barres de niveau
    outils = [
        ("Revit", "Modélisation BIM Architecture + Structure, extraction de quantités", 90, "Expert"),
        ("Navisworks", "Coordination et synthèse de maquettes, détection de clashs", 85, "Expert"),
        ("Dynamo", "Scripts visuels pour automatiser les workflows Revit", 80, "Autonome"),
        ("Python", "Automatisation de tâches BTP, traitement de données, scripts", 80, "Autonome"),
        ("C# / API Revit", "Développement de plugins Revit personnalisés", 75, "Autonome"),
        ("Excel avancé", "Métrés, estimations, bases de prix, tableaux de suivi financier", 95, "Expert"),
        ("AutoCAD", "Lecture et exploitation de plans, pièces dessinées", 85, "Expert"),
        ("MS Project", "Planification de travaux, suivi d'avancement", 70, "Autonome"),
    ]
    y = Cm(4)
    for name, desc, pct, level in outils:
        # Nom du logiciel
        add_text_box(slide, name, Cm(1.5), y, Cm(5), Cm(0.8),
                     font_size=12, color=BLANC, bold=True)
        add_text_box(slide, desc, Cm(1.5), y + Cm(0.8), Cm(5), Cm(0.6),
                     font_size=8, color=GRIS)
        # Barre de progression
        bar_w = Cm(10)
        bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(7.5), y + Cm(0.3), bar_w, Cm(0.6))
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(0x2D, 0x2D, 0x44)
        bg.line.fill.background()
        fill_w = int(bar_w * pct / 100)
        bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(7.5), y + Cm(0.3), fill_w, Cm(0.6))
        bar.fill.solid()
        bar.fill.fore_color.rgb = ROUGE
        bar.line.fill.background()
        # Niveau
        add_text_box(slide, f"{level} ({pct}%)", Cm(18), y + Cm(0.2), Cm(3), Cm(0.8),
                     font_size=8, color=ROUGE_CLAIR)
        y += Cm(2.8)
    # Encadré protocole BIM (conformité T23 - Appliquer un protocole de collaboration)
    bim_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(1), Cm(25.5), Cm(19), Cm(3.7))
    bim_bg.fill.solid()
    bim_bg.fill.fore_color.rgb = RGBColor(0x12, 0x12, 0x22)
    bim_bg.line.color.rgb = ROUGE
    bim_bg.line.width = Pt(1.5)
    add_text_box(slide, "PROTOCOLE DE COLLABORATION BIM APPLIQUÉ", Cm(1.5), Cm(25.7), Cm(18), Cm(0.7),
                 font_size=10, color=ROUGE_CLAIR, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, "Convention BIM : niveaux LOD 300 / LOI 3 | Échange via IFC 2x3 sur plateforme collaborative\n"
                 "Workflow : Modélisation Revit → Export IFC → Extraction quantités → Chiffrage DPGF → Reporting client\n"
                 "Interopérabilité Open BIM : coordination architecte / structure / économiste via maquette fédérée",
                 Cm(1.5), Cm(26.5), Cm(18), Cm(2.5), font_size=8, color=GRIS, align=PP_ALIGN.CENTER)
    add_page_number(slide, 11)
    print("  Page 11 : Outils techniques OK")

def page_12_app_gestion(prs):
    """Page 12 - Application Gestion Chantiers."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_rect(slide, BLANC)
    add_text_box(slide, "RÉALISATION PHARE", Cm(1), Cm(0.8), Cm(19), Cm(1),
                 font_size=14, color=GRIS, align=PP_ALIGN.CENTER)
    add_text_box(slide, "Application « Gestion Chantiers »", Cm(1), Cm(2), Cm(19), Cm(2),
                 font_size=24, color=ROUGE, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, "gestion.bimco-consulting.fr", Cm(1), Cm(4), Cm(19), Cm(0.6),
                 font_size=12, color=GRIS_FONCE, align=PP_ALIGN.CENTER)
    # Paragraphe de présentation
    add_text_box(slide, "Cette application est née d'un constat terrain : sur mes chantiers et marchés publics, "
                 "le suivi financier se faisait sur des tableurs Excel dispersés, sans vision consolidée. "
                 "J'ai conçu et développé intégralement cette solution qui couvre tout le cycle de gestion "
                 "d'une entreprise BTP : du devis initial à la facturation finale, en passant par le suivi "
                 "des situations de travaux, la gestion des équipes et le contrôle des approvisionnements. "
                 "Cette réalisation démontre ma capacité à analyser un besoin métier, à le formaliser en "
                 "spécifications fonctionnelles et à le développer en outil opérationnel.",
                 Cm(1.5), Cm(4.8), Cm(18), Cm(2.5), font_size=8, color=GRIS_FONCE, align=PP_ALIGN.CENTER)
    # Photo de l'appli (placeholder avec cadre)
    placeholder = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(2), Cm(6), Cm(17), Cm(9))
    placeholder.fill.solid()
    placeholder.fill.fore_color.rgb = GRIS_CLAIR
    placeholder.line.color.rgb = GRIS
    add_text_box(slide, "[CAPTURES D'ÉCRAN DE L'APPLICATION\nÀ FOURNIR PAR LE CANDIDAT]",
                 Cm(4), Cm(9), Cm(13), Cm(3), font_size=14, color=GRIS, align=PP_ALIGN.CENTER)
    # 8 modules en grille avec descriptions
    modules = [
        ("Devis", "Bibliothèque d'ouvrages,\nimport/export Excel"),
        ("Chantiers", "Multi-projets, cartes,\nKanban, budget 12 postes"),
        ("Facturation", "Situations de travaux,\navancement, OCR, relances"),
        ("Équipes", "Planning hebdo,\naffectations, compétences"),
        ("Main d'œuvre", "Pointage, planification,\nvariables de paie"),
        ("Finances", "30+ indicateurs,\nrentabilité, tableaux de bord"),
        ("Documents", "GED, fiches intervention\nphoto + signature"),
        ("Appro.", "Fournisseurs, catalogue,\ncommandes, stocks"),
    ]
    x = Cm(1.5)
    for i, (mod, desc) in enumerate(modules):
        if i == 4:
            x = Cm(1.5)
        y_mod = Cm(16.5) if i < 4 else Cm(20.5)
        c = add_circle(slide, x, y_mod, Cm(2), ROUGE)
        add_text_box(slide, mod, x, y_mod + Cm(0.3), Cm(2), Cm(1.2),
                     font_size=7, color=BLANC, bold=True, align=PP_ALIGN.CENTER)
        add_text_box(slide, desc, x - Cm(0.5), y_mod + Cm(2.2), Cm(3), Cm(1.5),
                     font_size=6, color=GRIS_FONCE, align=PP_ALIGN.CENTER)
        x += Cm(4.7)
    # Stack technique
    add_text_box(slide, "Stack technique : React/TypeScript | Node.js/Express | PostgreSQL | Electron | Docker",
                 Cm(1), Cm(25.5), Cm(19), Cm(0.8), font_size=9, color=GRIS_FONCE, align=PP_ALIGN.CENTER)
    add_text_box(slide, "Déployé sur serveur NAS via Docker | Version web + application desktop Windows",
                 Cm(1), Cm(26.5), Cm(19), Cm(0.8), font_size=9, color=GRIS, align=PP_ALIGN.CENTER)
    add_page_number(slide, 12)
    print("  Page 12 : App Gestion Chantiers OK")

def page_13_section2(prs):
    """Page 13 - Séparateur Partie 2."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_separator(slide, "02", "ACTIVITÉS ET PROJETS RÉALISÉS",
                          "Projet 1 : Mise à niveau 4 communes (53,5 M DH – 8 corps d'état)\nProjet 2 : Route Lehri-Kerrouchen 25 km (29 M DH – zone montagneuse)",
                          img("projet1"))
    print("  Page 13 : Séparateur Partie 2 OK")

def page_14_projet1_fiche(prs):
    """Page 14 - Projet 1 : Fiche d'identité."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_rect(slide, BLANC)
    add_text_box(slide, "PROJET 1", Cm(1), Cm(0.8), Cm(19), Cm(1),
                 font_size=14, color=GRIS, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, "Mise à niveau des centres\nde 4 communes", Cm(1), Cm(2), Cm(19), Cm(3),
                 font_size=26, color=ROUGE, bold=True, align=PP_ALIGN.CENTER)
    # Contexte du projet
    add_text_box(slide, "Ce projet s'inscrit dans le Programme de mise à niveau des centres émergents, lancé par le "
                 "Conseil Régional pour moderniser les infrastructures urbaines dans les communes rurales de la "
                 "province de Khénifra. L'objectif était de doter ces centres de voirie, trottoirs, assainissement, "
                 "éclairage public, signalisation, murs de soutènement, espaces verts et mobilier urbain. "
                 "La particularité de ce marché résidait dans son caractère multi-sites (4 communes) et "
                 "multi-techniques (8 corps d'état), ce qui exigeait une coordination rigoureuse entre les "
                 "métrés, les budgets et les plannings de chaque commune, tout en garantissant l'équité de "
                 "traitement entre les sites.",
                 Cm(0.5), Cm(5), Cm(20), Cm(2.2), font_size=8, color=GRIS_FONCE, align=PP_ALIGN.CENTER)
    # Photo
    add_image_safe(slide, img("chantier_urban"), Cm(0.5), Cm(7), Cm(10), Cm(6))
    # Fiche à droite
    fiche = [
        ("Marché", "n°38-RBK-2017 (Lot 4)"),
        ("Maître d'ouvrage", "Conseil Régional BMK"),
        ("Localisation", "Province de Khénifra"),
        ("4 communes", "El Hammam, Kerrouchen,\nOuaoumana, Sebt Ait Rahou"),
        ("Nature", "Aménagement urbain - VRD"),
        ("8 parties", "Assainissement, chaussée, trottoirs,\nsignalisation, éclairage, murs,\npaysager, mobilier urbain"),
    ]
    y = Cm(7.5)
    for label, value in fiche:
        add_text_box(slide, label, Cm(11.5), y, Cm(8), Cm(0.7),
                     font_size=9, color=ROUGE, bold=True)
        add_text_box(slide, value, Cm(11.5), y + Cm(0.7), Cm(8), Cm(1.5),
                     font_size=10, color=NOIR)
        y += Cm(2.5)
    # Grand chiffre
    big_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(2), Cm(22), Cm(17), Cm(4))
    big_bg.fill.solid()
    big_bg.fill.fore_color.rgb = ROUGE
    big_bg.line.fill.background()
    add_text_box(slide, "53,5 M DH TTC", Cm(2), Cm(22.3), Cm(17), Cm(2),
                 font_size=36, color=BLANC, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, "soit 4,86 M€ | 44,6 M DH HT + TVA 20% | Lot unique | 4 communes",
                 Cm(2), Cm(24.3), Cm(17), Cm(0.8), font_size=11, color=BLANC, align=PP_ALIGN.CENTER)
    add_text_box(slide, "Pièces du marché : CPS + RC + BPDE + Plans | Appel d'offres ouvert",
                 Cm(2), Cm(25.2), Cm(17), Cm(0.8), font_size=9, color=RGBColor(0xDD, 0xDD, 0xDD),
                 align=PP_ALIGN.CENTER)
    add_page_number(slide, 14)
    print("  Page 14 : Projet 1 fiche OK")

def page_15_repartition_communes(prs):
    """Page 15 - Répartition budgétaire par commune (graphique)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_rect(slide, BLANC)
    add_text_box(slide, "RÉPARTITION BUDGÉTAIRE PAR COMMUNE", Cm(1), Cm(1), Cm(19), Cm(1.5),
                 font_size=20, color=ROUGE, bold=True, align=PP_ALIGN.CENTER)
    # Graphique Pie Chart
    chart_data = CategoryChartData()
    chart_data.categories = ['El Hammam', 'Kerrouchen', 'Ouaoumana', 'Sebt Ait Rahou']
    chart_data.add_series('Budget HT (M DH)', (6.6, 7.3, 15.8, 14.8))
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, Cm(1), Cm(4), Cm(10), Cm(12), chart_data
    ).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.font.size = Pt(10)
    plot = chart.plots[0]
    # Couleurs des parts
    colors = [ROUGE_CLAIR, ROUGE, ROUGE_FONCE, RGBColor(0x44, 0x44, 0x44)]
    for i, color in enumerate(colors):
        point = plot.series[0].points[i]
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = color
    # Détails à droite
    communes = [
        ("El Hammam", "6,6 M DH HT", "14,9%", "7,96 M DH TTC"),
        ("Kerrouchen", "7,3 M DH HT", "16,5%", "8,80 M DH TTC"),
        ("Ouaoumana", "15,8 M DH HT", "35,5%", "19,0 M DH TTC"),
        ("Sebt Ait Rahou", "14,8 M DH HT", "33,2%", "17,8 M DH TTC"),
    ]
    y = Cm(5)
    for nom, montant, pct, ttc in communes:
        add_text_box(slide, nom, Cm(12), y, Cm(8), Cm(0.8),
                     font_size=12, color=NOIR, bold=True)
        add_text_box(slide, f"{montant} ({pct})", Cm(12), y + Cm(0.7), Cm(8), Cm(0.6),
                     font_size=10, color=GRIS_FONCE)
        add_text_box(slide, ttc, Cm(12), y + Cm(1.3), Cm(8), Cm(0.6),
                     font_size=8, color=GRIS)
        y += Cm(2.8)
    # Callout
    callout = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(2), Cm(20), Cm(17), Cm(3))
    callout.fill.solid()
    callout.fill.fore_color.rgb = ROUGE
    callout.line.fill.background()
    add_text_box(slide, "68% du budget concentré sur 2 communes\nOuaoumana + Sebt Ait Rahou = 30,6 M DH HT",
                 Cm(3), Cm(20.3), Cm(15), Cm(2.5), font_size=14, color=BLANC, bold=True, align=PP_ALIGN.CENTER)
    # Analyse
    add_text_box(slide, "Cette répartition inégale s'explique par l'étendue des centres urbains : Ouaoumana et Sebt Ait Rahou "
                 "disposent de linéaires de voirie plus importants et d'un réseau d'assainissement plus développé. "
                 "En tant que technicien, j'ai dû établir des métrés séparés par commune, car le Bordereau des Prix "
                 "Détail Estimatif (BPDE) est unique mais chaque enveloppe communale devait être respectée "
                 "individuellement. Tout transfert de crédit entre communes nécessitait un avenant, ce qui imposait "
                 "une grande rigueur dans le suivi financier. Cette contrainte m'a appris l'importance de la "
                 "traçabilité des engagements par poste budgétaire et par site géographique.",
                 Cm(1), Cm(23.5), Cm(19), Cm(3), font_size=8, color=GRIS_FONCE, align=PP_ALIGN.CENTER)
    # Montant total
    add_text_box(slide, "Montant total HT : 44 605 581 DH | TTC : 53 526 697 DH (TVA 20%) | soit 4,86 M€",
                 Cm(1), Cm(26.8), Cm(19), Cm(0.8), font_size=8, color=GRIS, align=PP_ALIGN.CENTER)
    add_page_number(slide, 15)
    print("  Page 15 : Répartition communes OK")

def page_16_parties_techniques(prs):
    """Page 16 - 8 parties techniques."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_rect(slide, BLANC)
    add_text_box(slide, "8 CORPS D'ÉTAT", Cm(1), Cm(1), Cm(19), Cm(1),
                 font_size=24, color=ROUGE, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, "Le CPS (Cahier des Prescriptions Spéciales) structurait le marché en 8 parties techniques, "
                 "chacune regroupant des ouvrages de même nature. Cette décomposition permettait un suivi financier "
                 "détaillé par poste et facilitait le contrôle des quantités exécutées. L'ordre d'exécution était "
                 "imposé par des contraintes techniques : l'assainissement (réseaux enterrés) devait impérativement "
                 "précéder la chaussée et les trottoirs. J'ai coordonné les métrés de ces 8 parties pour chacune "
                 "des 4 communes, soit 32 sous-ensembles à établir et à suivre tout au long du chantier.",
                 Cm(1), Cm(2.2), Cm(19), Cm(2), font_size=7, color=GRIS_FONCE, align=PP_ALIGN.CENTER)
    parties = [
        ("01", "Assainissement", "Tranchées, buses PEHD/PVC,\nregards, bouches d'égout"),
        ("02", "Chaussée", "Terrassement, GNF, GNA,\nenrobés bitumineux"),
        ("03", "Trottoirs", "Bordures T1/T3, carreaux\nstriés, pavés, béton"),
        ("04", "Signalisation", "Marquage au sol, panneaux,\npeinture bordures"),
        ("05", "Éclairage public", "Tranchées, tubes annelés,\nmassifs candélabres, câbles"),
        ("06", "Murs et ouvrages", "Béton armé, maçonnerie\nmoellons, gabions"),
        ("07", "Aménagement paysager", "Terre végétale, réseau\nd'arrosage, plantation"),
        ("08", "Mobilier urbain", "Corbeilles, bancs\nen granite"),
    ]
    for i, (num, title, desc) in enumerate(parties):
        col = i % 4
        row = i // 4
        x = Cm(0.8) + col * Cm(5)
        y = Cm(4) + row * Cm(11)
        # Carte
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Cm(4.5), Cm(9))
        card.fill.solid()
        nuance = max(0, 0xC0 - i * 0x10)
        card.fill.fore_color.rgb = RGBColor(nuance, 0x39 if nuance > 0x80 else 0x20, 0x2B if nuance > 0x80 else 0x18)
        card.line.fill.background()
        add_text_box(slide, num, x, y + Cm(1), Cm(4.5), Cm(2.5),
                     font_size=28, color=BLANC, bold=True, align=PP_ALIGN.CENTER)
        add_text_box(slide, title, x, y + Cm(4), Cm(4.5), Cm(2),
                     font_size=11, color=BLANC, bold=True, align=PP_ALIGN.CENTER)
        add_text_box(slide, desc, x, y + Cm(6), Cm(4.5), Cm(2),
                     font_size=8, color=BLANC, align=PP_ALIGN.CENTER)
    add_page_number(slide, 16)
    print("  Page 16 : 8 parties techniques OK")

def page_17_missions(prs):
    """Page 17 - Mes missions sur le Projet 1."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_rect(slide, GRIS_CLAIR)
    add_text_box(slide, "MES MISSIONS SUR LE PROJET 1", Cm(1), Cm(1), Cm(19), Cm(1.5),
                 font_size=22, color=ROUGE, bold=True, align=PP_ALIGN.CENTER)
    missions = [
        ("A", "Estimation de\nl'administration", "Métrés détaillés des 8 parties\npour 4 communes\nSurfaces, linéaires, cubatures\nPrix unitaires de référence\nTableurs Excel structurés"),
        ("B", "Préparation\ndu DCE", "Rédaction du CPS (clauses\nadmin. et techniques)\nRC, BPDE, plans et pièces\nCohérence CPS / BPDE\nNomenclature des prix"),
        ("C", "Analyse des\noffres", "Vérification arithmétique\nDétection prix anormaux\nTableau comparatif des offres\nGrille de notation (/100)\nParticipation à la CAO"),
        ("D", "Suivi\nfinancier", "Situations de travaux\nmensuelles vérifiées\nDécomptes provisoires\nGestion des avenants\nTableau de bord par commune"),
    ]
    x = Cm(0.5)
    for letter, title, desc in missions:
        # Carte avec ombre
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Cm(4), Cm(4.7), Cm(16))
        card.fill.solid()
        card.fill.fore_color.rgb = BLANC
        card.line.fill.background()
        add_shadow(card, alpha_pct=18)
        # Lettre en cercle rouge
        c = add_circle(slide, x + Cm(1.6), Cm(4.5), Cm(1.5), ROUGE)
        add_text_box(slide, letter, x + Cm(1.6), Cm(4.7), Cm(1.5), Cm(1.2),
                     font_size=16, color=BLANC, bold=True, align=PP_ALIGN.CENTER)
        # Titre
        add_text_box(slide, title, x + Cm(0.3), Cm(6.5), Cm(4.1), Cm(2.5),
                     font_size=12, color=ROUGE, bold=True, align=PP_ALIGN.CENTER)
        # Description
        add_text_box(slide, desc, x + Cm(0.3), Cm(9.5), Cm(4.1), Cm(8),
                     font_size=9, color=NOIR, align=PP_ALIGN.CENTER)
        # Flèche vers la droite (sauf dernier)
        if letter != "D":
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + Cm(4.7), Cm(11), Cm(0.7), Cm(1))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = ROUGE_CLAIR
            arrow.line.fill.background()
        x += Cm(5.1)
    # Légende enrichie en bas – conformité C18, C19-2
    add_multi_text(slide, [
        "Cycle complet : de l'estimation à la réception provisoire (PV signé) et levée des réserves",
        "Réunions de chantier hebdomadaires avec l'entreprise retenue pour valider l'avancement et les attachements",
        "Suivi des approvisionnements : vérification des livraisons (enrobés, buses PEHD, candélabres, bordures)",
        "Communication : reporting régulier au Directeur de l'Agence, comptes rendus écrits après chaque visite de site",
        "Confidentialité : respect du secret des délibérations de la CAO et de la confidentialité des offres de prix",
        "Cadre réglementaire : Décret n°2-12-349 du 20/03/2013 | 8 parties × 4 communes = 32 sous-ensembles",
    ], Cm(1), Cm(21.5), Cm(19), Cm(5.5), font_size=7, color=GRIS_FONCE)
    add_page_number(slide, 17)
    print("  Page 17 : Missions Projet 1 OK")

def page_18_metres_ouaoumana(prs):
    """Page 18 - Métrés et chiffres clés Ouaoumana."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_rect(slide, BLANC)
    add_text_box(slide, "FOCUS : MÉTRÉS - COMMUNE D'OUAOUMANA", Cm(1), Cm(0.8), Cm(19), Cm(1),
                 font_size=20, color=ROUGE, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, "Ouaoumana concentre 35,5% du budget total (15,8 M DH HT). Les métrés détaillés "
                 "ci-dessous couvrent les principaux ouvrages des 8 parties techniques. Chaque quantité "
                 "est calculée à partir des plans du bureau d'études et vérifiée par relevés sur site.",
                 Cm(1), Cm(2), Cm(19), Cm(1.5), font_size=8, color=GRIS_FONCE, align=PP_ALIGN.CENTER)
    chiffres = [
        ("3 620 m³", "Tranchées"),
        ("2 370 ml", "Buses PEHD"),
        ("56", "Regards de visite"),
        ("2 175 T", "Enrobés"),
        ("11 500 m²", "Carreaux striés"),
        ("163", "Massifs candélabres"),
        ("4 800 ml", "Bordures T3"),
    ]
    for i, (val, label) in enumerate(chiffres):
        col = i % 3
        row = i // 3
        x = Cm(1.5) + col * Cm(6.5)
        y = Cm(4) + row * Cm(7)
        # Grande valeur en rouge
        add_text_box(slide, val, x, y, Cm(5.5), Cm(2.5),
                     font_size=28, color=ROUGE, bold=True, align=PP_ALIGN.CENTER)
        add_text_box(slide, label, x, y + Cm(2.5), Cm(5.5), Cm(1.5),
                     font_size=12, color=NOIR, align=PP_ALIGN.CENTER)
        # Ligne de séparation
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Cm(1), y + Cm(4.5), Cm(3.5), Pt(1))
        line.fill.solid()
        line.fill.fore_color.rgb = GRIS_CLAIR
        line.line.fill.background()
    # Méthode de métré enrichie
    add_text_box(slide, "Méthodologie appliquée : les métrés ont été établis en deux phases. D'abord, un calcul sur "
                 "plans à partir des pièces graphiques fournies par le bureau d'études (plans de masse, profils en "
                 "long, coupes types). Ensuite, une vérification par relevés sur site pour ajuster les quantités "
                 "aux conditions réelles du terrain. Chaque métré est structuré dans un tableur Excel avec formules "
                 "vérifiables : le jury peut tracer chaque quantité depuis le plan jusqu'au prix total. Cette "
                 "traçabilité est fondamentale pour défendre l'estimation devant la Commission d'Appel d'Offres.",
                 Cm(1.5), Cm(19.5), Cm(18), Cm(3), font_size=8, color=GRIS_FONCE, align=PP_ALIGN.CENTER)
    # Citation en bas
    quote = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(2), Cm(22.5), Cm(17), Cm(2.5))
    quote.fill.solid()
    quote.fill.fore_color.rgb = GRIS_CLAIR
    quote.line.fill.background()
    add_text_box(slide, "« La précision des avant-métrés est fondamentale :\nune erreur de 5% sur les trottoirs peut représenter +1 million de DH d'écart. »",
                 Cm(3), Cm(22.7), Cm(15), Cm(2.2), font_size=11, color=ROUGE,
                 bold=True, align=PP_ALIGN.CENTER)
    # Photo en bande
    add_image_safe(slide, img("vrd"), Cm(0), Cm(26), A4_W, Cm(3.7))
    add_page_number(slide, 18)
    print("  Page 18 : Métrés Ouaoumana OK")

def page_19_analyse_offres(prs):
    """Page 19 - Analyse des offres et CAO."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_rect(slide, BLANC)
    add_text_box(slide, "ANALYSE DES OFFRES & CAO", Cm(1), Cm(1), Cm(19), Cm(1),
                 font_size=22, color=ROUGE, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, "Mon rôle lors de la Commission d'Appel d'Offres (CAO) : en tant que membre technique, "
                 "j'ai préparé l'analyse comparative des offres reçues. Concrètement, cela impliquait la "
                 "vérification arithmétique de chaque bordereau (détection d'erreurs de calcul), l'identification "
                 "des prix anormalement bas ou excessifs par comparaison avec l'estimation de l'administration "
                 "(que j'avais moi-même établie), et la notation sur 100 points selon les critères du règlement "
                 "de consultation. La confidentialité de cette analyse est impérative : les offres ne sont "
                 "divulguées qu'en séance publique lors de l'ouverture des plis.",
                 Cm(1), Cm(2.2), Cm(19), Cm(1.5), font_size=7, color=GRIS_FONCE, align=PP_ALIGN.CENTER)
    # Graphique barres - notation
    chart_data = CategoryChartData()
    chart_data.categories = ['Offre financière\n(/70)', 'Moyens humains\n(/10)', 'Moyens matériels\n(/10)', 'Références\n(/10)']
    chart_data.add_series('Entreprise A', (62, 7, 8, 6))
    chart_data.add_series('Entreprise B (retenue)', (68, 9, 9, 8))
    chart_data.add_series('Entreprise C', (55, 6, 7, 5))
    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, Cm(1), Cm(3.5), Cm(19), Cm(10), chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.font.size = Pt(9)
    # Couleurs des séries
    series_colors = [GRIS, ROUGE, GRIS_CLAIR]
    for i, color in enumerate(series_colors):
        chart.series[i].format.fill.solid()
        chart.series[i].format.fill.fore_color.rgb = color
    # Résultat
    result_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(2), Cm(14.5), Cm(17), Cm(2.5))
    result_bg.fill.solid()
    result_bg.fill.fore_color.rgb = ROUGE
    result_bg.line.fill.background()
    add_text_box(slide, "Entreprise B retenue avec 94/100", Cm(2), Cm(14.7), Cm(17), Cm(1),
                 font_size=18, color=BLANC, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, "Critères : offre financière (70 pts) + moyens humains (10) + matériels (10) + références (10)",
                 Cm(3), Cm(15.8), Cm(15), Cm(1), font_size=9, color=RGBColor(0xDD, 0xDD, 0xDD),
                 align=PP_ALIGN.CENTER)
    # Documents officiels
    add_text_box(slide, "Documents officiels de la procédure", Cm(1), Cm(17.8), Cm(19), Cm(1),
                 font_size=12, color=GRIS_FONCE, bold=True, align=PP_ALIGN.CENTER)
    add_image_safe(slide, img("cao"), Cm(1.5), Cm(19), Cm(8.5), Cm(7))
    add_image_safe(slide, img("rejet"), Cm(11), Cm(19), Cm(8.5), Cm(7))
    add_text_box(slide, "Convocation officielle de la\nCommission d'Appel d'Offres\nRégion Béni Mellal-Khénifra", Cm(1.5), Cm(26.2), Cm(8.5), Cm(1.5),
                 font_size=7, color=GRIS, align=PP_ALIGN.CENTER)
    add_text_box(slide, "Notification de rejet d'offre\nProcédure de marchés publics\nmarocains (Décret n°2-12-349)", Cm(11), Cm(26.2), Cm(8.5), Cm(1.5),
                 font_size=7, color=GRIS, align=PP_ALIGN.CENTER)
    add_page_number(slide, 19)
    print("  Page 19 : Analyse offres OK")

def page_20_suivi_kerrouchen(prs):
    """Page 20 - Suivi financier Kerrouchen."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_rect(slide, BLANC)
    add_text_box(slide, "SUIVI FINANCIER - KERROUCHEN", Cm(1), Cm(1), Cm(19), Cm(1),
                 font_size=22, color=ROUGE, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, "Pendant l'exécution, j'ai assuré la vérification des situations de travaux mensuelles, "
                 "le contrôle des quantités déclarées par rapprochement avec l'avancement réel sur chantier, "
                 "et le suivi des approvisionnements (livraisons d'enrobés, buses, candélabres). "
                 "À la fin des travaux, j'ai participé à la réception provisoire (PV signé avec réserves) "
                 "et au suivi de la levée des réserves avant le décompte définitif.",
                 Cm(1), Cm(2.2), Cm(19), Cm(1.8), font_size=8, color=GRIS_FONCE, align=PP_ALIGN.CENTER)
    # Graphique barres empilées
    chart_data = CategoryChartData()
    chart_data.categories = ['Assainis.', 'Chaussée', 'Trottoirs', 'Signal.', 'Éclairage', 'Murs', 'Paysager', 'Mobilier']
    chart_data.add_series('Marché (DH)', (856712, 1890645, 1756330, 134528, 902450, 345678, 289156, 161415))
    chart_data.add_series('Exécuté (DH)', (878234, 1923456, 1812567, 128976, 915234, 367890, 245678, 122460))
    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, Cm(1), Cm(4), Cm(19), Cm(12), chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.series[0].format.fill.solid()
    chart.series[0].format.fill.fore_color.rgb = GRIS
    chart.series[1].format.fill.solid()
    chart.series[1].format.fill.fore_color.rgb = ROUGE
    # Chiffres clés
    chiffres = [
        ("Marché", "7 336 914 DH HT"),
        ("Exécuté", "7 394 495 DH HT"),
        ("Écart", "+57 581 DH (+0,8%)"),
    ]
    x = Cm(1.5)
    for label, val in chiffres:
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Cm(18), Cm(6), Cm(3.5))
        card.fill.solid()
        card.fill.fore_color.rgb = GRIS_CLAIR
        card.line.fill.background()
        add_text_box(slide, label, x, Cm(18.3), Cm(6), Cm(1),
                     font_size=10, color=GRIS, bold=True, align=PP_ALIGN.CENTER)
        add_text_box(slide, val, x, Cm(19.3), Cm(6), Cm(1.5),
                     font_size=14, color=ROUGE if "+" in val else NOIR, bold=True, align=PP_ALIGN.CENTER)
        x += Cm(6.5)
    # Callout jauge
    gauge_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(2), Cm(22.5), Cm(17), Cm(4))
    gauge_bg.fill.solid()
    gauge_bg.fill.fore_color.rgb = VERT
    gauge_bg.line.fill.background()
    add_text_box(slide, "DÉPASSEMENT MAÎTRISÉ : +0,8%", Cm(2), Cm(22.5), Cm(17), Cm(1),
                 font_size=16, color=BLANC, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, "Le dépassement de +0,8% reste dans la tolérance contractuelle, évitant un avenant.\n"
                 "Les postes Chaussée et Murs ont été en plus-value en raison du terrain rocheux imprévu,\n"
                 "tandis que le poste Paysager a été réduit pour compenser partiellement.\n"
                 "Ce résultat illustre l'importance d'un suivi financier continu : sans mon tableau de bord\n"
                 "hebdomadaire, le dépassement aurait pu atteindre +5% et nécessiter un avenant coûteux\n"
                 "en temps administratif (3 à 6 mois de procédure).",
                 Cm(2.5), Cm(23.5), Cm(16), Cm(3.2), font_size=8, color=BLANC, align=PP_ALIGN.CENTER)
    add_page_number(slide, 20)
    print("  Page 20 : Suivi Kerrouchen OK")

def page_21_difficultes_p1(prs):
    """Page 21 - Difficultés et solutions Projet 1."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_rect(slide, BLANC)
    add_text_box(slide, "DIFFICULTÉS ET SOLUTIONS", Cm(1), Cm(1), Cm(19), Cm(1.5),
                 font_size=22, color=ROUGE, bold=True, align=PP_ALIGN.CENTER)
    defis = [
        ("8 parties techniques", "Coordination complexe entre\nassainissement, chaussée, trottoirs,\nsignalisation, éclairage, murs,\npaysager et mobilier urbain", "Tableau de bord Excel avec suivi\npar commune et par partie.\nMise à jour hebdomadaire après\nchaque visite de chantier"),
        ("Conditions géologiques", "Terrain rocheux imprévu à Kerrouchen\nnécessitant des plus-values de\nterrassement. Nature du sol variable\nentre les 4 communes", "Visites de chantier avec l'entreprise\nB pour constater les conditions.\nAttachements contradictoires signés\nsur place pour justifier les écarts"),
        ("4 chantiers simultanés", "Gestion de 4 sites répartis sur\nla province de Khénifra,\nsuivi dispersé géographiquement", "Réunion hebdomadaire avec le\nDirecteur de l'Agence :\ncompte rendu d'avancement\npar commune et par poste"),
        ("Écarts quantités", "Différences entre les quantités\nestimées au BPDE et les quantités\nréellement exécutées sur chantier", "Attachements contradictoires\nsignés avec l'entreprise B\nsur chantier + reporting\nécrit à la hiérarchie"),
    ]
    y = Cm(4)
    for defi, probleme, solution in defis:
        # Bloc problème (rouge) avec ombre
        pb = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(1), y, Cm(8), Cm(4.5))
        pb.fill.solid()
        pb.fill.fore_color.rgb = RGBColor(0xFF, 0xEB, 0xEB)
        pb.line.fill.background()
        add_shadow(pb, alpha_pct=12)
        add_text_box(slide, defi, Cm(1.5), y + Cm(0.3), Cm(7), Cm(1),
                     font_size=11, color=ROUGE, bold=True)
        add_text_box(slide, probleme, Cm(1.5), y + Cm(1.5), Cm(7), Cm(2.5),
                     font_size=10, color=NOIR)
        # Flèche
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Cm(9.3), y + Cm(1.5), Cm(1.5), Cm(1))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = VERT
        arrow.line.fill.background()
        # Bloc solution (vert) avec ombre
        sol = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(11.2), y, Cm(8.5), Cm(4.5))
        sol.fill.solid()
        sol.fill.fore_color.rgb = RGBColor(0xEB, 0xFF, 0xEB)
        sol.line.fill.background()
        add_shadow(sol, alpha_pct=12)
        add_text_box(slide, "SOLUTION", Cm(11.7), y + Cm(0.3), Cm(7.5), Cm(1),
                     font_size=9, color=VERT, bold=True)
        add_text_box(slide, solution, Cm(11.7), y + Cm(1.5), Cm(7.5), Cm(2.5),
                     font_size=10, color=NOIR)
        y += Cm(5.3)
    # Enseignements pour le jury
    add_text_box(slide, "Enseignements : ces difficultés m'ont appris que le suivi technique d'une opération ne "
                 "se limite pas au bureau. Il exige une présence régulière sur le terrain, une communication "
                 "transparente avec tous les acteurs (entreprise, hiérarchie, bureau d'études) et une "
                 "anticipation permanente des risques financiers. La capacité à rédiger des comptes rendus "
                 "clairs et factuels après chaque visite est un savoir-faire essentiel du technicien MEC.",
                 Cm(1), Cm(26), Cm(19), Cm(2), font_size=8, color=ROUGE, align=PP_ALIGN.CENTER)
    add_page_number(slide, 21)
    print("  Page 21 : Difficultés P1 OK")

def page_22_projet2_fiche(prs):
    """Page 22 - Projet 2 : Fiche d'identité."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Photo pleine page avec overlay semi-transparent
    add_image_safe(slide, img("route"), 0, 0, A4_W, A4_H)
    add_bg_rect(slide, ROUGE_FONCE, 0, 0, A4_W, A4_H, opacity=60)
    # Titre
    add_text_box(slide, "PROJET 2", Cm(2), Cm(3), Cm(17), Cm(2),
                 font_size=20, color=BLANC)
    add_text_box(slide, "ROUTE\nLEHRI-KERROUCHEN", Cm(2), Cm(5), Cm(17), Cm(5),
                 font_size=40, color=BLANC, bold=True)
    # Grands chiffres
    add_text_box(slide, "25 KM", Cm(2), Cm(12), Cm(8), Cm(3),
                 font_size=48, color=BLANC, bold=True)
    add_text_box(slide, "29 M DH", Cm(11), Cm(12), Cm(8), Cm(3),
                 font_size=48, color=ROUGE_CLAIR, bold=True)
    # Fiche avec fond semi-transparent et ombre
    fiche_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(2), Cm(17), Cm(17), Cm(8))
    fiche_bg.fill.solid()
    fiche_bg.fill.fore_color.rgb = BLANC
    fiche_bg.line.fill.background()
    set_transparency(fiche_bg, 90)
    add_shadow(fiche_bg, alpha_pct=25)
    fiche_items = [
        "Marché n°46-RBK-2017 – Programme National des Routes Rurales, 3e tranche (PRR3)",
        "Le PRR3 est un programme stratégique national visant à désenclaver les zones rurales du Maroc",
        "Cette route de 25 km relie Lehri à Kerrouchen en zone montagneuse (Moyen Atlas)",
        "Relief accidenté : 120 334 m³ de déblais, ouvrages hydrauliques, gabions de soutènement",
        "3 sections : Linéaire (23 prix) + Carrefour PK 0+000 (11 prix) + Bretelles (19 prix)",
        "Mon rôle : estimation de l'administration, préparation du DCE, suivi financier d'exécution",
        "Montant HT : 24 169 371 DH (~2,2 M€) | TVA 20% | Maîtrise d'ouvrage : Conseil Régional",
    ]
    y = Cm(17.5)
    for item in fiche_items:
        add_circle(slide, Cm(3), y + Cm(0.2), Cm(0.3), ROUGE)
        add_text_box(slide, item, Cm(4), y, Cm(14), Cm(1),
                     font_size=11, color=NOIR)
        y += Cm(1.4)
    # CPS signé
    add_image_safe(slide, img("cps"), Cm(13), Cm(22.5), Cm(5))
    add_text_box(slide, "CPS signé", Cm(13), Cm(26), Cm(5), Cm(0.8),
                 font_size=8, color=GRIS_FONCE, align=PP_ALIGN.CENTER)
    add_page_number(slide, 22)
    print("  Page 22 : Projet 2 fiche OK")

def page_23_metres_routiers(prs):
    """Page 23 - Métrés routiers."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_rect(slide, BLANC)
    add_text_box(slide, "PRINCIPAUX MÉTRÉS DE LA SECTION LINÉAIRE", Cm(1), Cm(0.8), Cm(19), Cm(1.2),
                 font_size=18, color=ROUGE, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, "Calculs à partir des profils en travers du bureau d'études routier + étude hydrologique",
                 Cm(1), Cm(2), Cm(19), Cm(0.6), font_size=9, color=GRIS_FONCE, align=PP_ALIGN.CENTER)
    add_text_box(slide, "Les cubatures de terrassement (déblais/remblais) ont été calculées à partir des profils en travers "
                 "type. Les quantités d'ouvrages hydrauliques (buses, dalots) ont été déterminées en fonction de "
                 "l'étude hydrologique et du nombre de talwegs traversés. Les quantités de corps de chaussée (GNB, "
                 "GNF2) sont basées sur le profil en travers type avec largeur de plateforme adaptée au trafic.",
                 Cm(1), Cm(2.6), Cm(19), Cm(2), font_size=7, color=GRIS_FONCE, align=PP_ALIGN.CENTER)
    metres = [
        ("120 334 m³", "Déblais"),
        ("76 735 m³", "Remblais"),
        ("19 989 m³", "Couche de base GNB"),
        ("34 481 m³", "Couche de fondation GNF2"),
        ("794 ml", "Buses Ø1000"),
        ("64,5 T", "Acier HA"),
        ("733 m³", "Béton B2"),
        ("4 579 m³", "Béton B3"),
        ("789 m³", "Gabions"),
    ]
    for i, (val, label) in enumerate(metres):
        col = i % 3
        row = i // 3
        x = Cm(1) + col * Cm(6.5)
        y = Cm(4) + row * Cm(6)
        add_text_box(slide, val, x, y, Cm(6), Cm(2),
                     font_size=24, color=ROUGE, bold=True, align=PP_ALIGN.CENTER)
        add_text_box(slide, label, x, y + Cm(2), Cm(6), Cm(1.5),
                     font_size=11, color=NOIR, align=PP_ALIGN.CENTER)
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Cm(1), y + Cm(3.8), Cm(4), Pt(1))
        line.fill.solid()
        line.fill.fore_color.rgb = GRIS_CLAIR
        line.line.fill.background()
    # Analyse comparative
    add_text_box(slide, "Contrairement au Projet 1 (aménagement urbain), les métrés routiers reposent essentiellement sur "
                 "des cubatures calculées à partir de profils en travers. La difficulté spécifique réside dans "
                 "l'incertitude géologique : les volumes de déblais rocheux ne peuvent être connus précisément qu'à "
                 "l'exécution. Mon estimation initiale s'est avérée conforme à ±8% des quantités réelles, ce qui "
                 "est considéré comme une bonne précision pour un projet routier en zone montagneuse.",
                 Cm(1), Cm(22.5), Cm(19), Cm(2), font_size=7, color=GRIS_FONCE, align=PP_ALIGN.CENTER)
    # Photo route en bande bas
    add_image_safe(slide, img("route2"), Cm(0), Cm(25), A4_W, Cm(4.7))
    add_page_number(slide, 23)
    print("  Page 23 : Métrés routiers OK")

def page_24_budget_route(prs):
    """Page 24 - Budget et suivi Projet 2."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_rect(slide, BLANC)
    add_text_box(slide, "RÉPARTITION BUDGÉTAIRE - PROJET 2", Cm(1), Cm(1), Cm(19), Cm(1.5),
                 font_size=20, color=ROUGE, bold=True, align=PP_ALIGN.CENTER)
    # Graphique Pie
    chart_data = CategoryChartData()
    chart_data.categories = ['Terrassement', 'Corps chaussée', 'Revêtement', 'Ouvrages hydr.', 'Soutènement', 'Bretelles']
    chart_data.add_series('Budget (M DH)', (4.3, 7.3, 2.6, 6.8, 0.46, 2.7))
    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, Cm(1), Cm(4), Cm(10), Cm(12), chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    colors_pie = [ROUGE_CLAIR, ROUGE, ROUGE_FONCE, RGBColor(0x44, 0x44, 0x44), GRIS, GRIS_FONCE]
    for i, c in enumerate(colors_pie):
        chart.plots[0].series[0].points[i].format.fill.solid()
        chart.plots[0].series[0].points[i].format.fill.fore_color.rgb = c
    # Détail à droite
    items = [
        ("Terrassement", "4,3 M DH", "17,9%"),
        ("Corps de chaussée", "7,3 M DH", "30,1%"),
        ("Revêtement", "2,6 M DH", "10,7%"),
        ("Ouvrages hydrauliques", "6,8 M DH", "28,1%"),
        ("Soutènement", "0,46 M DH", "1,9%"),
        ("Bretelles + carrefour", "2,7 M DH", "11,3%"),
    ]
    y = Cm(4.5)
    for nom, montant, pct in items:
        add_text_box(slide, nom, Cm(12), y, Cm(8), Cm(0.7), font_size=10, color=NOIR, bold=True)
        add_text_box(slide, f"{montant} ({pct})", Cm(12), y + Cm(0.6), Cm(8), Cm(0.7),
                     font_size=9, color=GRIS_FONCE)
        y += Cm(2)
    # Callout
    callout = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(2), Cm(19), Cm(17), Cm(3))
    callout.fill.solid()
    callout.fill.fore_color.rgb = ROUGE
    callout.line.fill.background()
    add_text_box(slide, "58% du budget sur 2 postes :\nCorps de chaussée + Ouvrages hydrauliques",
                 Cm(3), Cm(19.3), Cm(15), Cm(2.5), font_size=14, color=BLANC,
                 bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, "Analyse pour le jury : la répartition budgétaire d'un projet routier en montagne diffère "
                 "fortement d'un projet urbain. Ici, 58% du budget est absorbé par deux postes (corps de "
                 "chaussée + ouvrages hydrauliques), contre une répartition plus équilibrée en aménagement urbain. "
                 "Le coût kilométrique moyen de 967 000 DH HT/km (88 000 €/km) est cohérent avec les références "
                 "nationales du PRR3 pour les zones montagneuses. Le poste terrassement (120 334 m³ de déblais) "
                 "et les gabions de soutènement (789 m³) traduisent la difficulté du relief. Cette analyse m'a "
                 "appris à adapter mes estimations au contexte géographique et à ne pas appliquer des ratios "
                 "standards sans vérifier leur pertinence locale.",
                 Cm(1), Cm(22.5), Cm(19), Cm(3), font_size=7, color=GRIS_FONCE, align=PP_ALIGN.CENTER)
    add_text_box(slide, "Montant total HT : 24 169 371 DH | TTC : 29 003 246 DH | soit ~2,6 M€ sur 25 km",
                 Cm(1), Cm(25.2), Cm(19), Cm(0.8), font_size=8, color=GRIS, align=PP_ALIGN.CENTER)
    add_page_number(slide, 24)
    print("  Page 24 : Budget route OK")

def page_25_difficultes_route(prs):
    """Page 25 - Défis chantier routier en montagne."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Photo pleine page montagne avec overlay semi-transparent
    add_image_safe(slide, img("montagne"), 0, 0, A4_W, A4_H)
    add_bg_rect(slide, ROUGE_FONCE, 0, 0, A4_W, A4_H, opacity=55)
    add_text_box(slide, "DÉFIS D'UN CHANTIER ROUTIER\nEN ZONE MONTAGNEUSE", Cm(2), Cm(2), Cm(17), Cm(3),
                 font_size=24, color=BLANC, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, "Face à ces défis, j'ai mis en place un suivi rapproché des cubatures par attachements "
                 "contradictoires signés sur chantier avec l'entreprise. Les périodes d'arrêt hivernal étaient "
                 "anticipées dans le planning avec des ordres de service d'arrêt et de reprise. Les métrés "
                 "d'ouvrages hydrauliques ont été ajustés en concertation avec le bureau d'études après "
                 "observation des crues. Ces situations m'ont appris qu'un bon économiste doit savoir adapter "
                 "ses prévisions à la réalité du terrain et communiquer rapidement les écarts à sa hiérarchie.",
                 Cm(2), Cm(5.5), Cm(17), Cm(2.2), font_size=8, color=RGBColor(0xDD, 0xDD, 0xDD),
                 align=PP_ALIGN.CENTER)
    # 4 cartes flottantes avec ombre et transparence
    defis = [
        ("Relief montagneux", "Le tracé de 25 km traverse un\nterrain très accidenté, générant\n120 334 m³ de déblais et nécessitant\ndes gabions de soutènement (789 m³)"),
        ("Conditions climatiques", "Intempéries hivernales imposant\ndes arrêts de chantier : ordres de\nservice d'arrêt et de reprise émis\npour gel et fortes précipitations"),
        ("Écarts de cubatures", "Les cubatures de terrassement\nréelles se sont écartées des\nprévisions en raison de la géologie\nrencontrée (terrain rocheux imprévu)"),
        ("Ouvrages hydrauliques", "Le nombre et dimensionnement des\nouvrages de traversée ont été ajustés\naprès observation des crues pendant\nles travaux (buses, dalots, béton)"),
    ]
    y = Cm(8)
    for i, (title, desc) in enumerate(defis):
        x = Cm(1.5) if i % 2 == 0 else Cm(11)
        if i >= 2:
            y_pos = Cm(17)
        else:
            y_pos = Cm(8)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y_pos, Cm(8.5), Cm(7))
        card.fill.solid()
        card.fill.fore_color.rgb = BLANC
        card.line.fill.background()
        set_transparency(card, 92)
        add_shadow(card, alpha_pct=30)
        # Accent rouge en haut de la carte
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Cm(0.3), y_pos + Cm(0.3), Cm(7.9), Pt(4))
        accent.fill.solid()
        accent.fill.fore_color.rgb = ROUGE
        accent.line.fill.background()
        add_text_box(slide, title, x + Cm(0.5), y_pos + Cm(1), Cm(7.5), Cm(2),
                     font_size=14, color=ROUGE, bold=True, align=PP_ALIGN.CENTER)
        add_text_box(slide, desc, x + Cm(0.5), y_pos + Cm(3.5), Cm(7.5), Cm(3),
                     font_size=11, color=NOIR, align=PP_ALIGN.CENTER)
    add_page_number(slide, 25)
    print("  Page 25 : Défis route OK")

def page_26_complementaires(prs):
    """Page 26 - Activités complémentaires."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_rect(slide, BLANC)
    add_text_box(slide, "ACTIVITÉS COMPLÉMENTAIRES", Cm(1), Cm(1), Cm(19), Cm(1),
                 font_size=22, color=ROUGE, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, "Au-delà des deux projets principaux, la diversité de mes interventions (routes, pistes, "
                 "eau potable, voirie) m'a permis d'acquérir une vision élargie des différents types de "
                 "marchés publics et de leurs spécificités techniques et économiques.",
                 Cm(1), Cm(2.2), Cm(19), Cm(1.5), font_size=8, color=GRIS_FONCE, align=PP_ALIGN.CENTER)
    # Section Maroc
    add_text_box(slide, "5 AUTRES MARCHÉS AU CONSEIL RÉGIONAL", Cm(1), Cm(3.8), Cm(19), Cm(1),
                 font_size=14, color=ROUGE_FONCE, bold=True)
    marches = [
        ("27-RBK-2017", "Construction route village Sidi Bouabbad à Oued Grou", "Route / Piste rurale (CT Sidi Lamine)"),
        ("28-RBK-2017", "Route Ajdir-Ayoun Oum Errabia + Piste Lijon Kichchon", "Route + Piste (CT Sidi Lamine)"),
        ("30-RBK-2017", "Adduction en Eau Potable El Borj – El Hamam", "AEP (réseau d'eau potable)"),
        ("39-RBK-2017", "Pistes Hartaf – Sebt Ait Rahou", "Pistes rurales"),
        ("49-RBK-2016", "Aménagement voie Amghass – Bouchbel", "Voirie / Aménagement"),
    ]
    y = Cm(5)
    for ref, nom, typ in marches:
        add_circle(slide, Cm(2), y + Cm(0.15), Cm(0.3), ROUGE)
        add_text_box(slide, ref, Cm(3), y, Cm(4), Cm(0.7), font_size=9, color=ROUGE, bold=True)
        add_text_box(slide, nom, Cm(7), y, Cm(12), Cm(0.7), font_size=10, color=NOIR)
        add_text_box(slide, typ, Cm(7), y + Cm(0.7), Cm(12), Cm(0.5), font_size=8, color=GRIS)
        y += Cm(1.5)
    # Ligne séparatrice
    sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(2), Cm(12.5), Cm(17), Pt(2))
    sep.fill.solid()
    sep.fill.fore_color.rgb = ROUGE
    sep.line.fill.background()
    # Section France
    add_text_box(slide, "EXPÉRIENCE TERRAIN EN FRANCE", Cm(1), Cm(13.5), Cm(19), Cm(1),
                 font_size=14, color=ROUGE_FONCE, bold=True)
    exp_france = [
        ("2022-2023", "Chef d'équipe GO - Ergalis, Feurs (42)",
         "Lecture de plans, implantation et traçage, montage des banches,\nmise en place des armatures, coulage du béton, suivi des cycles de coffrage"),
        ("2024", "Chef de chantier - Minssieux et Fils, Mornant (69)",
         "Encadrement opérationnel d'équipe, organisation quotidienne,\nsuivi de l'avancement des travaux, contrôle qualité d'exécution"),
    ]
    y = Cm(15)
    for date, poste, detail in exp_france:
        add_text_box(slide, date, Cm(2), y, Cm(3), Cm(0.8), font_size=11, color=ROUGE, bold=True)
        add_text_box(slide, poste, Cm(5.5), y, Cm(14), Cm(0.8), font_size=11, color=NOIR, bold=True)
        add_text_box(slide, detail, Cm(5.5), y + Cm(0.9), Cm(14), Cm(1.5), font_size=9, color=GRIS_FONCE)
        y += Cm(3)
    # Photo gros oeuvre
    add_image_safe(slide, img("gros_oeuvre"), Cm(2), Cm(21.5), Cm(17), Cm(6))
    # Encadré apports enrichi
    add_text_box(slide, "Ce que cette expérience terrain apporte à mon profil d'économiste : avoir coulé du béton, "
                 "monté des banches et ferraillé des voiles me permet aujourd'hui de chiffrer un ouvrage en "
                 "connaissant les temps réels de main-d'œuvre, les rendements d'équipe, les pertes de matériaux "
                 "et les contraintes logistiques. Un économiste qui n'a jamais vu un chantier risque de produire "
                 "des estimations déconnectées de la réalité. Ma double expérience MOA + exécution me permet "
                 "de vérifier si un prix unitaire proposé par une entreprise est cohérent avec les coûts de "
                 "production réels, et d'argumenter cette analyse devant la Commission d'Appel d'Offres.",
                 Cm(1), Cm(27.2), Cm(19), Cm(2.3), font_size=7, color=GRIS_FONCE, align=PP_ALIGN.CENTER)
    add_page_number(slide, 26)
    print("  Page 26 : Complémentaires OK")

def page_27_competences(prs):
    """Page 27 - Compétences BTS MEC."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_rect(slide, BLANC)
    add_text_box(slide, "COMPÉTENCES BTS MEC ACQUISES", Cm(1), Cm(1), Cm(19), Cm(1.5),
                 font_size=22, color=ROUGE, bold=True, align=PP_ALIGN.CENTER)
    # Barres de progression
    competences = [
        ("Réaliser des métrés", 85, "Autonome", "7 marchés + extraction BIM"),
        ("Estimer un ouvrage", 80, "Autonome", "Estimations admin. + bordereaux"),
        ("Analyser des offres", 85, "Autonome", "Grilles notation, comparatifs, CAO"),
        ("Suivre un budget", 90, "Expert", "Décomptes, situations, 7 marchés"),
        ("Suivre un chantier", 95, "Expert", "5 ans terrain + 4,5 ans MOA"),
        ("Communiquer (C19)", 80, "Autonome", "Réunions, CR, reporting, CAO"),
        ("Rédiger pièces marchés", 80, "Autonome", "CPS, RC, BPDE, OS, PV"),
        ("Collaborer en BIM (T23)", 75, "Autonome", "IFC, convention BIM, LOD"),
        ("Modéliser en BIM", 85, "Autonome", "Revit, Navisworks, IFC"),
        ("Réglementation marchés", 60, "En progression", "Décret marocain + Code FR"),
    ]
    y = Cm(3.5)
    for comp, pct, level, detail in competences:
        add_text_box(slide, comp, Cm(1), y, Cm(5), Cm(0.6), font_size=9, color=NOIR, bold=True)
        add_text_box(slide, detail, Cm(1), y + Cm(0.6), Cm(5), Cm(0.5), font_size=6, color=GRIS)
        # Barre
        bar_left = Cm(7)
        bar_w = Cm(9.5)
        bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bar_left, y + Cm(0.15), bar_w, Cm(0.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = GRIS_CLAIR
        bg.line.fill.background()
        fill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bar_left, y + Cm(0.15), int(bar_w * pct / 100), Cm(0.5))
        fill.fill.solid()
        fill.fill.fore_color.rgb = ROUGE if pct >= 80 else ORANGE if pct >= 60 else GRIS
        fill.line.fill.background()
        add_text_box(slide, f"{level} ({pct}%)", Cm(17), y + Cm(0.05), Cm(3.5), Cm(0.6),
                     font_size=7, color=GRIS_FONCE)
        y += Cm(1.95)
    # Texte compétence distinctive
    add_text_box(slide, "COMPÉTENCE DISTINCTIVE : TRIPLE VISION MOA / EXÉCUTION / BIM",
                 Cm(1), Cm(23.5), Cm(19), Cm(0.8), font_size=10, color=ROUGE, bold=True, align=PP_ALIGN.CENTER)
    # 3 badges en bas
    badges = [
        ("MOA", "Maîtrise d'ouvrage\npublique (4,5 ans)"),
        ("EXÉCUTION", "Chantier gros\nœuvre (5 ans)"),
        ("BIM", "Modélisation &\nAutomatisation (8 mois)"),
    ]
    x = Cm(2)
    for badge, desc in badges:
        c = add_circle(slide, x, Cm(24.8), Cm(3.5), ROUGE)
        add_text_box(slide, badge, x, Cm(25.3), Cm(3.5), Cm(1),
                     font_size=11, color=BLANC, bold=True, align=PP_ALIGN.CENTER)
        add_text_box(slide, desc, x, Cm(26.3), Cm(3.5), Cm(1.5),
                     font_size=7, color=BLANC, align=PP_ALIGN.CENTER)
        x += Cm(6)
    add_page_number(slide, 27)
    print("  Page 27 : Compétences OK")

def page_28_comparaison(prs):
    """Page 28 - Comparaison Maroc / France."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_rect(slide, BLANC)
    add_text_box(slide, "ANALYSE COMPARATIVE", Cm(1), Cm(1), Cm(19), Cm(1),
                 font_size=14, color=GRIS)
    add_text_box(slide, "MAROC vs FRANCE", Cm(1), Cm(2), Cm(19), Cm(1.5),
                 font_size=28, color=ROUGE, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, "Mon parcours international me permet de comparer les pratiques de la construction et "
                 "des marchés publics entre les deux pays. Le système marocain est régi par un décret unique "
                 "tandis que le système français est plus complexe avec le Code de la commande publique.",
                 Cm(1), Cm(3.8), Cm(19), Cm(1.2), font_size=8, color=GRIS_FONCE, align=PP_ALIGN.CENTER)
    # En-têtes avec drapeaux (simulés)
    add_text_box(slide, "MAROC", Cm(4), Cm(5.5), Cm(6), Cm(1.5),
                 font_size=16, color=ROUGE, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, "FRANCE", Cm(11), Cm(5.5), Cm(6), Cm(1.5),
                 font_size=16, color=RGBColor(0x00, 0x2F, 0x87), bold=True, align=PP_ALIGN.CENTER)
    # Tableau comparatif
    rows = [
        ("Réglementation", "Décret n°2-12-349\ndu 20/03/2013", "Code de la commande\npublique"),
        ("Pièces du marché", "CPS + RC + BPDE\n+ Plans", "CCAP + CCTP + BPU/DQE\nou DPGF"),
        ("Normes", "Normes marocaines\nRPS 2000 (sismique)", "Eurocodes, DTU\nRE2020 (environnement)"),
        ("Suivi financier", "Attachements\ncontradictoires,\ndécomptes", "Situations de travaux\nmensuelles,\nrévision prix CCAP"),
        ("Passation", "Appel d'offres ouvert\nou restreint,\nmarché négocié", "Procédure formalisée\nou adaptée selon\nseuils européens"),
    ]
    y = Cm(8)
    for i, (critere, maroc, france) in enumerate(rows):
        bg_color = GRIS_CLAIR if i % 2 == 0 else BLANC
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(1), y, Cm(19), Cm(3))
        bg.fill.solid()
        bg.fill.fore_color.rgb = bg_color
        bg.line.fill.background()
        add_text_box(slide, critere, Cm(1.5), y + Cm(0.3), Cm(3.5), Cm(2.5),
                     font_size=9, color=ROUGE, bold=True)
        add_text_box(slide, maroc, Cm(5.5), y + Cm(0.3), Cm(5), Cm(2.5),
                     font_size=9, color=NOIR, align=PP_ALIGN.CENTER)
        add_text_box(slide, france, Cm(12), y + Cm(0.3), Cm(5), Cm(2.5),
                     font_size=9, color=NOIR, align=PP_ALIGN.CENTER)
        y += Cm(3.2)
    # Synthèse
    synth = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(2), Cm(24.5), Cm(17), Cm(3))
    synth.fill.solid()
    synth.fill.fore_color.rgb = GRIS_CLAIR
    synth.line.fill.background()
    add_text_box(slide, "Point commun : mêmes principes fondamentaux de la commande publique\nTransparence | Égalité de traitement | Mise en concurrence\nChoix de l'offre économiquement la plus avantageuse",
                 Cm(3), Cm(25), Cm(15), Cm(2), font_size=11, color=ROUGE,
                 bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, "Apport de cette double culture : capacité d'adaptation et compréhension globale des enjeux "
                 "de la commande publique. Dans les deux systèmes, le respect du secret professionnel, la "
                 "confidentialité des offres et des délibérations de la CAO, ainsi que les règles déontologiques "
                 "(impartialité, probité, transparence) sont des obligations impératives du technicien MEC.",
                 Cm(2), Cm(27.2), Cm(17), Cm(1.5), font_size=8, color=GRIS_FONCE, align=PP_ALIGN.CENTER)
    add_page_number(slide, 28)
    print("  Page 28 : Comparaison Maroc/France OK")

def page_29_projet_pro(prs):
    """Page 29 - Projet professionnel."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_rect(slide, BLANC)
    # Treillis islamique très discret en fond
    add_islamic_lattice(slide, 0, Cm(4), A4_W, Cm(16), Cm(1.8), ROUGE, opacity=4)
    add_text_box(slide, "MON PROJET PROFESSIONNEL", Cm(1), Cm(1), Cm(19), Cm(2),
                 font_size=26, color=ROUGE, bold=True, align=PP_ALIGN.CENTER)
    # Roadmap 3 horizons
    horizons = [
        ("COURT TERME", "2026", "Obtenir le BTS MEC\nLancer BIMCO (métrés, études\nde prix, suivi financier)\nDévelopper scripts et plugins\nRevit pour l'extraction\nautomatique de quantités"),
        ("MOYEN TERME", "2027-28", "Créer une gamme d'outils BIM\ndédiés au MEC :\n- Plugins extraction quantités\n- Chiffrage assisté maquette\n- Bases de prix connectées\n- Apps web suivi économique"),
        ("LONG TERME", "2029+", "BIMCO = cabinet d'ingénierie\nspécialisé BIM + Économie :\n- Prestations d'ingénierie\n  (métrés BIM, AMO, études)\n- Édition d'outils numériques\n  pour économistes"),
    ]
    x = Cm(0.5)
    for horizon, date, desc in horizons:
        # Carte avec ombre et fond blanc
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Cm(5), Cm(6.3), Cm(14))
        card.fill.solid()
        card.fill.fore_color.rgb = BLANC
        card.line.fill.background()
        add_shadow(card, alpha_pct=20)
        # Cercle en haut
        c_size = Cm(2.5)
        c = add_circle(slide, x + Cm(1.9), Cm(5.5), c_size, ROUGE)
        add_text_box(slide, date, x + Cm(1.9), Cm(5.9), c_size, Cm(1.5),
                     font_size=9, color=BLANC, bold=True, align=PP_ALIGN.CENTER)
        # Horizon
        add_text_box(slide, horizon, x, Cm(8.5), Cm(6.3), Cm(1.5),
                     font_size=13, color=ROUGE, bold=True, align=PP_ALIGN.CENTER)
        # Description
        add_text_box(slide, desc, x + Cm(0.3), Cm(10.5), Cm(5.7), Cm(7),
                     font_size=10, color=NOIR, align=PP_ALIGN.CENTER)
        # Flèche progressive
        if horizon != "LONG TERME":
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + Cm(6.3), Cm(11), Cm(0.8), Cm(1))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = ROUGE_CLAIR
            arrow.line.fill.background()
        x += Cm(6.8)
    # Citation en bas
    quote = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(2), Cm(22), Cm(17), Cm(4))
    quote.fill.solid()
    quote.fill.fore_color.rgb = ROUGE
    quote.line.fill.background()
    add_text_box(slide, "« Les outils numériques doivent être\nau service de l'économiste de la construction,\net non l'inverse. »",
                 Cm(3), Cm(22.3), Cm(15), Cm(2), font_size=16, color=BLANC,
                 bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, "Ma double compétence – économiste formé au terrain ET développeur\nmaîtrisant le BIM – constitue un avantage différenciant rare dans le secteur.",
                 Cm(3), Cm(24.5), Cm(15), Cm(1.5), font_size=10, color=RGBColor(0xDD, 0xDD, 0xDD),
                 align=PP_ALIGN.CENTER)
    add_page_number(slide, 29)
    print("  Page 29 : Projet professionnel OK")

def page_30_conclusion(prs):
    """Page 30 - Conclusion."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Photo pleine page avec overlay semi-transparent
    add_image_safe(slide, img("conclusion"), 0, 0, A4_W, A4_H)
    add_bg_rect(slide, ROUGE_FONCE, 0, 0, A4_W, A4_H, opacity=60)
    # --- Motifs islamiques ---
    add_islamic_border(slide, Cm(1.5), Cm(0.8), BLANC, opacity=12)
    add_islamic_star(slide, Cm(3), Cm(3), Cm(3.5), BLANC, opacity=10)
    add_islamic_star(slide, int(A4_W) - Cm(3), Cm(3), Cm(3.5), BLANC, opacity=10)
    add_islamic_star(slide, Cm(2), int(A4_H) - Cm(3), Cm(2.5), BLANC, opacity=8)
    add_islamic_star(slide, int(A4_W) - Cm(2), int(A4_H) - Cm(3), Cm(2.5), BLANC, opacity=8)
    add_islamic_border(slide, A4_H - Cm(2), Cm(0.8), BLANC, opacity=12)
    # Titre
    add_text_box(slide, "CONCLUSION", Cm(2), Cm(3), Cm(17), Cm(2),
                 font_size=36, color=BLANC, bold=True, align=PP_ALIGN.CENTER)
    # Points clés enrichis
    points = [
        ("8 ans de BTP entre Maroc et France",
         "De la formation initiale à la création de BIMCO, un parcours riche et complémentaire"),
        ("MOA publique + Exécution chantier + BIM",
         "Triple compétence rare : comprendre le maître d'ouvrage, le terrain et les outils numériques"),
        ("82,5 M DH de projets détaillés (7,4 M€)",
         "Mise à niveau de 4 communes (53,5 M DH) + Route Lehri-Kerrouchen 25 km (29 M DH)"),
        ("BIMCO : BIM + Économie de la Construction",
         "Des outils numériques au service de l'économiste : métrés, études de prix, suivi financier"),
    ]
    y = Cm(6)
    for title, desc in points:
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(4), y + Cm(0.3), Cm(0.8), Pt(2))
        line.fill.solid()
        line.fill.fore_color.rgb = ROUGE_CLAIR
        line.line.fill.background()
        add_text_box(slide, title, Cm(5.5), y, Cm(13), Cm(1),
                     font_size=14, color=BLANC, bold=True)
        add_text_box(slide, desc, Cm(5.5), y + Cm(1), Cm(13), Cm(1),
                     font_size=10, color=RGBColor(0xDD, 0xDD, 0xDD))
        y += Cm(3)
    # Grande citation
    add_text_box(slide, "« Le BTS MEC représente bien plus qu'un diplôme :\nc'est la validation officielle d'un parcours professionnel engagé\net le socle sur lequel je construirai des outils qui transformeront\nla pratique quotidienne de l'économie de la construction. »",
                 Cm(2), Cm(18.5), Cm(17), Cm(4.5), font_size=14, color=BLANC,
                 bold=True, align=PP_ALIGN.CENTER)
    # Logo et contact
    add_image_safe(slide, img("logo"), Cm(8), Cm(23.5), Cm(5))
    add_text_box(slide, "BAHAFID Mohamed | BIMCO | BTS MEC Session 2026 | Académie de Lyon",
                 Cm(1), Cm(26.5), Cm(19), Cm(1), font_size=10, color=BLANC, align=PP_ALIGN.CENTER)
    add_text_box(slide, "gestion.bimco-consulting.fr", Cm(1), Cm(27.5), Cm(19), Cm(1),
                 font_size=11, color=ROUGE_CLAIR, align=PP_ALIGN.CENTER)
    print("  Page 30 : Conclusion OK")


# =============================================================================
# MAIN
# =============================================================================
def main():
    print("=" * 60)
    print("GENERATION DU RAPPORT U62 - 30 PAGES")
    print("=" * 60)

    prs = Presentation()
    prs.slide_width = A4_W
    prs.slide_height = A4_H

    # Vérification des photos
    found = sum(1 for p in PHOTO.values() if os.path.exists(p))
    print(f"\nPhotos disponibles : {found}/{len(PHOTO)}")

    print("\nGeneration des 30 pages...")

    page_01_couverture(prs)
    page_02_fiche_candidat(prs)
    page_03_sommaire(prs)
    page_04_introduction(prs)
    page_05_section1(prs)
    page_06_parcours(prs)
    page_07_chiffres_cles(prs)
    page_08_conseil_regional(prs)
    page_09_organigramme(prs)
    page_10_bimco(prs)
    page_11_outils(prs)
    page_12_app_gestion(prs)
    page_13_section2(prs)
    page_14_projet1_fiche(prs)
    page_15_repartition_communes(prs)
    page_16_parties_techniques(prs)
    page_17_missions(prs)
    page_18_metres_ouaoumana(prs)
    page_19_analyse_offres(prs)
    page_20_suivi_kerrouchen(prs)
    page_21_difficultes_p1(prs)
    page_22_projet2_fiche(prs)
    page_23_metres_routiers(prs)
    page_24_budget_route(prs)
    page_25_difficultes_route(prs)
    page_26_complementaires(prs)
    page_27_competences(prs)
    page_28_comparaison(prs)
    page_29_projet_pro(prs)
    page_30_conclusion(prs)

    print(f"\nSauvegarde du fichier...")
    prs.save(OUTPUT)
    size = os.path.getsize(OUTPUT) / 1024
    print(f"\n{'=' * 60}")
    print(f"RAPPORT GENERE AVEC SUCCES !")
    print(f"  Fichier : {OUTPUT}")
    print(f"  Pages   : {len(prs.slides)}")
    print(f"  Taille  : {size:.0f} Ko")
    print(f"{'=' * 60}")

if __name__ == "__main__":
    main()
