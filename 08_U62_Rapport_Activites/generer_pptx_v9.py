"""
Génération du Rapport U62 v9 – 30 pages restructurées CPAR
BAHAFID Mohamed – BTS MEC Session 2026

CHANGEMENTS v9 vs v8 :
- 5 situations CPAR (Contexte/Problème/Action/Résultat) = cœur du rapport
- Pages fusionnées : 4+6+7, 8+9, 10+11+12, 15+16, 22+23
- Tableau synthèse activités → compétences
- Bilan réflexif + Protocole BIM
- Annexes : documents officiels + photos chantier

POUR MODIFIER LE TEXTE : éditez contenu_v2.py (pas ce fichier)
"""
import sys, os
sys.stdout.reconfigure(encoding="utf-8", errors="replace")

from contenu_v2 import (
    CANDIDAT, PAGE_01, PAGE_02, PAGE_03, PAGE_04, PAGE_05, PAGE_06, PAGE_07,
    PAGE_08, PAGE_09, PAGE_10, PAGE_15, PAGE_16, PAGE_17, PAGE_18,
    PAGE_20, PAGE_21, PAGE_22, PAGE_23, PAGE_24, PAGE_25, PAGE_26,
    PAGE_27, PAGE_28, PAGE_29, PAGE_30,
    SITUATION_1, SITUATION_2, SITUATION_3, SITUATION_4, SITUATION_5,
)

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData
from pptx.oxml.ns import qn
from lxml import etree
import math

# =============================================================================
# CONFIGURATION
# =============================================================================
BASE_DIR = r"D:\PREPA BTS MEC\08_U62_Rapport_Activites"
PHOTOS_DIR = os.path.join(BASE_DIR, "Documents_Entreprise", "Photos_Chantier")
ANNEXES_DIR = os.path.join(BASE_DIR, "Annexes", "Images")
OUTPUT = os.path.join(BASE_DIR, "Rapport_Redaction", "RAPPORT_U62_BAHAFID_30pages_v9.pptx")

# Couleurs
ROUGE = RGBColor(0xC0, 0x39, 0x2B)
ROUGE_CLAIR = RGBColor(0xE7, 0x4C, 0x3C)
ROUGE_FONCE = RGBColor(0x8B, 0x00, 0x00)
NOIR = RGBColor(0x33, 0x33, 0x33)
BLANC = RGBColor(0xFF, 0xFF, 0xFF)
GRIS_CLAIR = RGBColor(0xF5, 0xF5, 0xF5)
GRIS = RGBColor(0x99, 0x99, 0x99)
GRIS_FONCE = RGBColor(0x55, 0x55, 0x55)
VERT = RGBColor(0x27, 0xAE, 0x60)
VERT_CLAIR = RGBColor(0xEB, 0xFF, 0xEB)
ORANGE = RGBColor(0xF3, 0x9C, 0x12)
ORANGE_CLAIR = RGBColor(0xFF, 0xF3, 0xE0)
BLEU = RGBColor(0x29, 0x80, 0xB9)
GRIS_TRES_CLAIR = RGBColor(0xCC, 0xCC, 0xCC)

# Dimensions A4 portrait en EMU
A4_W = Cm(21)
A4_H = Cm(29.7)

# Photos sélectionnées
PHOTO = {
    "cover": os.path.join(PHOTOS_DIR, "tom-shamberger--KAoVUNv9W0-unsplash.jpg"),
    "bahafid": os.path.join(ANNEXES_DIR, "photo_bahafid.jpg"),
    "maroc": os.path.join(PHOTOS_DIR, "20180330_153333.jpg"),
    "france": os.path.join(PHOTOS_DIR, "IMG_20250718_155345[1].jpg"),
    "bimco_tech": os.path.join(PHOTOS_DIR, "d-c-NFVIq89r_8Q-unsplash.jpg"),
    "conseil": os.path.join(PHOTOS_DIR, "20181206_125854.jpg"),
    "chantier1": os.path.join(PHOTOS_DIR, "20180523_102152.jpg"),
    "chantier2": os.path.join(PHOTOS_DIR, "20180605_122447.jpg"),
    "chantier3": os.path.join(PHOTOS_DIR, "20180604_133353.jpg"),
    "chantier4": os.path.join(PHOTOS_DIR, "20180511_155634.jpg"),
    "chantier5": os.path.join(PHOTOS_DIR, "20180503_164033.jpg"),
    "chantier6": os.path.join(PHOTOS_DIR, "20180403_152746.jpg"),
    "chantier7": os.path.join(PHOTOS_DIR, "20180405_115900.jpg"),
    "projet1": os.path.join(PHOTOS_DIR, "20180403_123714.jpg"),
    "projet2": os.path.join(PHOTOS_DIR, "20181206_125636.jpg"),
    "route": os.path.join(PHOTOS_DIR, "20181206_130504.jpg"),
    "route2": os.path.join(PHOTOS_DIR, "20181206_141917.jpg"),
    "bilan": os.path.join(PHOTOS_DIR, "pexels-efeburakbaydar-35846752.jpg"),
    "conclusion": os.path.join(PHOTOS_DIR, "d-c-zEjnjuA_KBY-unsplash.jpg"),
    "gros_oeuvre": os.path.join(PHOTOS_DIR, "20180228_144045.jpg"),
    "vrd": os.path.join(PHOTOS_DIR, "20180330_154009.jpg"),
    "terrassement": os.path.join(PHOTOS_DIR, "20181206_133336.jpg"),
    "logo": os.path.join(ANNEXES_DIR, "logo_bimco.png"),
    "banniere": os.path.join(ANNEXES_DIR, "banniere_bimco.png"),
    "cao": os.path.join(ANNEXES_DIR, "convocation_cao.jpg"),
    "rejet": os.path.join(ANNEXES_DIR, "notification_rejet.jpg"),
    "cps": os.path.join(ANNEXES_DIR, "cps_marche46_signe.jpg"),
    "chantier_urban": os.path.join(PHOTOS_DIR, "20180503_105806.jpg"),
    "montagne": os.path.join(PHOTOS_DIR, "20181113_154208.jpg"),
    "france2": os.path.join(PHOTOS_DIR, "IMG_20250910_163813.jpg"),
}

# =============================================================================
# HELPERS (identiques à v8)
# =============================================================================
def img(name):
    p = PHOTO.get(name, "")
    return p if os.path.exists(p) else None

def set_transparency(shape, opacity_pct):
    fill_elem = shape._element.spPr.find(qn('a:solidFill'))
    if fill_elem is None:
        return
    srgb = fill_elem.find(qn('a:srgbClr'))
    if srgb is None:
        return
    for a in srgb.findall(qn('a:alpha')):
        srgb.remove(a)
    alpha = srgb.makeelement(qn('a:alpha'), {'val': str(int(opacity_pct * 1000))})
    srgb.append(alpha)

def add_shadow(shape, blur=50800, dist=38100, direction=5400000, alpha_pct=35):
    spPr = shape._element.spPr
    effectLst = spPr.find(qn('a:effectLst'))
    if effectLst is None:
        effectLst = spPr.makeelement(qn('a:effectLst'), {})
        spPr.append(effectLst)
    outerShdw = effectLst.makeelement(qn('a:outerShdw'), {
        'blurRad': str(blur), 'dist': str(dist),
        'dir': str(direction), 'algn': 'bl', 'rotWithShape': '0'
    })
    srgb = outerShdw.makeelement(qn('a:srgbClr'), {'val': '000000'})
    alpha_el = srgb.makeelement(qn('a:alpha'), {'val': str(int(alpha_pct * 1000))})
    srgb.append(alpha_el)
    outerShdw.append(srgb)
    effectLst.append(outerShdw)

def add_gradient_fill(shape, color1, color2, angle_deg=270):
    spPr = shape._element.spPr
    for tag in ['a:solidFill', 'a:noFill', 'a:gradFill']:
        for e in spPr.findall(qn(tag)):
            spPr.remove(e)
    gradFill = spPr.makeelement(qn('a:gradFill'), {'rotWithShape': '1'})
    gsLst = gradFill.makeelement(qn('a:gsLst'), {})
    for pos, color in [('0', color1), ('100000', color2)]:
        gs = gsLst.makeelement(qn('a:gs'), {'pos': pos})
        srgb = gs.makeelement(qn('a:srgbClr'), {
            'val': str(color)
        })
        gs.append(srgb)
        gsLst.append(gs)
    gradFill.append(gsLst)
    lin = gradFill.makeelement(qn('a:lin'), {
        'ang': str(angle_deg * 60000), 'scaled': '1'
    })
    gradFill.append(lin)
    spPr.insert(0, gradFill)

def add_bg_rect(slide, color, left=0, top=0, w=None, h=None, opacity=100):
    w = w or A4_W
    h = h or A4_H
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if opacity < 100:
        set_transparency(shape, opacity)
    return shape

def add_image_safe(slide, path, left, top, width=None, height=None):
    if path and os.path.exists(path):
        try:
            if width and height:
                return slide.shapes.add_picture(path, left, top, width, height)
            elif width:
                return slide.shapes.add_picture(path, left, top, width=width)
            elif height:
                return slide.shapes.add_picture(path, left, top, height=height)
            else:
                return slide.shapes.add_picture(path, left, top)
        except Exception as e:
            print(f"  Erreur image {path}: {e}")
    return None

def add_text_box(slide, text, left, top, width, height, font_size=12,
                 color=NOIR, bold=False, align=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_multi_text(slide, lines, left, top, width, height, font_size=11,
                   color=NOIR, spacing=1.2, bold_first=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(font_size * spacing * 0.3)
        if bold_first and i == 0:
            p.font.bold = True
    return txBox

def add_circle(slide, left, top, size, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_card(slide, left, top, w, h, title, text, icon_text="",
             bg_color=BLANC, title_color=ROUGE, text_color=NOIR, shadow=True):
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    rect.fill.solid()
    rect.fill.fore_color.rgb = bg_color
    rect.line.fill.background()
    if shadow:
        add_shadow(rect)
    if icon_text:
        add_text_box(slide, icon_text, left + Cm(0.3), top + Cm(0.2),
                     w - Cm(0.6), Cm(1.2), font_size=18, color=title_color,
                     bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, title, left + Cm(0.3), top + Cm(1.3),
                 w - Cm(0.6), Cm(1), font_size=9, color=title_color,
                 bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, text, left + Cm(0.3), top + Cm(2.1),
                 w - Cm(0.6), h - Cm(2.4), font_size=7, color=text_color,
                 align=PP_ALIGN.CENTER)

def add_bar_h(slide, left, top, max_w, h, pct, color, label="", value=""):
    bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, max_w, h)
    bg.fill.solid()
    bg.fill.fore_color.rgb = GRIS_CLAIR
    bg.line.fill.background()
    bar_w = int(max_w * pct / 100)
    if bar_w > 0:
        bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, bar_w, h)
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
    if label:
        add_text_box(slide, label, left - Cm(5), top, Cm(4.8), h,
                     font_size=8, color=NOIR, align=PP_ALIGN.RIGHT)
    if value:
        add_text_box(slide, value, left + max_w + Cm(0.2), top, Cm(2), h,
                     font_size=8, color=GRIS_FONCE)

def add_donut_with_center(slide, left, top, w, h, categories, values, colors,
                           center_text, center_sub=""):
    """Donut chart avec trou central, data labels et texte overlay."""
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series('Budget', tuple(values))
    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.DOUGHNUT, left, top, w, h, chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = False
    # Trou central 55% via XML
    for dc in chart._chartSpace.iter(qn('c:doughnutChart')):
        for hs in dc.findall(qn('c:holeSize')):
            dc.remove(hs)
        hs = etree.SubElement(dc, qn('c:holeSize'))
        hs.set('val', '55')
    # Couleurs des segments
    plot = chart.plots[0]
    for i, c in enumerate(colors):
        if i < len(values):
            plot.series[0].points[i].format.fill.solid()
            plot.series[0].points[i].format.fill.fore_color.rgb = c
    # Data labels (pourcentages)
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.show_percentage = True
    dl.show_value = False
    dl.show_category_name = False
    dl.font.size = Pt(9)
    dl.font.color.rgb = BLANC
    dl.font.bold = True
    # Texte overlay au centre du trou
    cx = left + int(w * 0.5) - Cm(2.5)
    cy = top + int(h * 0.45) - Cm(0.8)
    add_text_box(slide, center_text, cx, cy, Cm(5), Cm(1.2),
                 font_size=14, color=ROUGE, bold=True, align=PP_ALIGN.CENTER)
    if center_sub:
        add_text_box(slide, center_sub, cx, cy + Cm(1), Cm(5), Cm(0.6),
                     font_size=8, color=GRIS_FONCE, align=PP_ALIGN.CENTER)
    return chart_frame

def add_progress_bar(slide, left, top, width, height, pct, color, label=""):
    """Mini barre de progression pour tableaux de compétences."""
    bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    bg.line.fill.background()
    bar_w = int(width * pct / 100)
    if bar_w > 0:
        bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, bar_w, height)
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
    if label:
        add_text_box(slide, label, left + width + Cm(0.2), top - Cm(0.05),
                     Cm(2), height + Cm(0.1), font_size=8, color=color, bold=True)

def add_page_number(slide, num):
    add_islamic_border(slide, A4_H - Cm(1.5), Cm(0.35), ROUGE, opacity=25)
    add_text_box(slide, "BAHAFID Mohamed | U62 | BTS MEC 2026",
                 Cm(1), A4_H - Cm(1.2), Cm(12), Cm(0.8), font_size=7, color=GRIS)
    add_text_box(slide, str(num).zfill(2), A4_W - Cm(2), A4_H - Cm(1.2),
                 Cm(1.5), Cm(0.8), font_size=8, color=GRIS, align=PP_ALIGN.RIGHT)

def add_red_accent(slide, left=0, top=0, w=Cm(0.4), h=None):
    h = h or A4_H
    return add_bg_rect(slide, ROUGE, left, top, w, h)

def add_section_separator(slide, num, title, subtitle, photo_path):
    add_image_safe(slide, photo_path, 0, 0, A4_W, A4_H)
    overlay = add_bg_rect(slide, ROUGE_FONCE, 0, 0, A4_W, A4_H, opacity=70)
    add_islamic_border(slide, Cm(3.5), Cm(0.8), BLANC, opacity=12)
    add_islamic_star(slide, Cm(3), Cm(4), Cm(2.5), BLANC, opacity=10)
    add_islamic_star(slide, int(A4_W) - Cm(3), Cm(4), Cm(2.5), BLANC, opacity=10)
    add_islamic_border(slide, A4_H - Cm(3), Cm(0.8), BLANC, opacity=12)
    add_text_box(slide, num, Cm(2), Cm(5), Cm(17), Cm(10),
                 font_size=160, color=BLANC, bold=True, align=PP_ALIGN.CENTER)
    add_islamic_border(slide, Cm(15.8), Cm(0.5), ROUGE_CLAIR, opacity=30)
    add_text_box(slide, title, Cm(2), Cm(17), Cm(17), Cm(3),
                 font_size=28, color=BLANC, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, subtitle, Cm(2), Cm(20.5), Cm(17), Cm(2),
                 font_size=14, color=BLANC, align=PP_ALIGN.CENTER)

def add_fullpage_photo_overlay(slide, photo_path, overlay_color=ROUGE_FONCE, opacity=65):
    add_image_safe(slide, photo_path, 0, 0, A4_W, A4_H)
    add_bg_rect(slide, overlay_color, 0, 0, A4_W, A4_H, opacity=opacity)

# =============================================================================
# MOTIFS ISLAMIQUES
# =============================================================================
def add_islamic_star(slide, cx, cy, size, color=None, opacity=25):
    c = color or ROUGE
    r = size / 2
    inner_r = r * 0.38
    verts = []
    for i in range(16):
        angle = math.radians(i * 22.5 - 90)
        radius = r if i % 2 == 0 else inner_r
        verts.append((int(r + radius * math.cos(angle)),
                      int(r + radius * math.sin(angle))))
    try:
        builder = slide.shapes.build_freeform(verts[0][0], verts[0][1])
        for v in verts[1:]:
            builder.add_line_to(v[0], v[1])
        builder.add_line_to(verts[0][0], verts[0][1])
        shape = builder.convert_to_shape(int(cx - r), int(cy - r))
        shape.fill.solid()
        shape.fill.fore_color.rgb = c
        shape.line.fill.background()
        if opacity < 100:
            set_transparency(shape, opacity)
        return shape
    except Exception:
        d = slide.shapes.add_shape(MSO_SHAPE.DIAMOND,
                                   int(cx - r), int(cy - r), int(size), int(size))
        d.fill.solid()
        d.fill.fore_color.rgb = c
        d.line.fill.background()
        if opacity < 100:
            set_transparency(d, opacity)
        return d

def add_islamic_border(slide, y, band_h=None, color=None, opacity=18):
    c = color or ROUGE
    h = band_h or Cm(0.6)
    cell = int(h)
    x = 0
    while x < int(A4_W) + cell:
        d = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, x, int(y), cell, cell)
        d.fill.solid()
        d.fill.fore_color.rgb = c
        d.line.fill.background()
        set_transparency(d, opacity)
        x += cell

def add_islamic_corner(slide, corner, size=None, color=None, opacity=18):
    sz = size or Cm(3)
    c = color or ROUGE_CLAIR
    half = int(sz) // 2
    corners_pos = {
        'tl': (0, 0), 'tr': (int(A4_W), 0),
        'bl': (0, int(A4_H)), 'br': (int(A4_W), int(A4_H)),
    }
    cx, cy = corners_pos[corner]
    add_islamic_star(slide, cx, cy, int(sz), c, opacity)
    small = int(sz * 0.35)
    for dx, dy in [(half, 0), (0, half), (-half, 0), (0, -half)]:
        nx, ny = cx + dx, cy + dy
        d = slide.shapes.add_shape(MSO_SHAPE.DIAMOND,
                                   nx - small // 2, ny - small // 2, small, small)
        d.fill.solid()
        d.fill.fore_color.rgb = c
        d.line.fill.background()
        set_transparency(d, max(opacity - 5, 5))

def add_islamic_lattice(slide, left, top, width, height, cell_size=None,
                        color=None, opacity=10):
    c = color or ROUGE
    cs = cell_size or Cm(1.2)
    cell = int(cs)
    half = cell // 2
    y = int(top)
    row = 0
    while y < int(top + height):
        x = int(left) + (half if row % 2 else 0)
        while x < int(left + width):
            d = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, x, y, cell, cell)
            d.fill.solid()
            d.fill.fore_color.rgb = c
            d.line.fill.background()
            set_transparency(d, opacity)
            x += cell
        y += half
        row += 1


# =============================================================================
# NOUVEAU HELPER : PAGE SITUATION CPAR
# =============================================================================
def add_situation_page(prs, situation, page_num):
    """Génère une page CPAR (Contexte/Problème/Action/Résultat) pour une situation."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_rect(slide, BLANC)

    # --- Bandeau en-tête : SITUATION X ---
    bandeau = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, A4_W, Cm(3.8))
    bandeau.fill.solid()
    bandeau.fill.fore_color.rgb = ROUGE
    bandeau.line.fill.background()

    # Numéro de situation (grand)
    add_text_box(slide, str(situation["numero"]), Cm(0.5), Cm(0.2), Cm(3), Cm(3),
                 font_size=48, color=BLANC, bold=True, align=PP_ALIGN.CENTER)
    # Titre
    add_text_box(slide, f"SITUATION {situation['numero']}", Cm(3.5), Cm(0.3), Cm(13), Cm(1),
                 font_size=12, color=RGBColor(0xFF, 0xCC, 0xCC), bold=True)
    add_text_box(slide, situation["titre"], Cm(3.5), Cm(1.2), Cm(13), Cm(1.5),
                 font_size=16, color=BLANC, bold=True)

    # Chiffre clé en encadré à droite + ombre douce
    kpi_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Cm(16.5), Cm(0.4), Cm(4), Cm(2.8))
    kpi_bg.fill.solid()
    kpi_bg.fill.fore_color.rgb = BLANC
    kpi_bg.line.fill.background()
    set_transparency(kpi_bg, 95)
    add_shadow(kpi_bg, alpha_pct=8)
    add_text_box(slide, situation["chiffre_cle"], Cm(16.5), Cm(0.5), Cm(4), Cm(1.5),
                 font_size=16, color=ROUGE, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, situation["chiffre_label"], Cm(16.5), Cm(2), Cm(4), Cm(1),
                 font_size=8, color=GRIS_FONCE, align=PP_ALIGN.CENTER)

    # --- 4 blocs CPAR ---
    blocs = [
        ("C", "CONTEXTE", situation["contexte"], GRIS_CLAIR, GRIS_FONCE),
        ("P", "PROBLÈME", situation["probleme"], ORANGE_CLAIR, ORANGE),
        ("A", "ACTION", situation["action"], RGBColor(0xFF, 0xEB, 0xEB), ROUGE),
        ("R", "RÉSULTAT", situation["resultat"], VERT_CLAIR, VERT),
    ]

    y = Cm(4.3)
    bloc_h = Cm(5.0)

    for letter, label, text, bg_color, accent_color in blocs:
        # Fond du bloc
        bloc = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       Cm(1), y, Cm(19), bloc_h)
        bloc.fill.solid()
        bloc.fill.fore_color.rgb = bg_color
        bloc.line.fill.background()
        add_shadow(bloc, alpha_pct=8)

        # Barre latérale colorée
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(1), y, Cm(0.35), bloc_h)
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent_color
        bar.line.fill.background()

        # Lettre en cercle
        circle = add_circle(slide, Cm(1.8), y + Cm(0.5), Cm(1.3), accent_color)
        add_text_box(slide, letter, Cm(1.8), y + Cm(0.6), Cm(1.3), Cm(1.1),
                     font_size=14, color=BLANC, bold=True, align=PP_ALIGN.CENTER)

        # Label
        add_text_box(slide, label, Cm(3.5), y + Cm(0.5), Cm(5), Cm(1),
                     font_size=11, color=accent_color, bold=True)

        # Fine ligne séparatrice sous le label (accent color, 1pt)
        sep_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                           Cm(3.5), y + Cm(1.5), Cm(10), Pt(1))
        sep_line.fill.solid()
        sep_line.fill.fore_color.rgb = accent_color
        sep_line.line.fill.background()

        # Texte justifié
        add_text_box(slide, text, Cm(3.5), y + Cm(1.8), Cm(16), bloc_h - Cm(2.2),
                     font_size=10, color=NOIR, align=PP_ALIGN.JUSTIFY)

        y += bloc_h + Cm(0.35)

    # --- Badge compétence en bas ---
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Cm(1), A4_H - Cm(2.8), Cm(19), Cm(1.2))
    badge.fill.solid()
    badge.fill.fore_color.rgb = ROUGE_FONCE
    badge.line.fill.background()
    add_text_box(slide, situation['competence'].upper(),
                 Cm(1.5), A4_H - Cm(2.7), Cm(18), Cm(1),
                 font_size=11, color=BLANC, bold=True, align=PP_ALIGN.CENTER)

    add_page_number(slide, page_num)
    print(f"  Page {page_num} : Situation {situation['numero']} OK")


# =============================================================================
# PAGES
# =============================================================================

def page_01_couverture(prs):
    """Page 1 - Couverture."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_image_safe(slide, img("cover"), 0, 0, A4_W, A4_H)
    overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, A4_W, A4_H)
    overlay.line.fill.background()
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = RGBColor(0x2A, 0x00, 0x00)
    set_transparency(overlay, 65)
    add_islamic_star(slide, int(A4_W) - Cm(2), Cm(2), Cm(4), ROUGE_CLAIR, opacity=15)
    add_islamic_star(slide, int(A4_W) - Cm(0.5), Cm(5), Cm(1.8), ROUGE_CLAIR, opacity=12)
    add_islamic_star(slide, Cm(4), int(A4_H) - Cm(3), Cm(3), ROUGE_CLAIR, opacity=12)
    add_islamic_border(slide, A4_H - Cm(1.5), Cm(0.7), ROUGE_CLAIR, opacity=15)
    bande = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Cm(1.5), A4_H)
    bande.line.fill.background()
    bande.fill.solid()
    bande.fill.fore_color.rgb = ROUGE
    add_text_box(slide, "U62", Cm(0.1), Cm(12), Cm(1.3), Cm(4),
                 font_size=18, color=BLANC, bold=True, align=PP_ALIGN.CENTER)
    deco = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(3), Cm(2.5), Cm(5), Pt(3))
    deco.fill.solid()
    deco.fill.fore_color.rgb = ROUGE_CLAIR
    deco.line.fill.background()
    add_text_box(slide, PAGE_01["titre_1"], Cm(3), Cm(3.5), Cm(16), Cm(3),
                 font_size=48, color=BLANC, bold=True)
    add_text_box(slide, PAGE_01["titre_2"], Cm(3), Cm(6.5), Cm(16), Cm(3),
                 font_size=48, color=BLANC, bold=True)
    add_text_box(slide, PAGE_01["titre_3"], Cm(3), Cm(9.5), Cm(16), Cm(3),
                 font_size=36, color=ROUGE_CLAIR, bold=True)
    mec_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(3), Cm(13.5), Cm(13), Cm(2.5))
    mec_bg.fill.solid()
    mec_bg.fill.fore_color.rgb = ROUGE
    mec_bg.line.fill.background()
    set_transparency(mec_bg, 80)
    add_text_box(slide, PAGE_01["bts"], Cm(3.5), Cm(13.8), Cm(12), Cm(2),
                 font_size=16, color=BLANC, bold=True)
    add_text_box(slide, PAGE_01["session"], Cm(3), Cm(17), Cm(10), Cm(1.5),
                 font_size=20, color=ROUGE_CLAIR, bold=True)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(3), Cm(19), Cm(8), Pt(4))
    line.fill.solid()
    line.fill.fore_color.rgb = ROUGE_CLAIR
    line.line.fill.background()
    add_text_box(slide, CANDIDAT["nom"], Cm(3), Cm(20), Cm(15), Cm(2),
                 font_size=28, color=BLANC, bold=True)
    add_text_box(slide, PAGE_01["candidat_info"], Cm(3), Cm(22.5), Cm(15), Cm(1),
                 font_size=12, color=RGBColor(0xDD, 0xDD, 0xDD))
    logo_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(14.5), Cm(24.5), Cm(5.5), Cm(4))
    logo_bg.fill.solid()
    logo_bg.fill.fore_color.rgb = BLANC
    logo_bg.line.fill.background()
    set_transparency(logo_bg, 85)
    add_shadow(logo_bg, alpha_pct=20)
    add_image_safe(slide, img("logo"), Cm(15), Cm(25), Cm(4.5))
    print("  Page 1 : Couverture OK")


def page_02_fiche_candidat(prs):
    """Page 2 - Fiche d'identité du candidat."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_rect(slide, BLANC)
    add_islamic_border(slide, Cm(-0.1), Cm(0.5), ROUGE, opacity=35)
    add_islamic_border(slide, A4_H - Cm(0.4), Cm(0.5), ROUGE, opacity=35)
    add_text_box(slide, PAGE_02["titre"], Cm(1), Cm(1), Cm(19), Cm(1.5),
                 font_size=22, color=ROUGE, bold=True, align=PP_ALIGN.CENTER)
    circle = add_circle(slide, Cm(2), Cm(4), Cm(5), ROUGE)
    add_image_safe(slide, img("bahafid"), Cm(2.2), Cm(4.2), Cm(4.6), Cm(4.6))
    infos = PAGE_02["champs"]
    y = Cm(3.5)
    for label, value in infos:
        add_text_box(slide, label, Cm(9), y, Cm(5), Cm(0.6),
                     font_size=7, color=GRIS, bold=True)
        add_text_box(slide, value, Cm(9), y + Cm(0.4), Cm(10), Cm(0.8),
                     font_size=10, color=NOIR)
        y += Cm(1.7)
    add_image_safe(slide, img("banniere"), Cm(3), Cm(24), Cm(15))
    add_text_box(slide, PAGE_02["pied"],
                 Cm(1), Cm(27), Cm(19), Cm(1), font_size=9, color=ROUGE,
                 align=PP_ALIGN.CENTER)
    add_page_number(slide, 2)
    print("  Page 2 : Fiche candidat OK")


def page_03_sommaire(prs):
    """Page 3 - Sommaire mis à jour."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_rect(slide, BLANC)
    add_red_accent(slide, 0, 0, Cm(1), A4_H)
    add_text_box(slide, "SOMMAIRE", Cm(2), Cm(1.5), Cm(17), Cm(2),
                 font_size=30, color=ROUGE, bold=True)
    add_islamic_border(slide, Cm(3.5), Cm(0.4), ROUGE, opacity=12)
    add_islamic_star(slide, int(A4_W) - Cm(2.5), int(A4_H) - Cm(3.5), Cm(3),
                     ROUGE_CLAIR, opacity=10)
    sections = PAGE_03["sections"]
    y = Cm(4.5)
    for num, title, subtitle, page in sections:
        c = add_circle(slide, Cm(2.5), y, Cm(1.2), ROUGE)
        add_text_box(slide, num, Cm(2.5), y + Cm(0.15), Cm(1.2), Cm(1),
                     font_size=11, color=BLANC, bold=True, align=PP_ALIGN.CENTER)
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(4.2), y + Cm(0.55),
                                       Cm(12), Pt(1))
        line.fill.solid()
        line.fill.fore_color.rgb = GRIS_CLAIR
        line.line.fill.background()
        add_text_box(slide, title, Cm(4.2), y, Cm(12), Cm(1),
                     font_size=13, color=NOIR, bold=True)
        add_text_box(slide, subtitle, Cm(4.2), y + Cm(1), Cm(12), Cm(0.8),
                     font_size=8, color=GRIS_FONCE)
        add_text_box(slide, page, Cm(16.5), y + Cm(0.1), Cm(3), Cm(1),
                     font_size=10, color=GRIS)
        y += Cm(3)
    add_page_number(slide, 3)
    print("  Page 3 : Sommaire OK")


def page_04_intro_parcours(prs):
    """Page 4 - Introduction + Parcours (fusion p.4+6+7)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_rect(slide, BLANC)
    add_islamic_star(slide, int(A4_W) - Cm(2), Cm(2), Cm(2.5), ROUGE_CLAIR, opacity=8)

    add_text_box(slide, "INTRODUCTION ET PARCOURS", Cm(1), Cm(0.8), Cm(19), Cm(1.2),
                 font_size=24, color=ROUGE, bold=True, align=PP_ALIGN.CENTER)

    # Texte d'intro centré sur les compétences évaluées
    add_text_box(slide, PAGE_04["intro_texte"],
                 Cm(1.5), Cm(2.2), Cm(18), Cm(2), font_size=10, color=NOIR,
                 align=PP_ALIGN.CENTER)

    # Parcours horizontal compact en 5 cercles
    phases = PAGE_04["phases"]
    x = Cm(0.5)
    for i, (date, title, desc) in enumerate(phases):
        c = add_circle(slide, x, Cm(4.8), Cm(3), ROUGE if i < 4 else ROUGE_FONCE)
        add_text_box(slide, title.split(" – ")[0] if " – " in title else title.split("\n")[0],
                     x, Cm(5.2), Cm(3), Cm(1.5),
                     font_size=8, color=BLANC, bold=True, align=PP_ALIGN.CENTER)
        add_text_box(slide, date, x, Cm(8), Cm(3), Cm(0.6),
                     font_size=8, color=GRIS_FONCE, align=PP_ALIGN.CENTER)
        add_text_box(slide, desc, x - Cm(0.3), Cm(8.6), Cm(3.6), Cm(1.2),
                     font_size=8, color=GRIS_FONCE, align=PP_ALIGN.CENTER)
        if i < 4:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                            x + Cm(3), Cm(6.1), Cm(1), Pt(3))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = ROUGE_CLAIR
            arrow.line.fill.background()
        x += Cm(4)

    # 4 chiffres clés en cartes
    chiffres = PAGE_04["chiffres_cles"]
    x = Cm(1)
    for val, label in chiffres:
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Cm(10.5), Cm(4.5), Cm(3))
        card.fill.solid()
        card.fill.fore_color.rgb = GRIS_CLAIR
        card.line.fill.background()
        add_shadow(card, alpha_pct=12)
        add_text_box(slide, val, x, Cm(10.7), Cm(4.5), Cm(1.5),
                     font_size=20, color=ROUGE, bold=True, align=PP_ALIGN.CENTER)
        add_text_box(slide, label, x, Cm(12.2), Cm(4.5), Cm(1),
                     font_size=9, color=GRIS_FONCE, align=PP_ALIGN.CENTER)
        x += Cm(4.8)

    # Deux projets présentés
    add_text_box(slide, "Ce rapport présente 5 situations analysées issues de 2 projets majeurs :",
                 Cm(1.5), Cm(14.2), Cm(18), Cm(0.8), font_size=11, color=NOIR, bold=True)
    for i, projet in enumerate(PAGE_04["projets"]):
        add_circle(slide, Cm(2), Cm(15.2 + i * 1.3), Cm(0.3), ROUGE)
        add_text_box(slide, projet, Cm(3), Cm(15 + i * 1.3), Cm(16), Cm(1),
                     font_size=10, color=NOIR)

    # Citation
    quote_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       Cm(2), Cm(18), Cm(17), Cm(2))
    quote_bg.fill.solid()
    quote_bg.fill.fore_color.rgb = GRIS_CLAIR
    quote_bg.line.fill.background()
    add_shadow(quote_bg, alpha_pct=15)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(2), Cm(18), Cm(0.3), Cm(2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ROUGE
    bar.line.fill.background()
    add_text_box(slide, "« Cette double expérience, côté maîtrise d'ouvrage publique et côté\n"
                 "exécution, m'a offert une vision complète du cycle de vie d'un projet. »",
                 Cm(3), Cm(18.2), Cm(15), Cm(1.7), font_size=11, color=ROUGE,
                 bold=True, align=PP_ALIGN.CENTER)

    add_page_number(slide, 4)
    print("  Page 4 : Introduction + Parcours OK")


def page_05_separateur_cadre(prs):
    """Page 5 - Séparateur Cadre professionnel."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_separator(slide, PAGE_05["numero"], PAGE_05["titre"],
                          PAGE_05["sous_titre"], img("conseil"))
    print("  Page 5 : Séparateur Cadre professionnel OK")


def page_06_conseil_poste(prs):
    """Page 6 - Conseil Régional + Mon poste (fusion p.8+9)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_rect(slide, BLANC)
    add_islamic_star(slide, int(A4_W) - Cm(2.5), Cm(2), Cm(3), ROUGE_CLAIR, opacity=8)

    add_text_box(slide, "STRUCTURE D'ACCUEIL", Cm(1), Cm(0.5), Cm(19), Cm(0.8),
                 font_size=12, color=GRIS, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, "Conseil Régional de Béni Mellal-Khénifra",
                 Cm(1), Cm(1.3), Cm(19), Cm(1.5),
                 font_size=20, color=ROUGE, bold=True, align=PP_ALIGN.CENTER)

    # Facts compacts à gauche
    facts = PAGE_06["facts"]
    y = Cm(3.5)
    for title, desc in facts:
        add_text_box(slide, title, Cm(1), y, Cm(9), Cm(0.7),
                     font_size=11, color=ROUGE, bold=True)
        add_text_box(slide, desc, Cm(1), y + Cm(0.7), Cm(9), Cm(0.8),
                     font_size=9, color=NOIR)
        y += Cm(2)

    # Photo compacte
    add_image_safe(slide, img("conseil"), Cm(11), Cm(3.5), Cm(9), Cm(7))

    # Organigramme simplifié
    add_text_box(slide, "MON POSTE AU SEIN DE L'AGENCE", Cm(1), Cm(11.5), Cm(19), Cm(1),
                 font_size=14, color=ROUGE, bold=True, align=PP_ALIGN.CENTER)

    def add_org_box(text, x, y, w, h, is_me=False):
        rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        rect.fill.solid()
        rect.fill.fore_color.rgb = ROUGE if is_me else BLANC
        if not is_me:
            rect.line.color.rgb = GRIS
            rect.line.width = Pt(1)
        else:
            rect.line.fill.background()
        add_text_box(slide, text, x + Cm(0.2), y + Cm(0.1), w - Cm(0.4), h - Cm(0.2),
                     font_size=8, color=BLANC if is_me else NOIR, bold=True,
                     align=PP_ALIGN.CENTER)

    # Compact org chart
    add_org_box("Président", Cm(7.5), Cm(13), Cm(6), Cm(1.3))
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(10.35), Cm(14.3), Cm(0.15), Cm(0.8)).fill.solid()
    add_org_box("Directeur – M. DOGHMANI", Cm(6), Cm(15.1), Cm(9), Cm(1.3))
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(10.35), Cm(16.4), Cm(0.15), Cm(0.8)).fill.solid()

    services = PAGE_06["organigramme"]["services"]
    for i, svc in enumerate(services):
        x = Cm(1) + i * Cm(6.5)
        add_org_box(svc, x, Cm(17.8), Cm(6), Cm(1.5))

    # Horizontal line
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(4), Cm(17.2), Cm(13), Pt(1.5)).fill.solid()

    # My position
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(17.35), Cm(19.3), Cm(0.15), Cm(1)).fill.solid()
    add_org_box(PAGE_06["organigramme"]["mon_poste"], Cm(13.5), Cm(20.3), Cm(7), Cm(1.8), is_me=True)

    # Missions compactes
    add_text_box(slide, "Missions principales :", Cm(1), Cm(22.5), Cm(19), Cm(0.6),
                 font_size=9, color=ROUGE, bold=True)
    missions = PAGE_06["missions"]
    y = Cm(23.2)
    for mission in missions:
        add_text_box(slide, f"• {mission}", Cm(1.5), y, Cm(18), Cm(0.6),
                     font_size=8, color=GRIS_FONCE)
        y += Cm(0.9)

    add_page_number(slide, 6)
    print("  Page 6 : Conseil Régional + Poste OK")


def page_07_bimco_outils(prs):
    """Page 7 - BIMCO + Outils + App (fusion p.10+11+12)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_rect(slide, BLANC)

    # Logo + titre BIMCO
    add_image_safe(slide, img("logo"), Cm(1), Cm(0.5), Cm(4))
    add_text_box(slide, "BIMCO", Cm(5.5), Cm(0.5), Cm(8), Cm(1.5),
                 font_size=28, color=ROUGE, bold=True)
    add_text_box(slide, PAGE_07["slogan"], Cm(5.5), Cm(2), Cm(10), Cm(0.8),
                 font_size=11, color=GRIS_FONCE)

    # Infos compactes
    for i, info in enumerate(PAGE_07["infos"]):
        add_text_box(slide, info, Cm(5.5), Cm(2.8 + i * 0.7), Cm(14), Cm(0.6),
                     font_size=8, color=GRIS)

    # 5 domaines en ligne
    x = Cm(0.5)
    for dom, sub in PAGE_07["domaines"]:
        c = add_circle(slide, x, Cm(4.5), Cm(2.5), ROUGE)
        add_text_box(slide, dom, x, Cm(4.9), Cm(2.5), Cm(1.5),
                     font_size=8, color=BLANC, bold=True, align=PP_ALIGN.CENTER)
        add_text_box(slide, sub, x - Cm(0.2), Cm(7.2), Cm(2.9), Cm(0.8),
                     font_size=8, color=GRIS_FONCE, align=PP_ALIGN.CENTER)
        x += Cm(4)

    # Outils avec barres de progression (6 principaux)
    add_text_box(slide, "OUTILS MAÎTRISÉS", Cm(1), Cm(8.5), Cm(19), Cm(0.8),
                 font_size=14, color=ROUGE, bold=True)
    outils = PAGE_07["outils_principaux"]
    y = Cm(9.5)
    for name, pct in outils:
        add_text_box(slide, name, Cm(1), y, Cm(4), Cm(0.6),
                     font_size=10, color=NOIR, bold=True)
        bar_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                         Cm(5.5), y + Cm(0.1), Cm(12), Cm(0.4))
        bar_bg.fill.solid()
        bar_bg.fill.fore_color.rgb = GRIS_CLAIR
        bar_bg.line.fill.background()
        bar_w = int(Cm(12) * pct / 100)
        bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Cm(5.5), y + Cm(0.1), bar_w, Cm(0.4))
        bar.fill.solid()
        bar.fill.fore_color.rgb = ROUGE
        bar.line.fill.background()
        add_text_box(slide, f"{pct}%", Cm(18), y, Cm(2), Cm(0.6),
                     font_size=8, color=GRIS_FONCE)
        y += Cm(1.1)

    # App en encadré compact
    add_text_box(slide, "APPLICATION « GESTION CHANTIERS »", Cm(1), Cm(16.5), Cm(19), Cm(0.8),
                 font_size=14, color=ROUGE, bold=True)
    app_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Cm(1), Cm(17.5), Cm(19), Cm(4))
    app_bg.fill.solid()
    app_bg.fill.fore_color.rgb = GRIS_CLAIR
    app_bg.line.fill.background()
    add_text_box(slide, "[CAPTURES D'ÉCRAN DE L'APPLICATION À FOURNIR]",
                 Cm(2), Cm(18), Cm(17), Cm(2), font_size=12, color=GRIS,
                 align=PP_ALIGN.CENTER)
    add_text_box(slide, f"{PAGE_07['app']['url']} | {PAGE_07['app']['stack']}",
                 Cm(2), Cm(20.2), Cm(17), Cm(0.8), font_size=8, color=GRIS_FONCE,
                 align=PP_ALIGN.CENTER)

    # Protocole BIM compact
    bim_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Cm(1), Cm(22), Cm(19), Cm(2.5))
    bim_bg.fill.solid()
    bim_bg.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
    bim_bg.line.color.rgb = ROUGE
    bim_bg.line.width = Pt(1)
    add_text_box(slide, "PROTOCOLE BIM : Convention LOD 300 / LOI 3 | IFC 2x3 | Open BIM",
                 Cm(1.5), Cm(22.2), Cm(18), Cm(0.7),
                 font_size=9, color=ROUGE_CLAIR, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, "Workflow : Modélisation Revit → Export IFC → Extraction quantités → Chiffrage DPGF → Reporting",
                 Cm(1.5), Cm(23), Cm(18), Cm(1),
                 font_size=8, color=GRIS_TRES_CLAIR, align=PP_ALIGN.CENTER)

    add_page_number(slide, 7)
    print("  Page 7 : BIMCO + Outils + App OK")


def page_08_separateur_projet1(prs):
    """Page 8 - Séparateur Projet 1."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_separator(slide, PAGE_08["numero"], PAGE_08["titre"],
                          PAGE_08["sous_titre"], img("projet1"))
    print("  Page 8 : Séparateur Projet 1 OK")


def page_09_projet1_fiche(prs):
    """Page 9 - Projet 1 : Fiche d'identité."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_rect(slide, BLANC)
    add_text_box(slide, "PROJET 1", Cm(1), Cm(0.8), Cm(19), Cm(1),
                 font_size=14, color=GRIS, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, PAGE_09["titre"], Cm(1), Cm(2), Cm(19), Cm(3),
                 font_size=26, color=ROUGE, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, "Ce projet s'inscrit dans le Programme de mise à niveau des centres "
                 "émergents, lancé par le Conseil Régional pour moderniser les infrastructures "
                 "urbaines de la province de Khénifra. 4 communes, 8 corps d'état, lot unique.",
                 Cm(0.5), Cm(5), Cm(20), Cm(1.5), font_size=8, color=GRIS_FONCE,
                 align=PP_ALIGN.CENTER)
    add_image_safe(slide, img("chantier_urban"), Cm(0.5), Cm(7), Cm(10), Cm(6))
    fiche = PAGE_09["fiche"]
    y = Cm(7.5)
    for label, value in fiche:
        add_text_box(slide, label, Cm(11.5), y, Cm(8), Cm(0.7),
                     font_size=9, color=ROUGE, bold=True)
        add_text_box(slide, value, Cm(11.5), y + Cm(0.7), Cm(8), Cm(1.5),
                     font_size=10, color=NOIR)
        y += Cm(2.5)
    big_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Cm(2), Cm(22), Cm(17), Cm(4))
    big_bg.fill.solid()
    big_bg.fill.fore_color.rgb = ROUGE
    big_bg.line.fill.background()
    add_text_box(slide, PAGE_09["montant"], Cm(2), Cm(22.3), Cm(17), Cm(2),
                 font_size=36, color=BLANC, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, PAGE_09["detail_montant"], Cm(2), Cm(24.3), Cm(17), Cm(0.8),
                 font_size=11, color=BLANC, align=PP_ALIGN.CENTER)
    add_text_box(slide, PAGE_09["pied"], Cm(2), Cm(25.2), Cm(17), Cm(0.8),
                 font_size=9, color=RGBColor(0xDD, 0xDD, 0xDD), align=PP_ALIGN.CENTER)
    add_page_number(slide, 9)
    print("  Page 9 : Projet 1 fiche OK")


def page_10_budget_corps(prs):
    """Page 10 - Budget + 8 corps d'état — donut + barres + ombres."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_rect(slide, BLANC)
    add_text_box(slide, "BUDGET ET 8 CORPS D'ÉTAT", Cm(1), Cm(0.8), Cm(19), Cm(1),
                 font_size=20, color=ROUGE, bold=True, align=PP_ALIGN.CENTER)

    # Donut chart — 4 communes
    communes = PAGE_10["communes"]
    categories = [c[0] for c in communes]
    values = [float(c[1].replace(",", ".").split()[0]) for c in communes]
    colors_donut = [
        RGBColor(0xE7, 0x4C, 0x3C),  # corail
        RGBColor(0xC0, 0x39, 0x2B),  # rouge
        RGBColor(0xE8, 0x8D, 0x72),  # saumon
        RGBColor(0xA0, 0x8B, 0x7B),  # gris chaud
    ]
    add_donut_with_center(slide, Cm(0.5), Cm(2.5), Cm(9), Cm(9),
                          categories, values, colors_donut,
                          "44,6 M DH HT", "Budget total")

    # 4 barres horizontales proportionnelles à droite
    y = Cm(3)
    for i, (nom, montant, pct) in enumerate(communes):
        pct_val = float(pct.replace(",", ".").replace("%", ""))
        add_bar_h(slide, Cm(12), y, Cm(7), Cm(0.6), pct_val, colors_donut[i],
                  label=nom, value=f"{montant} ({pct})")
        y += Cm(2)

    # Callout — coins arrondis + ombre douce
    callout = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Cm(10.5), Cm(9.5), Cm(10), Cm(2))
    callout.fill.solid()
    callout.fill.fore_color.rgb = ROUGE
    callout.line.fill.background()
    add_shadow(callout, alpha_pct=12)
    add_text_box(slide, "68% du budget sur 2 communes\nOuaoumana + Sebt Ait Rahou",
                 Cm(10.8), Cm(9.7), Cm(9.5), Cm(1.8),
                 font_size=10, color=BLANC, bold=True, align=PP_ALIGN.CENTER)

    # 8 corps d'état en grille compacte (4×2) — ombre douce
    add_text_box(slide, "8 CORPS D'ÉTAT DU MARCHÉ", Cm(1), Cm(12.5), Cm(19), Cm(0.8),
                 font_size=14, color=ROUGE, bold=True, align=PP_ALIGN.CENTER)
    parties = PAGE_10["parties"]
    for i, (num, title, desc) in enumerate(parties):
        col = i % 4
        row = i // 4
        x = Cm(0.5) + col * Cm(5)
        y = Cm(14) + row * Cm(6.5)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Cm(4.7), Cm(5.5))
        card.fill.solid()
        nuance = max(0, 0xC0 - i * 0x10)
        card.fill.fore_color.rgb = RGBColor(nuance,
                                             0x39 if nuance > 0x80 else 0x20,
                                             0x2B if nuance > 0x80 else 0x18)
        card.line.fill.background()
        add_shadow(card, alpha_pct=8)
        add_text_box(slide, num, x, y + Cm(0.5), Cm(4.7), Cm(1.5),
                     font_size=22, color=BLANC, bold=True, align=PP_ALIGN.CENTER)
        add_text_box(slide, title, x, y + Cm(2.2), Cm(4.7), Cm(1),
                     font_size=10, color=BLANC, bold=True, align=PP_ALIGN.CENTER)
        add_text_box(slide, desc, x, y + Cm(3.5), Cm(4.7), Cm(1.5),
                     font_size=8, color=BLANC, align=PP_ALIGN.CENTER)

    add_page_number(slide, 10)
    print("  Page 10 : Budget donut + 8 corps d'état OK")


# Pages 11-14 : Les 4 situations CPAR du Projet 1
def page_11_situation_1(prs):
    add_situation_page(prs, SITUATION_1, 11)

def page_12_situation_2(prs):
    add_situation_page(prs, SITUATION_2, 12)

def page_13_situation_3(prs):
    add_situation_page(prs, SITUATION_3, 13)

def page_14_situation_4(prs):
    add_situation_page(prs, SITUATION_4, 14)


def page_15_difficultes_p1(prs):
    """Page 15 - Difficultés et solutions Projet 1."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_rect(slide, BLANC)
    add_text_box(slide, PAGE_15["titre"], Cm(1), Cm(1), Cm(19), Cm(1.5),
                 font_size=22, color=ROUGE, bold=True, align=PP_ALIGN.CENTER)
    defis = PAGE_15["defis"]
    y = Cm(4)
    for defi, probleme, solution in defis:
        pb = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(1), y, Cm(8), Cm(4.5))
        pb.fill.solid()
        pb.fill.fore_color.rgb = RGBColor(0xFF, 0xEB, 0xEB)
        pb.line.fill.background()
        add_shadow(pb, alpha_pct=12)
        add_text_box(slide, defi, Cm(1.5), y + Cm(0.3), Cm(7), Cm(1),
                     font_size=11, color=ROUGE, bold=True)
        add_text_box(slide, probleme, Cm(1.5), y + Cm(1.5), Cm(7), Cm(2.5),
                     font_size=10, color=NOIR)
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Cm(9.3), y + Cm(1.5),
                                        Cm(1.5), Cm(1))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = VERT
        arrow.line.fill.background()
        sol = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(11.2), y, Cm(8.5), Cm(4.5))
        sol.fill.solid()
        sol.fill.fore_color.rgb = VERT_CLAIR
        sol.line.fill.background()
        add_shadow(sol, alpha_pct=12)
        add_text_box(slide, "SOLUTION", Cm(11.7), y + Cm(0.3), Cm(7.5), Cm(1),
                     font_size=9, color=VERT, bold=True)
        add_text_box(slide, solution, Cm(11.7), y + Cm(1.5), Cm(7.5), Cm(2.5),
                     font_size=10, color=NOIR)
        y += Cm(5.3)
    add_text_box(slide, "Enseignements : le suivi technique exige présence terrain, communication "
                 "transparente et anticipation des risques financiers.",
                 Cm(1), Cm(26), Cm(19), Cm(1.5), font_size=8, color=ROUGE,
                 align=PP_ALIGN.CENTER)
    add_page_number(slide, 15)
    print("  Page 15 : Difficultés P1 OK")


def page_16_separateur_projet2(prs):
    """Page 16 - Séparateur Projet 2."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_separator(slide, PAGE_16["numero"], PAGE_16["titre"],
                          PAGE_16["sous_titre"], img("montagne"))
    print("  Page 16 : Séparateur Projet 2 OK")


def page_17_projet2_metres(prs):
    """Page 17 - Fiche P2 + Métrés routiers (fusion p.22+23)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Photo pleine page
    add_image_safe(slide, img("route"), 0, 0, A4_W, A4_H)
    add_bg_rect(slide, ROUGE_FONCE, 0, 0, A4_W, A4_H, opacity=60)

    add_text_box(slide, "PROJET 2", Cm(2), Cm(2), Cm(17), Cm(1.5),
                 font_size=16, color=BLANC)
    add_text_box(slide, PAGE_17["titre"], Cm(2), Cm(3.5), Cm(17), Cm(2.5),
                 font_size=32, color=BLANC, bold=True)

    # Montant
    add_text_box(slide, PAGE_17["montant"], Cm(2), Cm(7), Cm(8), Cm(2),
                 font_size=40, color=ROUGE_CLAIR, bold=True)

    # Fiche compacte
    fiche_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       Cm(11), Cm(6), Cm(9), Cm(6))
    fiche_bg.fill.solid()
    fiche_bg.fill.fore_color.rgb = BLANC
    fiche_bg.line.fill.background()
    set_transparency(fiche_bg, 90)
    y = Cm(6.3)
    for label, value in PAGE_17["fiche"]:
        add_text_box(slide, label, Cm(11.5), y, Cm(3), Cm(0.6),
                     font_size=8, color=ROUGE, bold=True)
        add_text_box(slide, value, Cm(14.5), y, Cm(5), Cm(0.6),
                     font_size=9, color=NOIR)
        y += Cm(1.3)

    # Métrés en grille
    add_text_box(slide, "PRINCIPAUX MÉTRÉS", Cm(2), Cm(13), Cm(17), Cm(1),
                 font_size=18, color=BLANC, bold=True, align=PP_ALIGN.CENTER)

    metres = PAGE_17["metres"]
    for i, (val, label) in enumerate(metres):
        col = i % 3
        row = i // 3
        x = Cm(1.5) + col * Cm(6.5)
        y = Cm(14.5) + row * Cm(5)
        # Fond semi-transparent
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       x, y, Cm(5.5), Cm(4))
        card.fill.solid()
        card.fill.fore_color.rgb = BLANC
        card.line.fill.background()
        set_transparency(card, 90)
        add_text_box(slide, val, x, y + Cm(0.3), Cm(5.5), Cm(1.8),
                     font_size=20, color=ROUGE, bold=True, align=PP_ALIGN.CENTER)
        add_text_box(slide, label, x, y + Cm(2.3), Cm(5.5), Cm(1),
                     font_size=10, color=NOIR, align=PP_ALIGN.CENTER)

    # CPS image
    add_image_safe(slide, img("cps"), Cm(15), Cm(24), Cm(5))
    add_text_box(slide, "CPS signé", Cm(15), Cm(27.5), Cm(5), Cm(0.8),
                 font_size=7, color=BLANC, align=PP_ALIGN.CENTER)

    add_page_number(slide, 17)
    print("  Page 17 : Fiche P2 + Métrés OK")


def page_18_budget_route(prs):
    """Page 18 - Budget route — donut + barres proportionnelles."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_rect(slide, BLANC)
    add_text_box(slide, PAGE_18["titre"], Cm(1), Cm(1), Cm(19), Cm(1.5),
                 font_size=20, color=ROUGE, bold=True, align=PP_ALIGN.CENTER)

    # Donut chart — 6 postes budgétaires
    items = PAGE_18["items"]
    categories = [it[0] for it in items]
    values = [float(it[1].replace(",", ".").split()[0]) for it in items]
    colors_donut = [
        RGBColor(0xE7, 0x4C, 0x3C),  # corail
        RGBColor(0xC0, 0x39, 0x2B),  # rouge
        RGBColor(0x8B, 0x00, 0x00),  # rouge foncé
        RGBColor(0xE8, 0x8D, 0x72),  # saumon
        RGBColor(0xA0, 0x8B, 0x7B),  # gris chaud
        RGBColor(0x66, 0x66, 0x66),  # gris moyen
    ]
    add_donut_with_center(slide, Cm(0.5), Cm(3), Cm(9.5), Cm(11),
                          categories, values, colors_donut,
                          "24,2 M DH HT", "Budget total")

    # 6 barres horizontales proportionnelles à droite
    y = Cm(3.5)
    for i, (nom, montant, pct) in enumerate(items):
        pct_val = float(pct.replace(",", ".").replace("%", ""))
        add_bar_h(slide, Cm(12.5), y, Cm(6.5), Cm(0.6), pct_val, colors_donut[i],
                  label=nom, value=f"{montant} ({pct})")
        y += Cm(2)

    # Callout — coins arrondis + ombre douce
    callout = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Cm(2), Cm(18), Cm(17), Cm(2.5))
    callout.fill.solid()
    callout.fill.fore_color.rgb = ROUGE
    callout.line.fill.background()
    add_shadow(callout, alpha_pct=12)
    add_text_box(slide, PAGE_18["callout"],
                 Cm(3), Cm(18.2), Cm(15), Cm(2.2),
                 font_size=14, color=BLANC, bold=True, align=PP_ALIGN.CENTER)

    # CPS image en bas
    add_image_safe(slide, img("cps"), Cm(14), Cm(21.5), Cm(5.5))
    add_text_box(slide, "CPS signé – Marché n°46", Cm(14), Cm(25), Cm(5.5), Cm(0.8),
                 font_size=7, color=GRIS, align=PP_ALIGN.CENTER)

    add_page_number(slide, 18)
    print("  Page 18 : Budget route donut OK")


def page_19_situation_5(prs):
    """Page 19 - Situation 5 : Cubatures."""
    add_situation_page(prs, SITUATION_5, 19)


def page_20_defis_route(prs):
    """Page 20 - Défis chantier routier — layout pb/sol vertical."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_image_safe(slide, img("montagne"), 0, 0, A4_W, A4_H)
    add_bg_rect(slide, ROUGE_FONCE, 0, 0, A4_W, A4_H, opacity=55)
    add_text_box(slide, PAGE_20["titre"], Cm(2), Cm(1), Cm(17), Cm(2.5),
                 font_size=22, color=BLANC, bold=True, align=PP_ALIGN.CENTER)
    defis = PAGE_20["defis"]
    y = Cm(4.5)
    for i, (title, probleme, solution) in enumerate(defis):
        pb = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(1), y, Cm(8.5), Cm(5))
        pb.fill.solid()
        pb.fill.fore_color.rgb = BLANC
        pb.line.fill.background()
        set_transparency(pb, 90)
        add_shadow(pb, alpha_pct=20)
        add_text_box(slide, title, Cm(1.5), y + Cm(0.2), Cm(7.5), Cm(1),
                     font_size=11, color=ROUGE, bold=True)
        add_text_box(slide, probleme[:180], Cm(1.5), y + Cm(1.3), Cm(7.5), Cm(3.3),
                     font_size=8.5, color=NOIR)

        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Cm(9.8), y + Cm(1.8),
                                        Cm(1.2), Cm(0.8))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = VERT
        arrow.line.fill.background()

        sol = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(11.2), y, Cm(8.5), Cm(5))
        sol.fill.solid()
        sol.fill.fore_color.rgb = VERT_CLAIR
        sol.line.fill.background()
        set_transparency(sol, 90)
        add_shadow(sol, alpha_pct=20)
        add_text_box(slide, "SOLUTION", Cm(11.7), y + Cm(0.2), Cm(7.5), Cm(1),
                     font_size=10, color=VERT, bold=True)
        add_text_box(slide, solution[:200], Cm(11.7), y + Cm(1.3), Cm(7.5), Cm(3.3),
                     font_size=8.5, color=NOIR)
        y += Cm(5.5)
    add_page_number(slide, 20)
    print("  Page 20 : Défis route OK")


def page_21_complementaires(prs):
    """Page 21 - Activités complémentaires."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_rect(slide, BLANC)
    add_text_box(slide, PAGE_21["titre"], Cm(1), Cm(1), Cm(19), Cm(1),
                 font_size=22, color=ROUGE, bold=True, align=PP_ALIGN.CENTER)

    add_text_box(slide, PAGE_21["maroc_titre"], Cm(1), Cm(3), Cm(19), Cm(1),
                 font_size=14, color=ROUGE_FONCE, bold=True)
    y = Cm(4.2)
    for ref, nom, typ in PAGE_21["marches"]:
        add_circle(slide, Cm(2), y + Cm(0.15), Cm(0.3), ROUGE)
        add_text_box(slide, ref, Cm(3), y, Cm(4), Cm(0.7),
                     font_size=9, color=ROUGE, bold=True)
        add_text_box(slide, nom, Cm(7), y, Cm(12), Cm(0.7), font_size=10, color=NOIR)
        add_text_box(slide, typ, Cm(7), y + Cm(0.7), Cm(12), Cm(0.5),
                     font_size=8, color=GRIS)
        y += Cm(1.5)

    sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(2), Cm(12), Cm(17), Pt(2))
    sep.fill.solid()
    sep.fill.fore_color.rgb = ROUGE
    sep.line.fill.background()

    add_text_box(slide, PAGE_21["france_titre"], Cm(1), Cm(13), Cm(19), Cm(1),
                 font_size=14, color=ROUGE_FONCE, bold=True)
    y = Cm(14.5)
    for date, poste, detail in PAGE_21["france"]:
        add_text_box(slide, date, Cm(2), y, Cm(3), Cm(0.8),
                     font_size=11, color=ROUGE, bold=True)
        add_text_box(slide, poste, Cm(5.5), y, Cm(14), Cm(0.8),
                     font_size=11, color=NOIR, bold=True)
        add_text_box(slide, detail, Cm(5.5), y + Cm(0.9), Cm(14), Cm(1.5),
                     font_size=9, color=GRIS_FONCE)
        y += Cm(3)

    add_image_safe(slide, img("gros_oeuvre"), Cm(2), Cm(21), Cm(17), Cm(6))
    add_page_number(slide, 21)
    print("  Page 21 : Complémentaires OK")


def page_22_separateur_bilan(prs):
    """Page 22 - Séparateur Bilan et Analyse."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_separator(slide, PAGE_22["numero"], PAGE_22["titre"],
                          PAGE_22["sous_titre"], img("bilan"))
    print("  Page 22 : Séparateur Bilan OK")


def page_23_synthese_competences(prs):
    """Page 23 - Tableau synthèse — jauges visuelles + lignes fines."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_rect(slide, BLANC)
    add_text_box(slide, PAGE_23["titre"], Cm(0.5), Cm(0.8), Cm(20), Cm(1.2),
                 font_size=18, color=ROUGE, bold=True, align=PP_ALIGN.CENTER)

    # En-tête du tableau
    headers = ["Activité réalisée", "Sous-compétence", "Situation", "Niveau"]
    col_x = [Cm(0.5), Cm(7.5), Cm(13), Cm(17)]
    col_w = [Cm(7), Cm(5.5), Cm(4), Cm(3.5)]

    y = Cm(3)
    # Header row + ombre douce
    header_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(0.5), y, Cm(20), Cm(1.2))
    header_bg.fill.solid()
    header_bg.fill.fore_color.rgb = ROUGE
    header_bg.line.fill.background()
    add_shadow(header_bg, alpha_pct=8)
    for i, h in enumerate(headers):
        add_text_box(slide, h, col_x[i], y + Cm(0.1), col_w[i], Cm(1),
                     font_size=10, color=BLANC, bold=True, align=PP_ALIGN.CENTER)

    # Data rows
    y = Cm(4.2)
    tableau = PAGE_23["tableau"]
    for idx, (activite, sous_comp, situation, niveau) in enumerate(tableau):
        row_h = Cm(2.5)
        # Alternating background
        if idx % 2 == 0:
            row_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                             Cm(0.5), y, Cm(20), row_h)
            row_bg.fill.solid()
            row_bg.fill.fore_color.rgb = GRIS_CLAIR
            row_bg.line.fill.background()

        # Ligne fine gris clair entre les rangées
        line_sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                           Cm(0.5), y + row_h - Pt(1), Cm(20), Pt(1))
        line_sep.fill.solid()
        line_sep.fill.fore_color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
        line_sep.line.fill.background()

        # Colonnes texte (activité, sous-compétence, situation)
        for i, val in enumerate([activite, sous_comp, situation]):
            color = NOIR
            bold = False
            if i == 2:
                color = ROUGE
                bold = True
            add_text_box(slide, val, col_x[i], y + Cm(0.3), col_w[i], row_h - Cm(0.6),
                         font_size=9, color=color, bold=bold,
                         align=PP_ALIGN.CENTER if i > 0 else PP_ALIGN.LEFT)

        # Colonne Niveau : barre de progression + label
        bar_color = VERT if niveau == "Expert" else BLEU
        bar_pct = 100 if niveau == "Expert" else 75
        add_progress_bar(slide, col_x[3], y + Cm(0.8), Cm(2), Cm(0.35),
                         bar_pct, bar_color, label=niveau)

        y += row_h

    # Légende
    add_text_box(slide, PAGE_23["legende"],
                 Cm(1), A4_H - Cm(3), Cm(19), Cm(1),
                 font_size=9, color=GRIS_FONCE, align=PP_ALIGN.CENTER)

    add_page_number(slide, 23)
    print("  Page 23 : Tableau synthèse + jauges OK")


def page_24_comparaison(prs):
    """Page 24 - Comparaison Maroc/France — cartes blanches + bandeaux colorés."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_rect(slide, BLANC)
    add_text_box(slide, "ANALYSE COMPARATIVE", Cm(1), Cm(0.5), Cm(19), Cm(1),
                 font_size=12, color=GRIS, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, "MAROC vs FRANCE", Cm(1), Cm(1.5), Cm(19), Cm(2),
                 font_size=28, color=ROUGE, bold=True, align=PP_ALIGN.CENTER)

    # Tableau comparatif — cartes blanches avec bandeaux
    comp = PAGE_24["comparaison"]
    y = Cm(4)
    for theme, maroc, france in comp:
        # Colonne thème : cercle rouge + initiale
        initial = theme[0]
        add_circle(slide, Cm(1.5), y + Cm(0.8), Cm(1.5), ROUGE)
        add_text_box(slide, initial, Cm(1.5), y + Cm(0.9), Cm(1.5), Cm(1.3),
                     font_size=14, color=BLANC, bold=True, align=PP_ALIGN.CENTER)
        add_text_box(slide, theme, Cm(3.2), y + Cm(0.8), Cm(3), Cm(2),
                     font_size=10, color=ROUGE, bold=True)

        # Carte Maroc : blanc + bandeau ROUGE
        maroc_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                             Cm(6.5), y, Cm(6.5), Cm(3.5))
        maroc_card.fill.solid()
        maroc_card.fill.fore_color.rgb = BLANC
        maroc_card.line.fill.background()
        add_shadow(maroc_card, alpha_pct=8)
        maroc_banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                               Cm(6.5), y, Cm(6.5), Cm(0.8))
        maroc_banner.fill.solid()
        maroc_banner.fill.fore_color.rgb = ROUGE
        maroc_banner.line.fill.background()
        add_text_box(slide, "MAROC", Cm(6.5), y + Cm(0.05), Cm(6.5), Cm(0.7),
                     font_size=9, color=BLANC, bold=True, align=PP_ALIGN.CENTER)
        add_text_box(slide, maroc, Cm(6.7), y + Cm(1), Cm(6.1), Cm(2.3),
                     font_size=9, color=NOIR, align=PP_ALIGN.CENTER)

        # Carte France : blanc + bandeau BLEU
        france_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                              Cm(13.5), y, Cm(7), Cm(3.5))
        france_card.fill.solid()
        france_card.fill.fore_color.rgb = BLANC
        france_card.line.fill.background()
        add_shadow(france_card, alpha_pct=8)
        france_banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                                Cm(13.5), y, Cm(7), Cm(0.8))
        france_banner.fill.solid()
        france_banner.fill.fore_color.rgb = BLEU
        france_banner.line.fill.background()
        add_text_box(slide, "FRANCE", Cm(13.5), y + Cm(0.05), Cm(7), Cm(0.7),
                     font_size=9, color=BLANC, bold=True, align=PP_ALIGN.CENTER)
        add_text_box(slide, france, Cm(13.7), y + Cm(1), Cm(6.6), Cm(2.3),
                     font_size=9, color=NOIR, align=PP_ALIGN.CENTER)

        y += Cm(3.8)

    # Synthèse — dégradé rouge → rouge foncé
    synth_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       Cm(2), Cm(23.5), Cm(17), Cm(2.5))
    synth_bg.line.fill.background()
    add_gradient_fill(synth_bg, ROUGE, ROUGE_FONCE)
    add_shadow(synth_bg, alpha_pct=10)
    add_text_box(slide, PAGE_24["synthese"],
                 Cm(2.5), Cm(23.7), Cm(16), Cm(2.2),
                 font_size=10, color=BLANC, bold=True, align=PP_ALIGN.CENTER)

    add_page_number(slide, 24)
    print("  Page 24 : Comparaison cartes OK")


def page_25_bilan_reflexif(prs):
    """Page 25 - Bilan réflexif (nouveau)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_rect(slide, BLANC)
    add_text_box(slide, PAGE_25["titre"], Cm(1), Cm(1), Cm(19), Cm(1.5),
                 font_size=24, color=ROUGE, bold=True, align=PP_ALIGN.CENTER)

    add_text_box(slide, "Réflexion personnelle sur mes activités professionnelles en lien avec le BTS MEC",
                 Cm(1), Cm(2.5), Cm(19), Cm(1),
                 font_size=10, color=GRIS_FONCE, align=PP_ALIGN.CENTER)

    blocs = PAGE_25["blocs"]
    colors = [ROUGE, ORANGE, VERT]
    bg_colors = [RGBColor(0xFF, 0xEB, 0xEB), ORANGE_CLAIR, VERT_CLAIR]
    icons = ["1", "2", "3"]

    y = Cm(4)
    for i, (titre, texte) in enumerate(blocs):
        bloc_h = Cm(7.5)
        # Background
        bloc = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       Cm(1), y, Cm(19), bloc_h)
        bloc.fill.solid()
        bloc.fill.fore_color.rgb = bg_colors[i]
        bloc.line.fill.background()
        add_shadow(bloc, alpha_pct=10)

        # Barre latérale
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(1), y, Cm(0.35), bloc_h)
        bar.fill.solid()
        bar.fill.fore_color.rgb = colors[i]
        bar.line.fill.background()

        # Icône
        circle = add_circle(slide, Cm(1.8), y + Cm(0.5), Cm(1.3), colors[i])
        add_text_box(slide, icons[i], Cm(1.8), y + Cm(0.6), Cm(1.3), Cm(1.1),
                     font_size=14, color=BLANC, bold=True, align=PP_ALIGN.CENTER)

        # Titre
        add_text_box(slide, titre, Cm(3.5), y + Cm(0.5), Cm(15), Cm(1.2),
                     font_size=14, color=colors[i], bold=True)

        # Texte
        add_text_box(slide, texte, Cm(3.5), y + Cm(2), Cm(16), bloc_h - Cm(2.5),
                     font_size=10, color=NOIR)

        y += bloc_h + Cm(0.5)

    add_page_number(slide, 25)
    print("  Page 25 : Bilan réflexif OK")


def page_26_protocole_bim(prs):
    """Page 26 - Protocole BIM (nouveau)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_rect(slide, RGBColor(0x1A, 0x1A, 0x2E))

    add_text_box(slide, PAGE_26["titre"], Cm(1), Cm(1), Cm(19), Cm(1.5),
                 font_size=22, color=BLANC, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, PAGE_26["sous_titre"], Cm(1), Cm(2.5), Cm(19), Cm(1),
                 font_size=12, color=ROUGE_CLAIR, align=PP_ALIGN.CENTER)

    # Convention BIM - items
    conv = PAGE_26["convention"]
    add_text_box(slide, conv["titre"], Cm(1), Cm(4.5), Cm(19), Cm(1),
                 font_size=16, color=ROUGE_CLAIR, bold=True)

    y = Cm(6)
    for label, desc in conv["items"]:
        # Label
        add_text_box(slide, label, Cm(1.5), y, Cm(6), Cm(0.7),
                     font_size=11, color=BLANC, bold=True)
        # Description
        add_text_box(slide, desc, Cm(8), y, Cm(12), Cm(0.7),
                     font_size=10, color=GRIS_TRES_CLAIR)
        y += Cm(1.5)

    # Workflow en étapes
    add_text_box(slide, "WORKFLOW BIM → ÉCONOMIE", Cm(1), Cm(13), Cm(19), Cm(1),
                 font_size=16, color=ROUGE_CLAIR, bold=True)

    workflow = PAGE_26["workflow"]
    x = Cm(0.5)
    for i, step in enumerate(workflow):
        step_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                          x, Cm(14.5), Cm(3.5), Cm(3.5))
        step_bg.fill.solid()
        step_bg.fill.fore_color.rgb = RGBColor(0x2D, 0x2D, 0x44)
        step_bg.line.color.rgb = ROUGE
        step_bg.line.width = Pt(1)
        # Step number
        add_text_box(slide, str(i + 1), x + Cm(1.2), Cm(14.7), Cm(1), Cm(1),
                     font_size=16, color=ROUGE_CLAIR, bold=True, align=PP_ALIGN.CENTER)
        # Step text (remove number prefix)
        step_text = step.split(". ", 1)[1] if ". " in step else step
        add_text_box(slide, step_text, x + Cm(0.2), Cm(15.8), Cm(3.1), Cm(2),
                     font_size=8, color=BLANC, align=PP_ALIGN.CENTER)
        if i < len(workflow) - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                            x + Cm(3.5), Cm(15.8), Cm(0.5), Cm(0.8))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = ROUGE
            arrow.line.fill.background()
        x += Cm(4)

    # Cas concret AFPA
    cas = PAGE_26.get("cas_concret", {})
    if cas:
        cas_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                         Cm(1), Cm(19), Cm(19), Cm(3.5))
        cas_bg.fill.solid()
        cas_bg.fill.fore_color.rgb = RGBColor(0x12, 0x12, 0x22)
        cas_bg.line.color.rgb = ORANGE
        cas_bg.line.width = Pt(1.5)
        add_text_box(slide, cas["titre"], Cm(1.5), Cm(19.2), Cm(18), Cm(0.8),
                     font_size=11, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
        add_text_box(slide, cas["details"],
                     Cm(2), Cm(20.2), Cm(17), Cm(2.2),
                     font_size=9, color=GRIS_TRES_CLAIR, align=PP_ALIGN.CENTER)

    # Apport MEC
    apport_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        Cm(1), Cm(23), Cm(19), Cm(3))
    apport_bg.fill.solid()
    apport_bg.fill.fore_color.rgb = RGBColor(0x12, 0x12, 0x22)
    apport_bg.line.color.rgb = VERT
    apport_bg.line.width = Pt(1.5)
    add_text_box(slide, "APPORT DU BIM POUR L'ÉCONOMISTE MEC", Cm(1.5), Cm(23.2), Cm(18), Cm(0.8),
                 font_size=11, color=VERT, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, PAGE_26["apport_mec"],
                 Cm(2), Cm(24.2), Cm(17), Cm(1.8),
                 font_size=9, color=GRIS_TRES_CLAIR, align=PP_ALIGN.CENTER)

    # Logo BIMCO
    add_image_safe(slide, img("logo"), Cm(8), Cm(26.5), Cm(5))

    add_page_number(slide, 26)
    print("  Page 26 : Protocole BIM OK")


def page_27_projet_pro(prs):
    """Page 27 - Projet professionnel."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_rect(slide, BLANC)
    add_text_box(slide, PAGE_27["titre"], Cm(1), Cm(1), Cm(19), Cm(1.5),
                 font_size=24, color=ROUGE, bold=True, align=PP_ALIGN.CENTER)

    horizons = PAGE_27["horizons"]
    x = Cm(0.5)
    for label, date, desc in horizons:
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       x, Cm(4), Cm(6.3), Cm(17))
        card.fill.solid()
        card.fill.fore_color.rgb = BLANC
        card.line.fill.background()
        add_shadow(card, alpha_pct=18)
        # Top accent
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                         x, Cm(4), Cm(6.3), Cm(1.5))
        accent.fill.solid()
        accent.fill.fore_color.rgb = ROUGE
        accent.line.fill.background()
        add_text_box(slide, label, x, Cm(4.1), Cm(6.3), Cm(0.8),
                     font_size=11, color=BLANC, bold=True, align=PP_ALIGN.CENTER)
        add_text_box(slide, date, x, Cm(4.9), Cm(6.3), Cm(0.6),
                     font_size=9, color=RGBColor(0xFF, 0xCC, 0xCC), align=PP_ALIGN.CENTER)
        add_text_box(slide, desc, x + Cm(0.3), Cm(6), Cm(5.7), Cm(14),
                     font_size=10, color=NOIR)
        x += Cm(6.8)

    # Citation
    quote_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       Cm(2), Cm(22), Cm(17), Cm(3))
    quote_bg.fill.solid()
    quote_bg.fill.fore_color.rgb = GRIS_CLAIR
    quote_bg.line.fill.background()
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(2), Cm(22), Cm(0.3), Cm(3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ROUGE
    bar.line.fill.background()
    add_text_box(slide, PAGE_27["citation"],
                 Cm(3), Cm(22.2), Cm(15), Cm(2.7),
                 font_size=12, color=ROUGE, bold=True, align=PP_ALIGN.CENTER)

    add_text_box(slide, "Ma double compétence – économiste formé au terrain ET développeur "
                 "maîtrisant le BIM – constitue un avantage différenciant rare dans le secteur.",
                 Cm(1), Cm(25.5), Cm(19), Cm(1.5), font_size=9, color=GRIS_FONCE,
                 align=PP_ALIGN.CENTER)

    add_page_number(slide, 27)
    print("  Page 27 : Projet professionnel OK")


def page_28_conclusion(prs):
    """Page 28 - Conclusion."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_fullpage_photo_overlay(slide, img("conclusion"), ROUGE_FONCE, opacity=70)
    add_islamic_star(slide, Cm(3), Cm(3), Cm(3), BLANC, opacity=10)
    add_islamic_star(slide, int(A4_W) - Cm(3), Cm(3), Cm(3), BLANC, opacity=10)

    add_text_box(slide, PAGE_28["titre"], Cm(2), Cm(2), Cm(17), Cm(2),
                 font_size=32, color=BLANC, bold=True, align=PP_ALIGN.CENTER)

    y = Cm(6)
    for title, desc in PAGE_28["points"]:
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       Cm(2), y, Cm(17), Cm(3.5))
        card.fill.solid()
        card.fill.fore_color.rgb = BLANC
        card.line.fill.background()
        set_transparency(card, 90)
        add_shadow(card, alpha_pct=20)
        add_text_box(slide, title, Cm(2.5), y + Cm(0.3), Cm(16), Cm(1.2),
                     font_size=14, color=ROUGE, bold=True)
        add_text_box(slide, desc, Cm(2.5), y + Cm(1.7), Cm(16), Cm(1.5),
                     font_size=11, color=NOIR)
        y += Cm(4)

    # Citation
    add_text_box(slide, PAGE_28["citation"],
                 Cm(2), Cm(22.5), Cm(17), Cm(3),
                 font_size=12, color=BLANC, bold=True, align=PP_ALIGN.CENTER)

    add_text_box(slide, PAGE_28["pied"],
                 Cm(2), Cm(26), Cm(17), Cm(1),
                 font_size=10, color=RGBColor(0xDD, 0xDD, 0xDD), align=PP_ALIGN.CENTER)

    add_page_number(slide, 28)
    print("  Page 28 : Conclusion OK")


def page_29_annexe_documents(prs):
    """Page 29 - Annexe 1 : Documents officiels."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_rect(slide, BLANC)
    add_text_box(slide, PAGE_29["titre"], Cm(1), Cm(1), Cm(19), Cm(1.5),
                 font_size=22, color=ROUGE, bold=True, align=PP_ALIGN.CENTER)

    add_text_box(slide, "Documents attestant de la réalité des activités décrites dans ce rapport",
                 Cm(1), Cm(2.5), Cm(19), Cm(1),
                 font_size=10, color=GRIS_FONCE, align=PP_ALIGN.CENTER)

    docs = PAGE_29["documents"]
    x = Cm(0.5)
    for img_key, titre, description in docs:
        # Cadre pour le document
        frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        x, Cm(4), Cm(6.5), Cm(20))
        frame.fill.solid()
        frame.fill.fore_color.rgb = GRIS_CLAIR
        frame.line.color.rgb = GRIS
        frame.line.width = Pt(0.5)

        # Image
        add_image_safe(slide, img(img_key), x + Cm(0.3), Cm(4.5), Cm(5.9), Cm(14))

        # Titre
        add_text_box(slide, titre, x, Cm(19), Cm(6.5), Cm(1),
                     font_size=12, color=ROUGE, bold=True, align=PP_ALIGN.CENTER)
        # Description
        add_text_box(slide, description, x, Cm(20.5), Cm(6.5), Cm(2),
                     font_size=9, color=GRIS_FONCE, align=PP_ALIGN.CENTER)

        x += Cm(7)

    add_page_number(slide, 29)
    print("  Page 29 : Annexe documents OK")


def page_30_annexe_photos(prs):
    """Page 30 - Annexe 2 : Photos de chantier."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_rect(slide, BLANC)
    add_text_box(slide, PAGE_30["titre"], Cm(1), Cm(1), Cm(19), Cm(1.5),
                 font_size=22, color=ROUGE, bold=True, align=PP_ALIGN.CENTER)

    add_text_box(slide, "Sélection de photos prises sur les chantiers suivis",
                 Cm(1), Cm(2.5), Cm(19), Cm(0.8),
                 font_size=10, color=GRIS_FONCE, align=PP_ALIGN.CENTER)

    photos = PAGE_30["photos"]
    for i, (photo_key, caption) in enumerate(photos):
        col = i % 3
        row = i // 3
        x = Cm(0.5) + col * Cm(7)
        y = Cm(4) + row * Cm(8)

        # Photo frame
        frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        x, y, Cm(6.5), Cm(7))
        frame.fill.solid()
        frame.fill.fore_color.rgb = GRIS_CLAIR
        frame.line.color.rgb = GRIS
        frame.line.width = Pt(0.5)

        # Photo
        add_image_safe(slide, img(photo_key), x + Cm(0.2), y + Cm(0.2),
                        Cm(6.1), Cm(5.5))

        # Caption
        add_text_box(slide, caption, x, y + Cm(5.8), Cm(6.5), Cm(1),
                     font_size=8, color=GRIS_FONCE, align=PP_ALIGN.CENTER)

    add_page_number(slide, 30)
    print("  Page 30 : Annexe photos OK")


# =============================================================================
# MAIN
# =============================================================================
def main():
    print("=" * 60)
    print("GENERATION DU RAPPORT U62 v9 – 30 PAGES RESTRUCTURÉES")
    print("  5 situations CPAR | Synthèse activités | Protocole BIM")
    print("=" * 60)

    prs = Presentation()
    prs.slide_width = A4_W
    prs.slide_height = A4_H

    found = sum(1 for p in PHOTO.values() if os.path.exists(p))
    print(f"\nPhotos disponibles : {found}/{len(PHOTO)}")

    print("\nGénération des 30 pages...")

    page_01_couverture(prs)        # 1  - Couverture
    page_02_fiche_candidat(prs)    # 2  - Fiche candidat
    page_03_sommaire(prs)          # 3  - Sommaire
    page_04_intro_parcours(prs)    # 4  - Introduction + Parcours (fusion)
    page_05_separateur_cadre(prs)  # 5  - Séparateur Cadre professionnel
    page_06_conseil_poste(prs)     # 6  - Conseil Régional + Poste (fusion)
    page_07_bimco_outils(prs)      # 7  - BIMCO + Outils + App (fusion)
    page_08_separateur_projet1(prs) # 8 - Séparateur Projet 1
    page_09_projet1_fiche(prs)     # 9  - Fiche Projet 1
    page_10_budget_corps(prs)      # 10 - Budget + 8 corps d'état (fusion)
    page_11_situation_1(prs)       # 11 - SITUATION 1 : Estimation
    page_12_situation_2(prs)       # 12 - SITUATION 2 : Analyse offres
    page_13_situation_3(prs)       # 13 - SITUATION 3 : Suivi financier
    page_14_situation_4(prs)       # 14 - SITUATION 4 : Communication
    page_15_difficultes_p1(prs)    # 15 - Difficultés P1
    page_16_separateur_projet2(prs) # 16 - Séparateur Projet 2
    page_17_projet2_metres(prs)    # 17 - Fiche P2 + Métrés (fusion)
    page_18_budget_route(prs)      # 18 - Budget route
    page_19_situation_5(prs)       # 19 - SITUATION 5 : Cubatures
    page_20_defis_route(prs)       # 20 - Défis route
    page_21_complementaires(prs)   # 21 - Activités complémentaires
    page_22_separateur_bilan(prs)  # 22 - Séparateur Bilan
    page_23_synthese_competences(prs) # 23 - Tableau synthèse
    page_24_comparaison(prs)       # 24 - Comparaison Maroc/France
    page_25_bilan_reflexif(prs)    # 25 - Bilan réflexif
    page_26_protocole_bim(prs)     # 26 - Protocole BIM
    page_27_projet_pro(prs)        # 27 - Projet professionnel
    page_28_conclusion(prs)        # 28 - Conclusion
    page_29_annexe_documents(prs)  # 29 - Annexe 1 : Documents
    page_30_annexe_photos(prs)     # 30 - Annexe 2 : Photos

    print(f"\nSauvegarde du fichier...")
    prs.save(OUTPUT)
    size = os.path.getsize(OUTPUT) / 1024
    print(f"\n{'=' * 60}")
    print(f"RAPPORT v9 GÉNÉRÉ AVEC SUCCÈS !")
    print(f"  Fichier : {OUTPUT}")
    print(f"  Pages   : {len(prs.slides)}")
    print(f"  Taille  : {size:.0f} Ko")
    print(f"{'=' * 60}")

if __name__ == "__main__":
    main()
